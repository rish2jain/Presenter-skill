"""Tests for edit_deck position-based slide operations."""
import sys
import zipfile
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
sys.path.insert(0, str(SCRIPTS))

from build_deck import build  # noqa: E402
from edit_deck import _resolve_position, duplicate, remove, list_slides  # noqa: E402


@pytest.fixture
def unpacked_three_slides(tmp_path):
    outline = ROOT / "assets" / "example-outline.md"
    pptx = tmp_path / "deck.pptx"
    assert build(outline, pptx)
    unpacked = tmp_path / "unpacked"
    from edit_deck import unpack

    unpack(pptx, unpacked)
    return unpacked


def test_resolve_position_matches_list_order(unpacked_three_slides):
    src = unpacked_three_slides
    assert _resolve_position(src, 1) == 1
    assert _resolve_position(src, 3) == 3
    assert _resolve_position(src, 99) is None


def test_remove_by_position(unpacked_three_slides):
    src = unpacked_three_slides
    assert remove(src, 1)
    assert _resolve_position(src, 1) == 2  # was slide 2, now first


def test_duplicate_by_position(unpacked_three_slides, tmp_path):
    src = unpacked_three_slides
    assert duplicate(src, 2)
    out = tmp_path / "packed.pptx"
    from edit_deck import pack

    assert pack(src, out)
    with zipfile.ZipFile(out) as zf:
        slide_parts = [n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
    assert len(slide_parts) == 7  # 6 original + 1 duplicate


def test_inventory_and_replace_roundtrip(tmp_path):
    import json
    from pptx import Presentation
    from pptx.util import Inches
    import edit_deck

    # fixture deck: one slide, one textbox
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "Old headline"
    run.font.bold = True
    src = tmp_path / "deck.pptx"
    prs.save(str(src))

    work = tmp_path / "unpacked"
    edit_deck.unpack(str(src), str(work))
    inv = edit_deck.inventory(str(work))
    assert inv == [{"slide": 1, "run": 0, "text": "Old headline"}]

    edits = tmp_path / "edits.json"
    edits.write_text(json.dumps(
        [{"slide": 1, "run": 0, "text": "New headline"}]))
    edit_deck.replace_runs(str(work), str(edits))
    out = tmp_path / "out.pptx"
    edit_deck.pack(str(work), str(out))

    prs2 = Presentation(str(out))
    runs = prs2.slides[0].shapes[0].text_frame.paragraphs[0].runs
    assert runs[0].text == "New headline"
    assert runs[0].font.bold  # formatting preserved


def test_replace_out_of_range_run_raises(tmp_path):
    import json
    import pytest
    from pptx import Presentation
    from pptx.util import Inches
    import edit_deck
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.text_frame.paragraphs[0].add_run().text = "x"
    src = tmp_path / "deck.pptx"
    prs.save(str(src))
    work = tmp_path / "unpacked"
    edit_deck.unpack(str(src), str(work))
    edits = tmp_path / "edits.json"
    edits.write_text(json.dumps([{"slide": 1, "run": 999, "text": "y"}]))
    with pytest.raises(SystemExit, match="out of range"):
        edit_deck.replace_runs(str(work), str(edits))


# ── selection parsing (extract / append --slides) ───────────────────────────

def test_parse_selection_range():
    from edit_deck import _parse_selection
    assert _parse_selection("3-7", 10) == [3, 4, 5, 6, 7]


def test_parse_selection_list_and_combo():
    from edit_deck import _parse_selection
    assert _parse_selection("2,5,9", 10) == [2, 5, 9]
    assert _parse_selection("1,3-5", 10) == [1, 3, 4, 5]
    assert _parse_selection("2,2,3-4,3", 10) == [2, 3, 4]  # deduped


def test_parse_selection_errors():
    from edit_deck import _parse_selection
    with pytest.raises(ValueError, match="empty"):
        _parse_selection("", 5)
    with pytest.raises(ValueError, match="out of range"):
        _parse_selection("6", 5)
    with pytest.raises(ValueError, match="out of range"):
        _parse_selection("0", 5)
    with pytest.raises(ValueError):
        _parse_selection("abc", 5)
    with pytest.raises(ValueError, match="reversed"):
        _parse_selection("5-3", 5)


# ── extract (split) ──────────────────────────────────────────────────────────

def _first_texts(pptx_path):
    """First non-empty text of each slide, in deck order."""
    from pptx import Presentation
    out = []
    for slide in Presentation(str(pptx_path)).slides:
        first = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                first = shape.text_frame.text.strip().splitlines()[0]
                break
        out.append(first)
    return out


@pytest.fixture(scope="module")
def example_deck(tmp_path_factory):
    """The 6-slide example deck built once for extract/append tests."""
    outline = ROOT / "assets" / "example-outline.md"
    pptx = tmp_path_factory.mktemp("example") / "deck.pptx"
    assert build(outline, pptx)
    return pptx


def test_extract_range_end_to_end(example_deck, tmp_path):
    from pptx import Presentation
    from edit_deck import extract
    out = tmp_path / "sub.pptx"
    assert extract(example_deck, "2-4", out)
    prs = Presentation(str(out))
    assert len(prs.slides) == 3
    originals = _first_texts(example_deck)
    assert _first_texts(out) == originals[1:4]


def test_extract_comma_list(example_deck, tmp_path):
    from pptx import Presentation
    from edit_deck import extract
    out = tmp_path / "sub.pptx"
    assert extract(example_deck, "1,6", out)
    prs = Presentation(str(out))
    assert len(prs.slides) == 2
    originals = _first_texts(example_deck)
    assert _first_texts(out) == [originals[0], originals[5]]


def test_extract_out_of_range_errors(example_deck, tmp_path):
    from edit_deck import extract
    assert not extract(example_deck, "5-9", tmp_path / "sub.pptx")
    assert not extract(example_deck, "", tmp_path / "sub.pptx")


def test_extract_output_equals_input_errors(example_deck):
    from edit_deck import extract
    assert not extract(example_deck, "1-2", example_deck)


# ── append (merge) ───────────────────────────────────────────────────────────

# Minimal valid 1x1 red PNG (no Pillow needed to create the fixture).
_PNG_1PX = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108020000009077"
    "53de0000000c4944415408d763f8cfc000000301010018dd8db00000000049"
    "454e44ae426082")


def _simple_src_deck(path, with_image=False):
    """Two-slide python-pptx deck: textboxes + optional picture on slide 2."""
    import io
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    for i, label in enumerate(("Alpha source slide", "Beta source slide")):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = label
        run.font.size = Pt(24)
        if with_image and i == 1:
            slide.shapes.add_picture(io.BytesIO(_PNG_1PX), Inches(1),
                                     Inches(3), Inches(2), Inches(2))
    prs.save(str(path))
    return path


def test_append_end_to_end(example_deck, tmp_path):
    import subprocess
    from pptx import Presentation
    from edit_deck import append_decks
    src = _simple_src_deck(tmp_path / "src.pptx", with_image=True)
    merged = tmp_path / "merged.pptx"
    assert append_decks(example_deck, src, None, merged)

    prs = Presentation(str(merged))
    assert len(prs.slides) == 8
    texts = _first_texts(merged)
    assert texts[6] == "Alpha source slide"
    assert texts[7] == "Beta source slide"
    # image copied along
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    assert any(s.shape_type == MSO_SHAPE_TYPE.PICTURE
               for s in prs.slides[7].shapes)

    # no duplicate part names in the package
    with zipfile.ZipFile(merged) as zf:
        names = zf.namelist()
        assert len(names) == len(set(names))
        # src master/theme imported alongside dst's own (faithful import)
        assert any("slideMaster2.xml" in n and "_rels" not in n for n in names)

    # qa_check opens the merged deck without integrity blowups
    proc = subprocess.run(
        [sys.executable, str(SCRIPTS / "qa_check.py"), str(merged)],
        capture_output=True, text=True)
    assert proc.returncode in (0, 1), proc.stderr
    assert "Traceback" not in proc.stderr
    assert "error(s)" in proc.stdout


def test_append_subset(example_deck, tmp_path):
    from pptx import Presentation
    from edit_deck import append_decks
    src = _simple_src_deck(tmp_path / "src.pptx")
    merged = tmp_path / "merged.pptx"
    assert append_decks(example_deck, src, "2", merged)
    prs = Presentation(str(merged))
    assert len(prs.slides) == 7
    assert _first_texts(merged)[6] == "Beta source slide"


def test_append_chart_slide(example_deck, tmp_path):
    """Chart-bearing source slide: chart part + embedded xlsx travel along."""
    from pptx import Presentation
    from edit_deck import append_decks
    dst = _simple_src_deck(tmp_path / "dst.pptx")
    merged = tmp_path / "merged.pptx"
    assert append_decks(dst, example_deck, "3", merged)
    prs = Presentation(str(merged))
    assert len(prs.slides) == 3
    copied = prs.slides[2]
    assert any(getattr(s, "has_chart", False) for s in copied.shapes)
    chart = next(s for s in copied.shapes if getattr(s, "has_chart", False)).chart
    assert chart.plots  # chart part loads
    with zipfile.ZipFile(merged) as zf:
        names = zf.namelist()
        assert any(n.startswith("ppt/charts/chart") for n in names)
        assert any(n.startswith("ppt/embeddings/") for n in names)


def test_append_dedupes_shared_layout_and_master(example_deck, tmp_path):
    """Two src slides sharing a master → master/theme copied exactly once."""
    from edit_deck import append_decks
    src = _simple_src_deck(tmp_path / "src.pptx")
    merged = tmp_path / "merged.pptx"
    assert append_decks(example_deck, src, "1-2", merged)
    with zipfile.ZipFile(merged) as zf:
        masters = [n for n in zf.namelist()
                   if n.startswith("ppt/slideMasters/slideMaster")
                   and n.endswith(".xml")]
    assert len(masters) == 2  # dst's own + one imported


def test_append_invalid_selection_errors(example_deck, tmp_path):
    from edit_deck import append_decks
    src = _simple_src_deck(tmp_path / "src.pptx")
    assert not append_decks(example_deck, src, "9", tmp_path / "m.pptx")


def test_append_multi_layout_master_pruned(example_deck, tmp_path):
    """Src slides on two different layouts of one master: both layouts
    copied, master's sldLayoutIdLst pruned to exactly those two, with
    fresh presentation-unique ids."""
    from lxml import etree
    from pptx import Presentation
    from edit_deck import append_decks, NS
    prs = Presentation()
    for li in (0, 5):
        slide = prs.slides.add_slide(prs.slide_layouts[li])
        if slide.shapes.title is not None:
            slide.shapes.title.text = f"Layout {li} slide"
    src = tmp_path / "src.pptx"
    prs.save(str(src))
    merged = tmp_path / "merged.pptx"
    assert append_decks(example_deck, src, None, merged)
    Presentation(str(merged))  # opens cleanly
    with zipfile.ZipFile(merged) as zf:
        master = etree.fromstring(zf.read("ppt/slideMasters/slideMaster2.xml"))
        layout_ids = master.findall("p:sldLayoutIdLst/p:sldLayoutId", NS)
        rels = etree.fromstring(
            zf.read("ppt/slideMasters/_rels/slideMaster2.xml.rels"))
        layout_rels = [r for r in rels
                       if r.get("Type", "").endswith("/slideLayout")]
    assert len(layout_ids) == 2
    assert len(layout_rels) == 2
    ids = [int(e.get("id")) for e in layout_ids]
    assert all(i > 2147483647 for i in ids) and len(set(ids)) == 2
    # ids must not collide with dst's existing master/layout ids
    with zipfile.ZipFile(example_deck) as zf:
        dst_master = etree.fromstring(
            zf.read("ppt/slideMasters/slideMaster1.xml"))
        dst_ids = {int(e.get("id")) for e in
                   dst_master.findall("p:sldLayoutIdLst/p:sldLayoutId", NS)}
    assert not set(ids) & dst_ids


# ── dangling slide-jump hyperlinks (extract / append) ───────────────────────

def _rewrite_zip(path, replacements):
    """Rewrite zip members in place. replacements: {name: bytes}."""
    import io
    buf = io.BytesIO()
    with zipfile.ZipFile(path) as zin, \
            zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
        pending = dict(replacements)
        for item in zin.infolist():
            data = pending.pop(item.filename, None)
            if data is None:
                data = zin.read(item.filename)
            zout.writestr(item.filename, data)
        for name, data in pending.items():
            zout.writestr(name, data)
    Path(path).write_bytes(buf.getvalue())


def _inject_slide_jump(pptx, slide_name, rid, target_slide):
    """Add an in-XML slide-jump hlinkClick (rid) + matching Relationship."""
    from lxml import etree
    from edit_deck import NS, SLIDE_RT
    with zipfile.ZipFile(pptx) as zf:
        slide_xml = etree.fromstring(zf.read(f"ppt/slides/{slide_name}"))
        rels_xml = etree.fromstring(
            zf.read(f"ppt/slides/_rels/{slide_name}.rels"))
    rpr = slide_xml.find(".//a:rPr", NS)
    assert rpr is not None
    h = etree.SubElement(rpr, f"{{{NS['a']}}}hlinkClick")
    h.set(f"{{{NS['r']}}}id", rid)
    h.set("action", "ppaction://hlinksldjump")
    rel = etree.SubElement(rels_xml, f"{{{NS['rel']}}}Relationship")
    rel.set("Id", rid)
    rel.set("Type", SLIDE_RT)
    rel.set("Target", target_slide)
    _rewrite_zip(pptx, {
        f"ppt/slides/{slide_name}": etree.tostring(
            slide_xml, xml_declaration=True, encoding="UTF-8",
            standalone=True),
        f"ppt/slides/_rels/{slide_name}.rels": etree.tostring(
            rels_xml, xml_declaration=True, encoding="UTF-8",
            standalone=True),
    })


def _assert_no_dangling_hlinks(pptx):
    """Every hlinkClick/hlinkHover r:id in every slide exists in its rels."""
    from lxml import etree
    from edit_deck import NS
    with zipfile.ZipFile(pptx) as zf:
        slides = [n for n in zf.namelist()
                  if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        for n in slides:
            tree = etree.fromstring(zf.read(n))
            rels_name = f"ppt/slides/_rels/{Path(n).name}.rels"
            rel_ids = set()
            if rels_name in zf.namelist():
                rels = etree.fromstring(zf.read(rels_name))
                rel_ids = {r.get("Id") for r in rels}
            for tag in ("hlinkClick", "hlinkHover"):
                for h in tree.iter(f"{{{NS['a']}}}{tag}"):
                    rid = h.get(f"{{{NS['r']}}}id")
                    assert not rid or rid in rel_ids, \
                        f"{n}: {tag} references missing {rid}"


def test_append_strips_dangling_hlink_click(tmp_path):
    """Slide-jump to an unselected slide: rel dropped AND hlink stripped."""
    from pptx import Presentation
    from edit_deck import append_decks
    dst = _simple_src_deck(tmp_path / "dst.pptx")
    src = _simple_src_deck(tmp_path / "src.pptx")
    _inject_slide_jump(src, "slide1.xml", "rId99", "slide2.xml")
    merged = tmp_path / "merged.pptx"
    assert append_decks(dst, src, "1", merged)  # slide 2 unselected
    _assert_no_dangling_hlinks(merged)
    prs = Presentation(str(merged))
    assert len(prs.slides) == 3


def test_extract_strips_dangling_hlink_click_and_warns(tmp_path, capsys):
    """Kept slide jumping to a removed slide: rel + hlink stripped, WARN."""
    from lxml import etree
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from edit_deck import extract, NS
    prs = Presentation()
    for label in ("One", "Two", "Three"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1), Inches(1),
                                      Inches(6), Inches(1))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = label
        run.font.size = Pt(24)
    deck = tmp_path / "deck.pptx"
    prs.save(str(deck))
    _inject_slide_jump(deck, "slide1.xml", "rId99", "slide3.xml")

    out = tmp_path / "sub.pptx"
    assert extract(deck, "1", out)
    captured = capsys.readouterr()
    assert "WARN" in captured.err

    with zipfile.ZipFile(out) as zf:
        rels_name = next(n for n in zf.namelist()
                         if n.startswith("ppt/slides/_rels/"))
        rels = etree.fromstring(zf.read(rels_name))
        assert not any("slide3" in (r.get("Target") or "")
                       for r in rels), "rels still target removed slide"
    _assert_no_dangling_hlinks(out)
    assert len(Presentation(str(out)).slides) == 1


def test_append_media_name_collision_preserved(tmp_path):
    """dst and src both carry ppt/media/image1.png with different bytes:
    dst's original bytes stay under its name; src's land under a new one."""
    from pptx import Presentation
    from edit_deck import append_decks
    dst = _simple_src_deck(tmp_path / "dst.pptx", with_image=True)
    src = _simple_src_deck(tmp_path / "src.pptx", with_image=True)
    src_bytes = _PNG_1PX + b"src-variant-trailing-bytes"
    _rewrite_zip(src, {"ppt/media/image1.png": src_bytes})
    with zipfile.ZipFile(dst) as zf:
        dst_bytes = zf.read("ppt/media/image1.png")
    assert dst_bytes != src_bytes

    merged = tmp_path / "merged.pptx"
    assert append_decks(dst, src, None, merged)
    with zipfile.ZipFile(merged) as zf:
        assert zf.read("ppt/media/image1.png") == dst_bytes
        others = [n for n in zf.namelist()
                  if n.startswith("ppt/media/") and n != "ppt/media/image1.png"]
        assert any(zf.read(n) == src_bytes for n in others), \
            "src image bytes missing from merged deck"
    Presentation(str(merged))  # reopens cleanly


def test_append_double_append_unique_parts(example_deck, tmp_path):
    """Appending the same src twice: unique part names, unique master ids."""
    from lxml import etree
    from pptx import Presentation
    from edit_deck import append_decks, NS
    src = _simple_src_deck(tmp_path / "src.pptx")
    m1 = tmp_path / "m1.pptx"
    m2 = tmp_path / "m2.pptx"
    assert append_decks(example_deck, src, None, m1)
    assert append_decks(m1, src, None, m2)

    with zipfile.ZipFile(m2) as zf:
        names = zf.namelist()
        assert len(names) == len(set(names))
        pres = etree.fromstring(zf.read("ppt/presentation.xml"))
        master_ids = [int(e.get("id")) for e in
                      pres.findall("p:sldMasterIdLst/p:sldMasterId", NS)]
        layout_ids = []
        for n in names:
            if (n.startswith("ppt/slideMasters/slideMaster")
                    and n.endswith(".xml") and "_rels" not in n):
                m = etree.fromstring(zf.read(n))
                layout_ids += [int(e.get("id")) for e in
                               m.findall("p:sldLayoutIdLst/p:sldLayoutId", NS)]
    assert len(master_ids) == 3  # dst's own + two imports
    all_ids = master_ids + layout_ids
    assert len(all_ids) == len(set(all_ids)), "master/layout ids collide"
    prs = Presentation(str(m2))
    assert len(prs.slides) == 10  # 6 + 2 + 2
