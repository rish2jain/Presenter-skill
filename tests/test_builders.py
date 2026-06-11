"""Tests for layout builders."""
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
sys.path.insert(0, str(SCRIPTS))

from builders import _slide_title, build_title_slide, build_closing_slide  # noqa: E402
from build_deck import parse_outline  # noqa: E402
from pptx import Presentation  # noqa: E402
from palettes import get_palette  # noqa: E402


def _slide_text(slide):
    return [s.text_frame.text for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def test_title_uses_slide_heading_when_no_title_field():
    _, slides = parse_outline("## Slide 1: My Deck Title\n**Layout:** title\n")
    prs = Presentation()
    pal = get_palette("midnight-executive")
    slide = build_title_slide(prs, slides[0], pal, {})
    assert "My Deck Title" in _slide_text(slide)
    assert "Presentation Title" not in _slide_text(slide)


def test_closing_uses_slide_heading():
    _, slides = parse_outline("## Slide 1: Questions?\n**Layout:** closing\n")
    prs = Presentation()
    pal = get_palette("midnight-executive")
    slide = build_closing_slide(prs, slides[0], pal, {})
    assert "Questions?" in _slide_text(slide)
    assert "Thank You" not in _slide_text(slide)


def test_slide_title_helper():
    assert _slide_title({"title": "A", "heading": "B"}, "default") == "A"
    assert _slide_title({"heading": "B"}, "default") == "B"
    assert _slide_title({}, "default") == "default"
