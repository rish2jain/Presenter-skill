"""Tests for geometry_report.py per-slide layout metrics."""
import io
import struct
import sys
import zlib
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))

from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

import geometry_report as geo  # noqa: E402


def _prs():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
    return prs


def _blank(prs):
    layout = min(prs.slide_layouts, key=lambda l: len(l.placeholders))
    return prs.slides.add_slide(layout)


def _rect(slide, left, top, w, h):
    return slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))


def _solid_rect(slide, left, top, w, h):
    shape = _rect(slide, left, top, w, h)
    shape.fill.solid()
    return shape


def _textbox(slide, left, top, w, h, text="word " * 40, pt=18):
    """Text box whose estimated glyph area fills the box."""
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(w), Inches(h))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(pt)
    return tb


def _png_bytes():
    """Minimal 1x1 RGB PNG, no Pillow dependency."""
    def chunk(tag, data):
        body = tag + data
        return struct.pack(">I", len(data)) + body \
            + struct.pack(">I", zlib.crc32(body))
    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    idat = zlib.compress(b"\x00\xff\x00\x00")
    return (b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr)
            + chunk(b"IDAT", idat) + chunk(b"IEND", b""))


def _picture(slide, left, top, w, h):
    return slide.shapes.add_picture(
        io.BytesIO(_png_bytes()),
        Inches(left), Inches(top), Inches(w), Inches(h))


def _slide_metrics(prs):
    return geo.analyze_deck(prs)[1]


def test_text_on_text_overlap_reported():
    prs = _prs()
    slide = _blank(prs)
    _textbox(slide, 1.0, 1.0, 2.0, 2.0)
    _textbox(slide, 2.0, 2.0, 2.0, 2.0)  # 1 sq in intersection
    m = _slide_metrics(prs)
    assert len(m["overlaps"]) == 1, m
    assert abs(m["overlaps"][0]["area_sqin"] - 1.0) < 0.05


def test_text_box_sliver_overlap_not_reported():
    """Boxes intersect but the rendered glyphs cannot: one short line in
    a tall box, the other box only grazing the empty bottom region."""
    prs = _prs()
    slide = _blank(prs)
    _textbox(slide, 1.0, 1.0, 3.0, 1.5, text="One label", pt=12)
    _textbox(slide, 1.0, 2.3, 3.0, 0.5, text="Footnote", pt=10)
    m = _slide_metrics(prs)
    assert m["overlaps"] == [], m


def test_text_under_solid_shape_reported():
    prs = _prs()
    slide = _blank(prs)
    _textbox(slide, 1.0, 1.0, 3.0, 1.0)          # text below
    _solid_rect(slide, 1.0, 1.4, 3.0, 1.0)       # solid band over the text
    m = _slide_metrics(prs)
    assert len(m["overlaps"]) == 1, m


def test_picture_over_text_reported():
    """A picture has no .fill but always paints — drawn over earlier
    body text it occludes the glyphs and must be reported."""
    prs = _prs()
    slide = _blank(prs)
    _textbox(slide, 1.0, 1.0, 3.0, 1.0)          # text below
    _picture(slide, 2.0, 0.8, 4.0, 2.0)          # picture on top of glyphs
    m = _slide_metrics(prs)
    assert len(m["overlaps"]) == 1, m


def test_picture_under_label_suppressed():
    """Label drawn over an image is the standard caption pattern."""
    prs = _prs()
    slide = _blank(prs)
    _picture(slide, 1.0, 1.0, 4.0, 3.0)          # image below
    _textbox(slide, 3.5, 1.5, 2.5, 0.6, text="Chart label", pt=12)
    m = _slide_metrics(prs)
    assert m["overlaps"] == [], m


def test_text_buried_under_later_solid_reported():
    """Earlier text 100% contained in a later-drawn solid shape is
    completely hidden — containment must not suppress it."""
    prs = _prs()
    slide = _blank(prs)
    _textbox(slide, 2.0, 2.0, 4.0, 1.5)          # text drawn first
    _solid_rect(slide, 1.8, 1.8, 4.4, 1.9)       # solid fully covers it
    m = _slide_metrics(prs)
    assert len(m["overlaps"]) == 1, m


def test_text_on_own_card_suppressed():
    """Text box placed on its own card (card drawn first) is the
    standard card pattern — child-on-card containment stays quiet."""
    prs = _prs()
    slide = _blank(prs)
    _solid_rect(slide, 1.0, 1.0, 4.0, 2.0)       # card drawn first
    _textbox(slide, 1.2, 1.2, 3.6, 1.6, text="Card body copy", pt=12)
    m = _slide_metrics(prs)
    assert m["overlaps"] == [], m


def test_label_above_bar_suppressed():
    prs = _prs()
    slide = _blank(prs)
    _solid_rect(slide, 1.0, 1.0, 3.0, 1.0)       # bar below
    _textbox(slide, 2.5, 1.2, 2.0, 0.6)          # label overhanging the bar
    m = _slide_metrics(prs)
    assert m["overlaps"] == [], m


def test_solid_layering_suppressed():
    prs = _prs()
    slide = _blank(prs)
    _solid_rect(slide, 1.0, 1.0, 2.0, 2.0)
    _solid_rect(slide, 2.0, 2.0, 2.0, 2.0)  # deliberate layering on top
    m = _slide_metrics(prs)
    assert m["overlaps"] == [], m


def test_containment_not_reported_as_overlap():
    prs = _prs()
    slide = _blank(prs)
    _rect(slide, 1.0, 1.0, 5.0, 3.0)   # card
    _rect(slide, 1.5, 1.5, 2.0, 1.0)   # fully inside = intentional
    m = _slide_metrics(prs)
    assert m["overlaps"] == [], m


def test_uneven_row_gaps_flagged():
    prs = _prs()
    slide = _blank(prs)
    for left in (1.0, 3.31, 5.60, 8.15):  # gaps 0.31 / 0.29 / 0.55
        _rect(slide, left, 2.0, 2.0, 1.0)
    m = _slide_metrics(prs)
    assert any(g["axis"] == "row" for g in m["uneven_gaps"]), m


def test_even_row_gaps_pass():
    prs = _prs()
    slide = _blank(prs)
    for left in (1.0, 3.3, 5.6, 7.9):  # gaps all 0.30
        _rect(slide, left, 2.0, 2.0, 1.0)
    m = _slide_metrics(prs)
    assert m["uneven_gaps"] == [], m


def test_mixed_size_column_chain_not_a_gap_finding():
    """Title/body/footer stacked at one left edge is not a card column."""
    prs = _prs()
    slide = _blank(prs)
    _rect(slide, 1.0, 0.5, 8.0, 1.1)   # title-sized
    _rect(slide, 1.0, 1.8, 8.0, 0.12)  # divider
    _rect(slide, 1.0, 2.5, 8.0, 2.4)   # body block
    _rect(slide, 1.0, 6.5, 8.0, 0.5)   # footer band
    m = _slide_metrics(prs)
    assert m["uneven_gaps"] == [], m


def test_near_miss_flagged():
    prs = _prs()
    slide = _blank(prs)
    _rect(slide, 1.0, 1.0, 2.0, 1.0)
    _rect(slide, 6.0, 1.05, 2.0, 1.0)  # top edge off by 0.05in
    m = _slide_metrics(prs)
    assert any(n["edge"] == "top" and abs(n["off_in"] - 0.05) < 0.011
               for n in m["near_misses"]), m


def test_aligned_shapes_no_near_miss():
    prs = _prs()
    slide = _blank(prs)
    _rect(slide, 1.0, 1.0, 2.0, 1.0)
    _rect(slide, 6.0, 1.0, 2.0, 1.0)
    m = _slide_metrics(prs)
    assert m["near_misses"] == [], m


def test_co_aligned_opposite_edge_near_miss_suppressed():
    """Shared bottom edge: the 0.05in top difference is a data-driven
    height (bars off a baseline), not a misalignment."""
    prs = _prs()
    slide = _blank(prs)
    _rect(slide, 1.0, 1.05, 2.0, 0.95)  # bottom 2.00, top 1.05
    _rect(slide, 6.0, 1.00, 2.0, 1.00)  # bottom 2.00, top 1.00
    m = _slide_metrics(prs)
    assert m["near_misses"] == [], m


def test_non_co_aligned_near_miss_still_reported():
    """Equal heights shifted 0.05in: tops and bottoms both off, neither
    axis shares an opposite edge — a genuine misalignment."""
    prs = _prs()
    slide = _blank(prs)
    _rect(slide, 1.0, 1.0, 2.0, 1.0)
    _rect(slide, 6.0, 1.05, 2.0, 1.0)
    m = _slide_metrics(prs)
    assert {n["edge"] for n in m["near_misses"]} == {"top", "bottom"}, m


def test_text_box_near_miss_suppressed():
    """Text frames align by glyph edge, so box-edge near-misses with a
    text shape in the pair are skipped."""
    prs = _prs()
    slide = _blank(prs)
    _rect(slide, 1.0, 1.0, 2.0, 1.0)
    _textbox(slide, 6.0, 1.05, 2.0, 1.0, text="Label", pt=12)
    m = _slide_metrics(prs)
    assert m["near_misses"] == [], m


def test_solid_text_card_near_miss_reported():
    """Solid-filled cards bearing text align by their painted box edge,
    so a 0.06in left-edge offset between two cards is a real defect."""
    prs = _prs()
    slide = _blank(prs)
    card1 = _solid_rect(slide, 1.00, 1.0, 3.00, 1.5)
    card1.text_frame.text = "Workstream alpha status"
    card2 = _solid_rect(slide, 1.06, 3.0, 3.10, 1.5)
    card2.text_frame.text = "Workstream beta status"
    m = _slide_metrics(prs)
    assert any(n["edge"] == "left" and abs(n["off_in"] - 0.06) < 0.011
               for n in m["near_misses"]), m


def test_plain_textbox_near_miss_still_suppressed():
    """Borderless unfilled text boxes align by glyph edge — a 0.06in
    box-edge offset between them stays suppressed."""
    prs = _prs()
    slide = _blank(prs)
    _textbox(slide, 1.00, 1.0, 3.00, 1.0, text="Alpha summary", pt=12)
    _textbox(slide, 1.06, 3.0, 3.10, 1.0, text="Beta summary", pt=12)
    m = _slide_metrics(prs)
    assert m["near_misses"] == [], m


def test_near_miss_cap_and_summary_line():
    prs = _prs()
    slide = _blank(prs)
    for i in range(5):                       # aligned, evenly spaced stack
        _rect(slide, 1.0, 1.0 + i * 1.2, 2.0, 1.0)
    _rect(slide, 1.05, 0.0, 2.0, 0.9)        # both x edges off by 0.05
    m = _slide_metrics(prs)
    assert len(m["near_misses"]) > geo.NEAR_MISS_CAP, m
    lines = geo.findings(m)
    near = [f for f in lines if f.startswith("almost aligned")]
    assert len(near) == geo.NEAR_MISS_CAP, lines
    extra = len(m["near_misses"]) - geo.NEAR_MISS_CAP
    assert f"(+{extra} more near-miss alignments suppressed)" in lines, lines


def test_gap_run_split_at_large_jump():
    """A >1.5in jump separates visual groups; the even group on each
    side must not be read as one wildly uneven run."""
    prs = _prs()
    slide = _blank(prs)
    for left in (1.0, 3.3, 5.6, 11.0):  # gaps 0.30 / 0.30 / 3.40
        _rect(slide, left, 2.0, 2.0, 1.0)
    m = _slide_metrics(prs)
    assert m["uneven_gaps"] == [], m


def test_gap_split_groups_still_assessed():
    """Unevenness inside a group survives the split."""
    prs = _prs()
    slide = _blank(prs)
    for left in (0.5, 2.81, 5.10, 7.65):  # gaps 0.31 / 0.29 / 0.55
        _rect(slide, left, 2.0, 2.0, 1.0)
    _rect(slide, 11.3, 2.0, 2.0, 1.0)     # 1.65in jump = separate group
    m = _slide_metrics(prs)
    assert any(g["axis"] == "row" for g in m["uneven_gaps"]), m


def test_word_overload_finding():
    prs = _prs()
    slide = _blank(prs)
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(5))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "word " * 95
    run.font.size = Pt(14)
    m = _slide_metrics(prs)
    assert m["words"] == 95
    assert any("text overload" in f for f in geo.findings(m)), geo.findings(m)


def test_whitespace_ratio_bounds():
    prs = _prs()
    _blank(prs)                                  # empty slide
    _rect(_blank(prs), 0.0, 0.0, 13.33, 7.5)     # fully covered slide
    report = geo.analyze_deck(prs)
    assert report[1]["whitespace_ratio"] == 1.0
    assert report[2]["whitespace_ratio"] < 0.05


def test_imbalance_metrics_present_and_flagged_when_extreme():
    prs = _prs()
    slide = _blank(prs)
    _rect(slide, 0.2, 0.2, 6.0, 7.0)  # everything on the left half
    m = _slide_metrics(prs)
    assert m["lr_imbalance_pp"] > geo.IMBALANCE_PP
    assert any("left/right" in f for f in geo.findings(m)), geo.findings(m)


def test_analyze_deck_slide_filter():
    prs = _prs()
    _blank(prs)
    _blank(prs)
    _blank(prs)
    report = geo.analyze_deck(prs, only={2})
    assert set(report) == {2}


def test_clean_grid_slide_has_no_findings():
    prs = _prs()
    slide = _blank(prs)
    _rect(slide, 0.7, 0.5, 11.9, 1.0)  # title band
    for i in range(3):                 # even card row
        _rect(slide, 0.7 + i * 4.1, 2.0, 3.8, 2.0)
    m = _slide_metrics(prs)
    assert geo.findings(m) == [], geo.findings(m)
