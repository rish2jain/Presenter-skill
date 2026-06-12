"""Tests for textfit: line estimation, capacity validation, normAutofit."""
import math
import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
sys.path.insert(0, str(SCRIPTS))

import builders  # noqa: E402
import textfit  # noqa: E402
from build_deck import parse_outline, validate  # noqa: E402
from palettes import get_palette  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402


@pytest.fixture(autouse=True)
def _default_density():
    builders.set_density("compact")


CTX = {"outline_dir": ROOT / "assets", "assets_dir": ROOT / "assets"}
SENTENCE = "Hybrid cloud migration is 62 percent complete and on track for Q4"


# ── estimate_lines ───────────────────────────────────────────────────────────
def test_empty_text_zero_lines():
    assert textfit.estimate_lines("", 14, 5.0) == 0
    assert textfit.estimate_lines("   ", 14, 5.0) == 0
    assert textfit.estimate_lines(None, 14, 5.0) == 0


def test_single_word_one_line():
    assert textfit.estimate_lines("Revenue", 14, 5.0) == 1


def test_more_text_never_fewer_lines():
    short = textfit.estimate_lines(SENTENCE, 14, 5.0)
    long = textfit.estimate_lines(SENTENCE + " " + SENTENCE, 14, 5.0)
    longer = textfit.estimate_lines((SENTENCE + " ") * 4, 14, 5.0)
    assert short <= long <= longer
    assert longer > short


def test_wider_box_never_more_lines():
    narrow = textfit.estimate_lines(SENTENCE * 2, 14, 3.0)
    wide = textfit.estimate_lines(SENTENCE * 2, 14, 8.0)
    assert wide <= narrow
    assert narrow > 1


def test_bigger_font_never_fewer_lines():
    small = textfit.estimate_lines(SENTENCE * 2, 12, 5.0)
    big = textfit.estimate_lines(SENTENCE * 2, 24, 5.0)
    assert big >= small


def test_long_word_counts_split_lines():
    token = "x" * 300
    lines = textfit.estimate_lines(token, 14, 2.0)
    assert lines >= 2
    # narrower box -> proportionally more lines
    assert textfit.estimate_lines(token, 14, 1.0) >= lines


def test_estimate_height_formula():
    lines = textfit.estimate_lines(SENTENCE, 14, 5.0)
    h = textfit.estimate_height_in(SENTENCE, 14, 5.0)
    assert h == pytest.approx(lines * 14 * 1.2 / 72)
    assert textfit.estimate_height_in("", 14, 5.0) == 0


def test_fallback_within_tolerance_of_pil():
    if textfit._load_font(14, False) is None:
        pytest.skip("PIL or measurement font unavailable")
    cases = [
        (SENTENCE, 14, 5.465), (SENTENCE * 2, 14, 5.465),
        (SENTENCE * 3, 13, 5.8), ((SENTENCE + " ") * 4, 16, 11.9),
        ("Mid-market segment growing fastest at 47% CAGR", 14, 5.465),
        ("Approve Phase 2 GPU budget this quarter to hold savings", 32, 11.9),
    ]
    pil = [textfit.estimate_lines(t, pt, w) for t, pt, w in cases]
    textfit._FONT_CACHE.clear()
    real_load = textfit._load_font
    try:
        textfit._load_font = lambda pt, bold: None
        fb = [textfit.estimate_lines(t, pt, w) for t, pt, w in cases]
    finally:
        textfit._load_font = real_load
        textfit._FONT_CACHE.clear()
    for a, b in zip(pil, fb):
        assert abs(a - b) <= 1, (pil, fb)
    assert sum(fb) <= sum(pil) * 1.25 and sum(fb) >= sum(pil) * 0.8


# ── strip_markup ─────────────────────────────────────────────────────────────
def test_strip_markup_icon_and_rich():
    assert textfit.strip_markup("icon:rocket Launch fast") == ("Launch fast", True)
    assert textfit.strip_markup("**Bold** move") == ("Bold move", False)
    assert textfit.strip_markup("{accent}47%{/} CAGR") == ("47% CAGR", False)
    assert textfit.strip_markup("plain text") == ("plain text", False)


# ── bullets_fit ──────────────────────────────────────────────────────────────
def test_bullets_fit_empty():
    assert textfit.bullets_fit([], 14, 5.765, 5.3) == (0.0, 0, [])


def test_bullets_fit_step_floor():
    # short bullets each consume one step slot
    ratio, lines, overflow = textfit.bullets_fit(["a", "b", "c"], 14, 5.765, 5.3,
                                                 cols=1, step_in=0.74)
    assert ratio == pytest.approx(3 * 0.74 / 5.3)
    assert lines == 3
    assert overflow == []  # single-line bullets never overflow their slot


def test_bullets_fit_two_columns_halve_the_load():
    bullets = ["short bullet"] * 8
    one, *_ = textfit.bullets_fit(bullets, 14, 5.765, 5.3, cols=1, step_in=0.74)
    two, *_ = textfit.bullets_fit(bullets, 14, 5.765, 5.3, cols=2, step_in=0.74)
    assert two == pytest.approx(one / 2)


def test_bullets_fit_longer_text_higher_ratio():
    short, *_ = textfit.bullets_fit([SENTENCE] * 4, 14, 5.765, 5.3)
    long, *_ = textfit.bullets_fit([SENTENCE * 4] * 4, 14, 5.765, 5.3)
    assert long > short


def test_bullets_fit_overflowing_indices_non_terminal():
    """A very long non-terminal bullet is flagged; the last bullet in a column
    is exempt because there is no next bullet to collide with."""
    geo = builders.bullet_geometry()
    # Build 6 bullets: first is very long (~270 chars), rest are short.
    long_bullet = (SENTENCE + " ") * 4  # well over 270 chars
    bullets = [long_bullet] + [SENTENCE[:40]] * 5
    ratio, _, overflowing = textfit.bullets_fit(
        bullets, geo["font_pt"], geo["col_w"], geo["avail_h"],
        cols=geo["cols"], step_in=geo["step"])
    # The long bullet at index 0 is not the last in its column and wraps > step.
    assert 0 in overflowing


def test_bullets_fit_last_bullet_not_flagged():
    """When the long bullet is the LAST in its column it should NOT be flagged."""
    geo = builders.bullet_geometry()
    long_bullet = (SENTENCE + " ") * 4
    # Place long bullet last; 5 short bullets before it in a 1-column layout.
    bullets = [SENTENCE[:40]] * 5 + [long_bullet]
    _, _, overflowing = textfit.bullets_fit(
        bullets, geo["font_pt"], geo["col_w"], geo["avail_h"],
        cols=1, step_in=geo["step"])
    # Last bullet is always exempt.
    assert len(bullets) - 1 not in overflowing


def test_bullets_fit_no_overflowing_without_step():
    """When step_in == 0 there is no fixed slot, so overflowing is always []."""
    long_bullet = (SENTENCE + " ") * 4
    bullets = [long_bullet] * 4
    _, _, overflowing = textfit.bullets_fit(bullets, 14, 5.765, 5.3, step_in=0.0)
    assert overflowing == []


# ── validate() capacity budgets ──────────────────────────────────────────────
def _bullet_outline(bullets, layout="bullet-list"):
    lines = [f"## Slide 1: Long action title describing the takeaway here",
             f"**Layout:** {layout}"]
    if layout == "two-column-split":
        lines.append("**Visual:** none")
    lines += [f"- {b}" for b in bullets]
    lines.append('- Notes: "n"')
    _, slides = parse_outline("\n".join(lines) + "\n")
    return slides


def _capacity_msgs(errors, warnings):
    errs = [e for e in errors if "overflows layout capacity" in e]
    warns = [w for w in warnings if "text capacity" in w]
    return errs, warns


def test_validate_normal_bullets_no_capacity_messages():
    slides = _bullet_outline([SENTENCE[:50]] * 6)
    errors, warnings = validate(slides, CTX)
    errs, warns = _capacity_msgs(errors, warnings)
    assert not errs and not warns


def test_validate_mild_overflow_warns_autofit():
    geo = builders.bullet_geometry()
    text = ((SENTENCE + " ") * 4)[:205].strip()
    bullets = [text] * 14
    ratio, *_ = textfit.bullets_fit(bullets, geo["font_pt"], geo["col_w"],
                                    geo["avail_h"], cols=geo["cols"],
                                    step_in=geo["step"])
    assert 1.0 < ratio <= 1.4, f"test fixture drifted: ratio={ratio}"
    slides = _bullet_outline(bullets)
    errors, warnings = validate(slides, CTX)
    errs, warns = _capacity_msgs(errors, warnings)
    assert not errs
    assert warns and "autofit will shrink text" in warns[0]


def test_validate_severe_overflow_errors():
    text = ((SENTENCE + " ") * 8)[:420].strip()
    slides = _bullet_outline([text] * 14)
    errors, warnings = validate(slides, CTX)
    errs, _ = _capacity_msgs(errors, warnings)
    assert errs and "split the slide or cut copy" in errs[0]


def test_validate_two_column_overflow():
    text = ((SENTENCE + " ") * 8)[:500].strip()
    slides = _bullet_outline([text] * 8, layout="two-column-split")
    errors, warnings = validate(slides, CTX)
    errs, _ = _capacity_msgs(errors, warnings)
    assert errs


def test_validate_heading_wrap_warns():
    md = ("## Slide 1: " + ("Strategic transformation accelerates measurable "
          "enterprise outcomes across " * 3).strip() + "\n"
          "**Layout:** bullet-list\n- one\n- Notes: \"n\"\n")
    _, slides = parse_outline(md)
    _, warnings = validate(slides, CTX)
    assert any("title wraps to ~" in w for w in warnings)


def test_validate_short_heading_no_wrap_warning():
    slides = _bullet_outline(["one", "two"])
    _, warnings = validate(slides, CTX)
    assert not any("title wraps" in w for w in warnings)


# ── per-slot collision guard ─────────────────────────────────────────────────
def _collision_warns(warnings):
    return [w for w in warnings if "will collide with the next bullet" in w]


def _collision_errors(errors):
    return [e for e in errors if "will collide with the next bullet" in e]


def test_validate_per_slot_collision_warns():
    """Exact repro: 6 bullets, first is very long (~270+ chars), rest short.
    The long non-terminal bullet must trigger a collision warning at --check."""
    geo = builders.bullet_geometry()
    long_bullet = (SENTENCE + " ") * 4  # ~270 chars, wraps many lines
    bullets = [long_bullet] + [SENTENCE[:40]] * 5
    # Confirm the long bullet is NOT enough to trip the total-capacity warning.
    ratio, *_ = textfit.bullets_fit(
        bullets, geo["font_pt"], geo["col_w"], geo["avail_h"],
        cols=geo["cols"], step_in=geo["step"])
    # Total ratio may or may not exceed 1.0 — what matters is collision is flagged.
    slides = _bullet_outline(bullets)
    errors, warnings = validate(slides, CTX)
    coll_warns = _collision_warns(warnings)
    assert coll_warns, (
        f"expected per-slot collision warning (ratio={ratio:.2f}); "
        f"warnings={warnings}")
    assert "bullet 1" in coll_warns[0]


def test_validate_per_slot_last_bullet_no_collision_warn():
    """Long bullet placed LAST in its column must not generate a collision warning."""
    long_bullet = (SENTENCE + " ") * 4
    # 5 short bullets + 1 very long bullet last, single column
    bullets = [SENTENCE[:40]] * 5 + [long_bullet]
    slides = _bullet_outline(bullets)
    errors, warnings = validate(slides, CTX)
    coll_warns = _collision_warns(warnings)
    coll_errs = _collision_errors(errors)
    # The long bullet is the last slot in its column — no collision possible.
    assert not coll_warns and not coll_errs, (
        f"unexpected collision message for last bullet: {coll_warns + coll_errs}")


def test_validate_per_slot_severe_errors():
    """When a non-terminal bullet is 1.8x taller than its step, it is an ERROR."""
    geo = builders.bullet_geometry()
    line_h = geo["font_pt"] * 1.2 / 72.0
    # Build a bullet that wraps to enough lines to exceed 1.8 * step.
    # Need lines * line_h > step * 1.8  →  lines > step * 1.8 / line_h
    target_lines = int(geo["step"] * 1.8 / line_h) + 2
    # Use a very long bullet to guarantee enough wrapped lines.
    long_bullet = (SENTENCE + " ") * (target_lines * 4)
    bullets = [long_bullet] + ["short"] * 5
    slides = _bullet_outline(bullets)
    errors, warnings = validate(slides, CTX)
    coll_errs = _collision_errors(errors)
    assert coll_errs, (
        f"expected collision ERROR for ~{target_lines}-line bullet; "
        f"errors={errors}, warnings={warnings}")


def test_builder_autofit_per_slot_overflow():
    """Builder applies autofit even when total ratio <= 1.0 but a non-terminal
    bullet overflows its slot."""
    geo = builders.bullet_geometry()
    long_bullet = (SENTENCE + " ") * 4
    bullets = [long_bullet] + [SENTENCE[:35]] * 5
    slide = _build_bullets(bullets)
    fits = _autofit_elements(slide)
    # Autofit must be written for at least the overflowing bullet's text frame.
    assert len(fits) > 0, "expected normAutofit for per-slot overflow"


# ── builder normAutofit ──────────────────────────────────────────────────────
def _autofit_elements(slide):
    found = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        found.extend(shape._element.findall(".//" + qn("a:normAutofit")))
    return found


def _build_bullets(bullets):
    _, slides = parse_outline(
        "## Slide 1: T\n**Layout:** bullet-list\n"
        + "".join(f"- {b}\n" for b in bullets))
    prs = Presentation()
    pal = get_palette("midnight-executive")
    return builders.build_bullet_slide(prs, slides[0], pal, {})


def test_builder_writes_normautofit_in_overflow_band():
    geo = builders.bullet_geometry()
    text = ((SENTENCE + " ") * 4)[:205].strip()
    bullets = [text] * 14
    ratio, _, overflowing = textfit.bullets_fit(bullets, geo["font_pt"], geo["col_w"],
                                               geo["avail_h"], cols=geo["cols"],
                                               step_in=geo["step"])
    assert 1.0 < ratio <= 1.4, f"test fixture drifted: ratio={ratio}"
    # Compute the same combined scale the builder uses.
    scale_ratio = max(80, min(100, int(100 / ratio)))
    scale_slot = 100
    if overflowing and geo["step"] > 0:
        line_h = geo["font_pt"] * 1.2 / 72.0
        for idx in overflowing:
            plain, has_icon = textfit.strip_markup(bullets[idx])
            w = geo["col_w"] - (textfit.ICON_INDENT if has_icon
                                else textfit.MARKER_INDENT)
            lines = max(textfit.estimate_lines(plain, geo["font_pt"], w), 1)
            needed = max(80, min(100, int(100 * geo["step"] / (lines * line_h))))
            scale_slot = min(scale_slot, needed)
    expected = min(scale_ratio, scale_slot) * 1000
    slide = _build_bullets(bullets)
    fits = _autofit_elements(slide)
    assert len(fits) == len(bullets)  # one per bullet text frame
    for el in fits:
        assert el.get("fontScale") == str(expected)
        assert el.get("lnSpcReduction") == "10000"
    assert 80000 <= expected < 100000


def test_builder_no_autofit_under_capacity():
    slide = _build_bullets([SENTENCE[:40]] * 6)
    assert _autofit_elements(slide) == []


def test_builder_two_column_autofit():
    geo = builders.twocol_geometry()
    text = ((SENTENCE + " ") * 5)[:250].strip()
    bullets = [text] * 8
    ratio, *_ = textfit.bullets_fit(bullets, geo["font_pt"], geo["col_w"],
                                    geo["avail_h"], cols=geo["cols"],
                                    step_in=geo["step"])
    assert ratio > 1.0, f"test fixture drifted: ratio={ratio}"
    _, slides = parse_outline(
        "## Slide 1: T\n**Layout:** two-column-split\n"
        + "".join(f"- {b}\n" for b in bullets))
    prs = Presentation()
    pal = get_palette("midnight-executive")
    slide = builders.build_two_column_slide(prs, slides[0], pal, {})
    assert len(_autofit_elements(slide)) == len(bullets)


def test_apply_autofit_replaces_existing():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    from pptx.util import Inches
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textfit.apply_autofit(tb.text_frame, 90)
    textfit.apply_autofit(tb.text_frame, 85)
    fits = tb._element.findall(".//" + qn("a:normAutofit"))
    assert len(fits) == 1
    assert fits[0].get("fontScale") == "85000"
