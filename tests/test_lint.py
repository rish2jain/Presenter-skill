"""Tests for pptx_lint.py cross-slide consistency checks."""
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))

from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

from pptx_lint import lint_deck  # noqa: E402


def _prs():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
    return prs


def _blank(prs):
    layout = min(prs.slide_layouts, key=lambda l: len(l.placeholders))
    return prs.slides.add_slide(layout)


def _tb(slide, text, left, top, w=0.8, h=0.35, size=11, color=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size)
    if color:
        from pptx.dml.color import RGBColor
        run.font.color.rgb = RGBColor.from_string(color)
    return tb


def test_jiggle_flagged_for_misaligned_page_numbers():
    prs = _prs()
    for n, left in enumerate((11.9, 11.9, 11.4)):  # third one jiggles
        _tb(_blank(prs), str(n + 1), left, 7.08)
    issues = lint_deck(prs)
    assert any("jiggle" in e for e in issues["error"]), issues


def test_aligned_page_numbers_pass():
    prs = _prs()
    for n in range(3):
        _tb(_blank(prs), str(n + 1), 11.9, 7.08)
    issues = lint_deck(prs)
    assert not any("jiggle" in e for e in issues["error"]), issues


def test_page_sequence_gap_flagged():
    prs = _prs()
    for label in ("1", "2", "4"):
        _tb(_blank(prs), label, 11.9, 7.08)
    issues = lint_deck(prs)
    assert any("sequence" in e for e in issues["error"]), issues


def test_font_explosion_warned():
    prs = _prs()
    slide = _blank(prs)
    for i, fname in enumerate(
            ("Calibri", "Arial", "Georgia", "Verdana", "Impact", "Tahoma")):
        tb = _tb(slide, f"text {i}", 1.0, 1.0 + i * 0.5)
        tb.text_frame.paragraphs[0].runs[0].font.name = fname
    issues = lint_deck(prs)
    assert any("font" in w.lower() for w in issues["warn"]), issues


def test_palette_whitelist_flags_off_palette_color():
    prs = _prs()
    slide = _blank(prs)
    _tb(slide, "on palette", 1.0, 1.0, color="C9A84C")   # midnight accent1
    _tb(slide, "rogue", 1.0, 2.0, color="FF00FF")
    issues = lint_deck(prs, palette_key="midnight-executive")
    assert any("FF00FF" in e for e in issues["error"]), issues
    assert not any("C9A84C" in e for e in issues["error"]), issues


def test_missing_font_warned(monkeypatch):
    import pptx_lint
    monkeypatch.setattr(pptx_lint, "installed_fonts",
                        lambda: {"calibri", "arial"})
    prs = _prs()
    slide = _blank(prs)
    tb = _tb(slide, "hello", 1.0, 1.0)
    tb.text_frame.paragraphs[0].runs[0].font.name = "Gill Sans MT"
    issues = pptx_lint.lint_deck(prs)
    assert any("Gill Sans MT" in w for w in issues["warn"]), issues


def test_font_check_skipped_when_inventory_unavailable(monkeypatch):
    import pptx_lint
    monkeypatch.setattr(pptx_lint, "installed_fonts", lambda: None)
    prs = _prs()
    _tb(_blank(prs), "hello", 1.0, 1.0)
    issues = pptx_lint.lint_deck(prs)
    assert not any("not installed" in w for w in issues["warn"]), issues


def test_installed_fonts_returns_none_or_lowercase_set():
    from pptx_lint import installed_fonts
    fonts = installed_fonts()
    assert fonts is None or (
        isinstance(fonts, set) and all(f == f.lower() for f in fonts))


def test_jiggle_blames_outlier_not_majority():
    """Slide 1 is the outlier; only slide 1 should be blamed."""
    prs = _prs()
    # slide 1: left=11.4 (outlier); slides 2-4: left=11.9 (majority)
    for i, left in enumerate((11.4, 11.9, 11.9, 11.9), start=1):
        _tb(_blank(prs), str(i), left, 7.08)
    issues = lint_deck(prs)
    jiggle_errors = [e for e in issues["error"] if "jiggle" in e]
    assert len(jiggle_errors) == 1, jiggle_errors
    assert "Slide 1" in jiggle_errors[0], jiggle_errors[0]


def test_jiggle_no_false_positive_when_all_aligned():
    """Four perfectly aligned page numbers must produce zero jiggle errors."""
    prs = _prs()
    for i in range(1, 5):
        _tb(_blank(prs), str(i), 11.9, 7.08)
    issues = lint_deck(prs)
    assert not any("jiggle" in e for e in issues["error"]), issues


def test_page_gap_over_unnumbered_slide_passes():
    """A gap in page numbers spanning an unnumbered divider slide must not be flagged.

    Slides 1,2,4 carry numbers; slide 3 is an unnumbered layout.
    Number delta (4-2=2) equals slide-position delta (4-2=2), so no error.
    """
    prs = _prs()
    _tb(_blank(prs), "1", 11.9, 7.08)
    _tb(_blank(prs), "2", 11.9, 7.08)
    _blank(prs)                      # divider: no page number
    _tb(_blank(prs), "4", 11.9, 7.08)
    issues = lint_deck(prs)
    assert not any("sequence" in e for e in issues["error"]), issues


def test_lint_with_custom_palette(tmp_path):
    import json
    from palettes import PALETTES, load_custom_palettes
    pdir = tmp_path / "palettes"
    pdir.mkdir()
    (pdir / "lintbrand.json").write_text(json.dumps({
        "bg": "101820", "bg_deep": "0A0F14", "surface": "1E2A33",
        "accent1": "FEE715", "accent2": "8DA9C4", "accent3": "5C946E",
        "text": "F4F4F4", "text_muted": "9DB2BF", "dark": True}))
    load_custom_palettes(pdir)
    try:
        prs = _prs()
        _tb(_blank(prs), "branded", 1.0, 1.0, color="FEE715")
        issues = lint_deck(prs, palette_key="lintbrand")
        assert not any("unknown palette" in w for w in issues["warn"]), issues
        assert not any("FEE715" in e for e in issues["error"]), issues
    finally:
        PALETTES.pop("lintbrand", None)


def test_lint_tolerates_malformed_chart_series():
    from palettes import PALETTES
    PALETTES["malformed"] = {
        **PALETTES["midnight-executive"],
        "chart_series": [123, None, "FEE715", "not-hex"],
    }
    try:
        prs = _prs()
        _tb(_blank(prs), "branded", 1.0, 1.0, color="FEE715")
        issues = lint_deck(prs, palette_key="malformed")
        assert not any("unknown palette" in w for w in issues["warn"]), issues
        assert not any("FEE715" in e for e in issues["error"]), issues
        assert any("chart_series" in w for w in issues["warn"]), issues
    finally:
        PALETTES.pop("malformed", None)


def test_lint_tolerates_non_list_chart_series():
    from palettes import PALETTES
    PALETTES["badseries"] = {**PALETTES["midnight-executive"], "chart_series": 42}
    try:
        issues = lint_deck(_prs(), palette_key="badseries")
        assert any("chart_series is not a list" in w for w in issues["warn"]), issues
    finally:
        PALETTES.pop("badseries", None)


# ── AI-tell checks ───────────────────────────────────────────────────────────
def _rect(slide, left, top, w, h):
    from pptx.enum.shapes import MSO_SHAPE
    return slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))


def test_accent_line_under_title_flagged():
    prs = _prs()
    slide = _blank(prs)
    _tb(slide, "Strong takeaway title", 0.7, 0.5, w=8.0, h=0.6, size=28)
    _rect(slide, 0.7, 1.15, 3.0, 0.05)  # thin bar just under title bottom
    issues = lint_deck(prs)
    assert any("accent line under title" in w for w in issues["warn"]), issues


def test_thin_bar_far_from_title_passes():
    prs = _prs()
    slide = _blank(prs)
    _tb(slide, "Strong takeaway title", 0.7, 0.5, w=8.0, h=0.6, size=28)
    _rect(slide, 0.7, 4.0, 3.0, 0.05)  # divider mid-slide, not under title
    issues = lint_deck(prs)
    assert not any("accent line" in w for w in issues["warn"]), issues


def _centered_body_slide(prs, align, n_extra=3):
    from pptx.enum.text import PP_ALIGN
    slide = _blank(prs)
    _tb(slide, "Title", 0.7, 0.3, w=8.0, h=0.6, size=28)
    for i in range(n_extra - 1):
        _tb(slide, f"label {i}", 0.7 + i * 3, 1.0)
    tb = _tb(slide, "Centered paragraphs of body copy are a hallmark of "
             "AI-generated slides and hurt scanability for readers.",
             2.0, 2.0, w=9.0, h=1.5, size=14)
    tb.text_frame.paragraphs[0].alignment = align
    return slide


def test_long_centered_body_flagged():
    from pptx.enum.text import PP_ALIGN
    prs = _prs()
    _centered_body_slide(prs, PP_ALIGN.CENTER)
    issues = lint_deck(prs)
    assert any("centered body" in w for w in issues["warn"]), issues


def test_left_aligned_body_passes():
    from pptx.enum.text import PP_ALIGN
    prs = _prs()
    _centered_body_slide(prs, PP_ALIGN.LEFT)
    issues = lint_deck(prs)
    assert not any("centered body" in w for w in issues["warn"]), issues


def test_centered_text_on_sparse_slide_passes():
    """<=3 shapes = quote/big-number territory; centered is intentional."""
    from pptx.enum.text import PP_ALIGN
    prs = _prs()
    slide = _blank(prs)
    tb = _tb(slide, "A pull quote that is deliberately centered and long "
             "enough to clear the eighty character threshold easily.",
             2.0, 2.5, w=9.0, h=1.5, size=20)
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    issues = lint_deck(prs)
    assert not any("centered body" in w for w in issues["warn"]), issues


def test_placeholder_rx_catches_ai_tell_text():
    from qa_check import PLACEHOLDER_RX
    assert PLACEHOLDER_RX.search("This slide layout works for comparisons")
    assert PLACEHOLDER_RX.search("This page layout is versatile")
    assert PLACEHOLDER_RX.search("Lorem Ipsum dolor sit amet")
    assert not PLACEHOLDER_RX.search("Margins expanded 400bps")


# ── value-axis scale comparison ──────────────────────────────────────────────
def _chart_slide(prs, values, max_scale=None):
    from palettes import get_palette
    from charts import add_native_chart
    slide = _blank(prs)
    chart = add_native_chart(slide, get_palette("midnight-executive"),
                             "column", ["A", "B"], values, 1.0, 1.0, 6.0, 4.0)
    if max_scale is not None:
        chart.value_axis.maximum_scale = float(max_scale)
    return chart


def test_axis_scale_mismatch_warned():
    prs = _prs()
    _chart_slide(prs, [3, 5], max_scale=10)
    _chart_slide(prs, [300, 500], max_scale=500)
    issues = lint_deck(prs)
    assert any("dishonest scale" in w for w in issues["warn"]), issues


def test_axis_scale_same_magnitude_passes():
    prs = _prs()
    _chart_slide(prs, [3, 5], max_scale=40)
    _chart_slide(prs, [30, 44], max_scale=50)
    issues = lint_deck(prs)
    assert not any("dishonest scale" in w for w in issues["warn"]), issues


def test_axis_scale_auto_axes_ignored():
    prs = _prs()
    _chart_slide(prs, [3, 5], max_scale=10)
    _chart_slide(prs, [300, 500])  # auto axis -> not comparable
    issues = lint_deck(prs)
    assert not any("dishonest scale" in w for w in issues["warn"]), issues


# ── Google Slides compatibility (--gslides) ──────────────────────────────────
def test_gslides_unsafe_font_warned_only_with_flag():
    prs = _prs()
    slide = _blank(prs)
    tb = _tb(slide, "hello", 1.0, 1.0)
    tb.text_frame.paragraphs[0].runs[0].font.name = "Gill Sans MT"
    issues = lint_deck(prs, gslides=True)
    assert any("gslides" in w and "Gill Sans MT" in w
               for w in issues["warn"]), issues
    issues = lint_deck(prs)  # flag off -> no gslides output
    assert not any("gslides" in w for w in issues["warn"]), issues


def test_gslides_safe_font_passes():
    prs = _prs()
    slide = _blank(prs)
    tb = _tb(slide, "hello", 1.0, 1.0)
    tb.text_frame.paragraphs[0].runs[0].font.name = "Arial"
    issues = lint_deck(prs, gslides=True)
    assert not any("gslides: font" in w for w in issues["warn"]), issues


def test_gslides_nonfade_transition_warned():
    from lxml import etree
    from pptx.oxml.ns import qn
    prs = _prs()
    slide = _blank(prs)
    trans = etree.SubElement(slide.element, qn("p:transition"))
    etree.SubElement(trans, qn("p:wipe"))
    issues = lint_deck(prs, gslides=True)
    assert any("transition 'wipe'" in w for w in issues["warn"]), issues


def test_gslides_fade_transition_passes():
    from lxml import etree
    from pptx.oxml.ns import qn
    prs = _prs()
    slide = _blank(prs)
    trans = etree.SubElement(slide.element, qn("p:transition"))
    etree.SubElement(trans, qn("p:fade"))
    issues = lint_deck(prs, gslides=True)
    assert not any("transition" in w for w in issues["warn"]), issues


def test_gslides_smartart_warned():
    from lxml import etree
    prs = _prs()
    slide = _blank(prs)
    etree.SubElement(
        slide.shapes._spTree,
        "{http://schemas.openxmlformats.org/drawingml/2006/diagram}relIds")
    issues = lint_deck(prs, gslides=True)
    assert any("SmartArt" in w for w in issues["warn"]), issues


def test_gslides_embedded_media_warned(tmp_path):
    from pptx.util import Inches
    movie = tmp_path / "clip.mp4"
    movie.write_bytes(b"fake video bytes")
    prs = _prs()
    slide = _blank(prs)
    slide.shapes.add_movie(str(movie), Inches(1), Inches(1),
                           Inches(2), Inches(2), mime_type="video/mp4")
    issues = lint_deck(prs, gslides=True)
    assert any("embedded media" in w for w in issues["warn"]), issues


# ── accessibility extensions (qa_check --accessibility) ─────────────────────
def _qa(prs, tmp_path, accessibility=True):
    from qa_check import check_deck
    p = tmp_path / "deck.pptx"
    prs.save(p)
    return check_deck(p, accessibility=accessibility)


def test_a11y_missing_title_escalates_to_error(tmp_path):
    prs = _prs()
    _rect(_blank(prs), 1.0, 1.0, 2.0, 2.0)  # no text anywhere
    issues = _qa(prs, tmp_path)
    assert any("no detectable slide title" in e for e in issues["error"]), issues
    issues = _qa(prs, tmp_path, accessibility=False)
    assert any("no detectable slide title" in w for w in issues["warn"]), issues
    assert not any("no detectable slide title" in e for e in issues["error"])


def test_a11y_duplicate_titles_one_consolidated_issue(tmp_path):
    prs = _prs()
    for text in ("Revenue Growth", "  revenue growth"):  # case-folded match
        _tb(_blank(prs), text, 0.7, 0.5, w=8.0, h=0.8, size=24)
    issues = _qa(prs, tmp_path)
    dups = [e for e in issues["error"] if "duplicate slide title" in e]
    assert len(dups) == 1, issues
    assert "slides 1, 2" in dups[0], dups[0]


def _table_slide(prs, rows=2, cols=2):
    from pptx.util import Inches
    slide = _blank(prs)
    return slide.shapes.add_table(rows, cols, Inches(1), Inches(1),
                                  Inches(4), Inches(2))


def test_a11y_table_without_header_row_error(tmp_path):
    prs = _prs()
    gf = _table_slide(prs)
    del gf.table._tbl.tblPr.attrib["firstRow"]
    issues = _qa(prs, tmp_path)
    assert any("header row" in e for e in issues["error"]), issues


def test_a11y_table_with_header_row_passes(tmp_path):
    prs = _prs()
    _table_slide(prs)  # python-pptx default sets firstRow="1"
    issues = _qa(prs, tmp_path)
    assert not any("header row" in e for e in issues["error"]), issues


def test_a11y_merged_cells_warned(tmp_path):
    prs = _prs()
    gf = _table_slide(prs)
    gf.table.cell(0, 0).merge(gf.table.cell(0, 1))
    issues = _qa(prs, tmp_path)
    assert any("merged cells" in w for w in issues["warn"]), issues


def test_a11y_table_checks_off_in_default_mode(tmp_path):
    prs = _prs()
    gf = _table_slide(prs)
    del gf.table._tbl.tblPr.attrib["firstRow"]
    gf.table.cell(0, 0).merge(gf.table.cell(0, 1))
    issues = _qa(prs, tmp_path, accessibility=False)
    assert not any("header row" in m or "merged cells" in m
                   for m in issues["error"] + issues["warn"]), issues


def _picture_slide(prs, tmp_path, descr):
    from PIL import Image
    from pptx.util import Inches
    img = tmp_path / "pic.png"
    Image.new("RGB", (8, 8), (40, 90, 160)).save(img)
    slide = _blank(prs)
    pic = slide.shapes.add_picture(str(img), Inches(1), Inches(1))
    pic._element._nvXxPr.cNvPr.set("descr", descr)
    return slide


def test_a11y_filename_alt_text_flagged(tmp_path):
    prs = _prs()
    _picture_slide(prs, tmp_path, "hero-shot.png")
    issues = _qa(prs, tmp_path)
    assert any("alt text is a filename" in e for e in issues["error"]), issues


def test_a11y_descriptive_alt_text_passes(tmp_path):
    prs = _prs()
    _picture_slide(prs, tmp_path, "Team on stage at the product launch")
    issues = _qa(prs, tmp_path)
    assert not any("alt text is a filename" in e for e in issues["error"])


def test_a11y_reading_order_title_not_first_warned(tmp_path):
    prs = _prs()
    slide = _blank(prs)
    _tb(slide, "body copy appears first", 0.7, 2.0, w=6.0, h=0.6, size=20)
    _tb(slide, "The Actual Title", 0.7, 0.5, w=8.0, h=0.8, size=30)
    issues = _qa(prs, tmp_path)
    assert any("reading order" in w for w in issues["warn"]), issues


def test_a11y_reading_order_title_first_passes(tmp_path):
    prs = _prs()
    slide = _blank(prs)
    _tb(slide, "The Actual Title", 0.7, 0.5, w=8.0, h=0.8, size=30)
    _tb(slide, "body copy after the title", 0.7, 2.0, w=6.0, h=0.6, size=20)
    issues = _qa(prs, tmp_path)
    assert not any("reading order" in w for w in issues["warn"]), issues
