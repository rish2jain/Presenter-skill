"""Tests for chart annotations: waterfall difference bracket, CAGR arrow."""
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))

from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

import builders  # noqa: E402
from build_deck import parse_outline  # noqa: E402
from builders_consulting import build_waterfall_slide  # noqa: E402
from palettes import get_palette  # noqa: E402

PAL = get_palette("midnight-executive")
CTX = {"outline_dir": Path("."), "assets_dir": Path("assets")}


def _prs():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
    builders.set_canvas(prs)
    return prs


def _texts(slide):
    return [sh.text_frame.text for sh in slide.shapes
            if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip()]


def _cagr_labels(slide):
    return [t for t in _texts(slide) if t.strip().startswith("CAGR ")]


BRACKET_MD = """## Slide 1: Three levers bridge run-rate down to target
**Layout:** waterfall
**Data:**
- FY25: 46
- Tiering: -8
- Exit: -6
- Discounts: -4
- FY27: total
- Bracket: FY25, FY27
"""


def test_bracket_directive_parses():
    _, slides = parse_outline(BRACKET_MD)
    assert slides[0]["bracket"] == "FY25, FY27"


def test_bracket_renders_auto_pct_label():
    _, slides = parse_outline(BRACKET_MD)
    slide = build_waterfall_slide(_prs(), slides[0], PAL, CTX)
    texts = _texts(slide)
    # 46 -> 28 is -39%
    assert any("-39%" in t for t in texts), texts


def test_bracket_custom_label():
    md = BRACKET_MD.replace('- Bracket: FY25, FY27',
                            '- Bracket: FY25, FY27, "Run-rate reset"')
    _, slides = parse_outline(md)
    slide = build_waterfall_slide(_prs(), slides[0], PAL, CTX)
    assert any("Run-rate reset" in t for t in _texts(slide))


CAGR_MD = """## Slide 1: Revenue compounds at double digits
**Layout:** two-column-split
**Visual:** chart:bar
- CAGR: on
**Data:**
- 2021: 10
- 2022: 13
- 2023: 17
- 2024: 22
- Strong compounding story
"""


def test_cagr_arrow_label_on_chart_slide():
    _, slides = parse_outline(CAGR_MD)
    assert slides[0]["cagr"] == "on"
    prs = _prs()
    slide = builders.LAYOUT_MAP["two-column-split"](prs, slides[0], PAL, CTX)
    # (22/10)^(1/3)-1 = 30.1%
    assert any("CAGR" in t and "30" in t for t in _texts(slide)), _texts(slide)


AXIS_MD = """## Slide 1: EMEA revenue lags at identical scale
**Layout:** two-column-split
**Visual:** chart:bar
- Axis-Max: 50
**Data:**
- 2023: 17
- 2024: 22
- Same scale as the Americas chart
"""


def test_axis_max_pins_value_axis():
    _, slides = parse_outline(AXIS_MD)
    assert slides[0]["axis_max"] == "50"
    prs = _prs()
    slide = builders.LAYOUT_MAP["two-column-split"](prs, slides[0], PAL, CTX)
    chart = next(sh.chart for sh in slide.shapes if getattr(sh, "has_chart", False))
    assert chart.value_axis.maximum_scale == 50.0


CAGR_AXIS_MD = """## Slide 1: Revenue compounds at double digits on a fixed scale
**Layout:** two-column-split
**Visual:** chart:bar
- CAGR: on
- Axis-Max: 50
**Data:**
- 2021: 10
- 2022: 13
- 2023: 17
- 2024: 22
- Same scale as sibling slides
"""


def test_cagr_arrow_respects_axis_max():
    _, slides = parse_outline(CAGR_AXIS_MD)
    prs = _prs()
    slide = builders.LAYOUT_MAP["two-column-split"](prs, slides[0], PAL, CTX)
    chart = next(sh.chart for sh in slide.shapes if getattr(sh, "has_chart", False))
    assert chart.value_axis.maximum_scale == 50.0
    assert any("CAGR" in t for t in _texts(slide))


def test_cagr_invalid_axis_max_warns(capsys):
    md = CAGR_AXIS_MD.replace("- Axis-Max: 50", "- Axis-Max: fifty")
    _, slides = parse_outline(md)
    prs = _prs()
    builders.LAYOUT_MAP["two-column-split"](prs, slides[0], PAL, CTX)
    err = capsys.readouterr().err
    assert "Axis-Max not numeric" in err and "fifty" in err and "CAGR" in err


def test_cagr_skips_with_fewer_than_two_data_points():
    md = """## Slide 1: One year only
**Layout:** two-column-split
**Visual:** chart:bar
- CAGR: on
**Data:**
- 2024: 22
- Not enough history for compounding
"""
    _, slides = parse_outline(md)
    prs = _prs()
    slide = builders.LAYOUT_MAP["two-column-split"](prs, slides[0], PAL, CTX)
    assert not _cagr_labels(slide), _texts(slide)


def test_cagr_skips_with_non_positive_start_or_end():
    md = CAGR_MD.replace("- 2021: 10", "- 2021: 0")
    _, slides = parse_outline(md)
    prs = _prs()
    slide = builders.LAYOUT_MAP["two-column-split"](prs, slides[0], PAL, CTX)
    assert not _cagr_labels(slide), _texts(slide)

    md_end = CAGR_MD.replace("- 2024: 22", "- 2024: -1")
    _, slides_end = parse_outline(md_end)
    slide_end = builders.LAYOUT_MAP["two-column-split"](_prs(), slides_end[0], PAL, CTX)
    assert not _cagr_labels(slide_end), _texts(slide_end)


def test_cagr_label_on_declining_series():
    md = CAGR_MD.replace("- 2021: 10", "- 2021: 22").replace("- 2024: 22",
                                                             "- 2024: 10")
    _, slides = parse_outline(md)
    prs = _prs()
    slide = builders.LAYOUT_MAP["two-column-split"](prs, slides[0], PAL, CTX)
    # (10/22)^(1/3)-1 = -23.1%
    assert any("CAGR" in t and "-23" in t for t in _texts(slide)), _texts(slide)


def test_add_arrowhead_to_connector_adds_tail_end():
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn
    from charts import add_arrowhead_to_connector

    prs = _prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(3), Inches(2))
    assert add_arrowhead_to_connector(conn) is True
    ln = conn.line._get_or_add_ln()
    tail = ln.find(qn("a:tailEnd"))
    assert tail is not None
    assert tail.get("type") == "triangle"


def test_add_arrowhead_to_connector_noop_without_private_api(caplog):
    from types import SimpleNamespace
    from charts import add_arrowhead_to_connector

    fake_line = SimpleNamespace()  # no _get_or_add_ln
    conn = SimpleNamespace(line=fake_line)
    assert add_arrowhead_to_connector(conn, warn=True) is False
    assert any("Skipping connector arrowhead" in r.message for r in caplog.records)


# ── largest-remainder rounding ───────────────────────────────────────────────
def test_round_to_sum_largest_remainder():
    from charts import round_to_sum
    assert round_to_sum([33.33, 33.33, 33.33]) == [33, 33, 34]


def test_round_to_sum_preserves_order():
    from charts import round_to_sum
    assert round_to_sum([20.5, 60.2, 19.3]) == [21, 60, 19]


def test_round_to_sum_negative_passthrough():
    from charts import round_to_sum
    assert round_to_sum([-10.4, 110.6]) == [-10, 111]


def test_round_to_sum_decimals():
    from charts import round_to_sum
    assert round_to_sum([33.333, 33.333, 33.334], decimals=1) == \
        [33.3, 33.3, 33.4]


# ── value line annotation ────────────────────────────────────────────────────
VLINE_MD = """## Slide 1: Pipeline coverage clears the 40-unit target
**Layout:** two-column-split
**Visual:** chart:bar
- Value-Line: Target, 40
**Data:**
- 2023: 30
- 2024: 44
- Coverage improves every quarter
"""


def test_value_line_parses_and_renders_label():
    _, slides = parse_outline(VLINE_MD)
    assert slides[0]["value_line"] == "Target, 40"
    slide = builders.LAYOUT_MAP["two-column-split"](_prs(), slides[0], PAL, CTX)
    assert any(t == "Target" for t in _texts(slide)), _texts(slide)


def test_value_line_outside_axis_range_warns_and_skips(capsys):
    md = VLINE_MD.replace("Target, 40", "Target, 400")
    _, slides = parse_outline(md)
    slide = builders.LAYOUT_MAP["two-column-split"](_prs(), slides[0], PAL, CTX)
    err = capsys.readouterr().err
    assert "Value-Line" in err
    assert not any(t == "Target" for t in _texts(slide)), _texts(slide)


def test_value_line_unparsable_warns(capsys):
    md = VLINE_MD.replace("Target, 40", "Target, forty")
    _, slides = parse_outline(md)
    slide = builders.LAYOUT_MAP["two-column-split"](_prs(), slides[0], PAL, CTX)
    err = capsys.readouterr().err
    assert "Value-Line" in err
    assert not any(t == "Target" for t in _texts(slide)), _texts(slide)


def test_value_and_benchmark_lines_all_none_series_no_crash():
    """Regression: series with exclusively None values must not raise ValueError."""
    from types import SimpleNamespace
    from charts import add_benchmark_line, add_value_line
    from pptx.util import Inches

    prs = _prs()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Minimal mock chart whose single series is all-None
    series = SimpleNamespace(values=[None, None, None])
    plot = SimpleNamespace(series=[series])
    va = SimpleNamespace(maximum_scale=None, minimum_scale=None)
    chart = SimpleNamespace(plots=[plot], value_axis=va)

    left, top, w, h = 1.0, 1.0, 8.0, 5.0
    # add_benchmark_line: all-None series — data_max stays at the spec value (10)
    add_benchmark_line(slide, chart, PAL, "10 Benchmark", left, top, w, h)
    # add_value_line: all-None series — axis derives from data_max=0; value 5
    # falls outside 0.._nice_ceil(0*1.05) so line is skipped (no crash)
    add_value_line(slide, chart, PAL, "Target, 5", left, top, w, h)


# ── stacked-100 dual labels ──────────────────────────────────────────────────
STACK_MD = """## Slide 1: Mix shifts decisively toward subscriptions
**Layout:** two-column-split
**Visual:** chart:stacked-100
**Series:** Subs, Services, Hardware
- Labels: pct
**Data:**
- 2023: 1, 1, 1
- 2024: 8.4, 7.6, 4
- Subscriptions now carry the mix
"""


def _chart(slide):
    return next(sh.chart for sh in slide.shapes
                if getattr(sh, "has_chart", False))


def _label_col(chart, j):
    return [s.points[j].data_label.text_frame.text
            for s in chart.plots[0].series]


def test_stacked_100_chart_type_registered():
    from pptx.enum.chart import XL_CHART_TYPE
    from charts import CHART_TYPES
    assert CHART_TYPES["stacked-100"] == XL_CHART_TYPE.COLUMN_STACKED_100


def test_stacked_100_pct_labels_sum_to_100():
    _, slides = parse_outline(STACK_MD)
    assert slides[0]["labels_mode"] == "pct"
    slide = builders.LAYOUT_MAP["two-column-split"](_prs(), slides[0], PAL, CTX)
    chart = _chart(slide)
    # 1/1/1 splits 33/33/34 via largest-remainder, not 33/33/33
    assert _label_col(chart, 0) == ["33%", "33%", "34%"]
    assert _label_col(chart, 1) == ["42%", "38%", "20%"]


def test_stacked_100_abs_labels():
    md = STACK_MD.replace("- Labels: pct", "- Labels: abs")
    _, slides = parse_outline(md)
    slide = builders.LAYOUT_MAP["two-column-split"](_prs(), slides[0], PAL, CTX)
    assert _label_col(_chart(slide), 1) == ["8.4", "7.6", "4"]


def test_stacked_100_both_labels():
    md = STACK_MD.replace("- Labels: pct", "- Labels: both")
    _, slides = parse_outline(md)
    slide = builders.LAYOUT_MAP["two-column-split"](_prs(), slides[0], PAL, CTX)
    assert _label_col(_chart(slide), 1) == \
        ["42% (8.4)", "38% (7.6)", "20% (4)"]


def test_stacked_100_bad_labels_mode_warns(capsys):
    md = STACK_MD.replace("- Labels: pct", "- Labels: nope")
    _, slides = parse_outline(md)
    builders.LAYOUT_MAP["two-column-split"](_prs(), slides[0], PAL, CTX)
    assert "Labels" in capsys.readouterr().err


# ── same-scale groups ────────────────────────────────────────────────────────
SCALE_MD = """**Scale-Group:** auto

## Slide 1: Americas revenue compounds fastest
**Layout:** two-column-split
**Visual:** chart:column
**Data:**
- 2023: 30
- 2024: 42
- Americas leads growth

## Slide 2: EMEA lags at the identical scale
**Layout:** two-column-split
**Visual:** chart:column
**Data:**
- 2023: 12
- 2024: 17
- EMEA trails on volume

## Slide 3: Margin trend stands alone
**Layout:** two-column-split
**Visual:** chart:line
**Data:**
- 2023: 5
- 2024: 6
- Single line chart stays auto-scaled
"""


def test_scale_group_auto_shares_axis_max(capsys):
    from build_deck import apply_scale_groups
    meta, slides = parse_outline(SCALE_MD)
    assert meta["scale_group"] == "auto"
    slides = apply_scale_groups(meta, slides)
    assert slides[0]["axis_max"] == 50  # _nice_ceil(42)
    assert slides[1]["axis_max"] == 50
    assert "axis_max" not in slides[2]  # group of one -> untouched
    out = capsys.readouterr().out
    assert "Scale-group" in out and "slides 1, 2" in out, out


def test_scale_group_skips_explicit_axis_max():
    from build_deck import apply_scale_groups
    md = SCALE_MD.replace("- 2024: 17\n", "- 2024: 17\n- Axis-Max: 100\n")
    meta, slides = parse_outline(md)
    slides = apply_scale_groups(meta, slides)
    assert slides[1]["axis_max"] == "100"   # explicit wins, slide excluded
    assert "axis_max" not in slides[0]      # remaining group of one -> no-op


def test_scale_group_off_by_default():
    from build_deck import apply_scale_groups
    md = SCALE_MD.replace("**Scale-Group:** auto\n\n", "")
    meta, slides = parse_outline(md)
    slides = apply_scale_groups(meta, slides)
    assert all("axis_max" not in s for s in slides)
