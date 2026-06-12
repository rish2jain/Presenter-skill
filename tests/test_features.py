"""Tests for narrative, smart layout, diff, and template helpers."""
import re
import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
sys.path.insert(0, str(SCRIPTS))

from build_deck import parse_outline, validate, build  # noqa: E402
from narrative import normalize_purpose, validate_narrative, generate_appendix_outline  # noqa: E402
from smart_layout import auto_layout  # noqa: E402
from template_helpers import build_template_map, palette_from_theme  # noqa: E402
import builders  # noqa: E402


@pytest.fixture(autouse=True)
def _density():
    builders.set_density("compact")


CTX = {"outline_dir": ROOT / "assets", "assets_dir": ROOT / "assets"}


def test_auto_layout_stat_callout():
    _, slides = parse_outline("## Slide 1: Metrics\n- Value=\"1\" Label=\"X\" Sublabel=\"Y\"\n")
    assert slides[0]["layout"] == "stat-callout"


def test_auto_layout_closing():
    _, slides = parse_outline("## Slide 1: Questions?\n**Layout:** closing\n")
    assert slides[0]["layout"] == "closing"


def test_narrative_purpose_validation():
    md = """**Purpose:** pitch
**Takeaway:** "We save enterprises $14M"

## Slide 1: Title
**Layout:** title
- Notes: "Open"

## Slide 2: Close
**Layout:** closing
- Notes: "End"
"""
    meta, slides = parse_outline(md)
    warnings = validate_narrative(meta, slides)
    assert not any("Takeaway" in w for w in warnings)
    assert any("stat-callout" in w for w in warnings)


def test_variant_preset_b():
    outline = ROOT / "assets" / "example-outline.md"
    out = Path("/tmp/variant-b-deck.pptx")
    assert build(outline, out, variant="b")
    assert out.is_file()


def test_palette_from_theme():
    pal = palette_from_theme({"accent1": "C9A84C", "dk1": "0A0F1E", "lt1": "F1F5F9"})
    assert pal["accent1"] == "C9A84C"
    assert pal["dark"] is True


def test_generate_appendix():
    meta, slides = parse_outline("**Purpose:** pitch\n## Slide 1: T\n**Layout:** title\n")
    text = generate_appendix_outline(meta, slides)
    assert "Appendix" in text
    assert "Financial Detail" in text


def test_diff_deck_example(tmp_path):
    outline = ROOT / "assets" / "example-outline.md"
    out = tmp_path / "deck.pptx"
    assert build(outline, out)
    from diff_deck import diff_outline
    errors, warnings = diff_outline(outline, out)
    assert not errors
    assert not any("expected text not found" in w and "'title'" in w for w in warnings)
    assert not any("expected text not found" in w and "'closing'" in w for w in warnings)


def test_build_template_map():
    from pptx import Presentation
    prs = Presentation()
    m = build_template_map(prs)
    assert "layout_map" in m
    assert "title" in m["layout_map"] or m["layout_scores"].get("title", 0) >= 0


def test_thumbnail_grid_is_labeled(tmp_path):
    from PIL import Image
    from render_slides import create_thumbnail_grid
    imgs = []
    for i in range(3):
        p = tmp_path / f"slide-{i + 1:02d}.png"
        Image.new("RGB", (400, 225), (40, 40, 60)).save(p)
        imgs.append(p)
    out = tmp_path / "grid.png"
    create_thumbnail_grid(imgs, out, cols=2)
    grid = Image.open(out)
    # badge: each cell's top-left corner carries an opaque label strip
    assert grid.getpixel((4, 4)) == (0, 0, 0), "badge background not black at (4,4)"


def test_every_palette_has_chart_series_token():
    from palettes import PALETTES
    for key, pal in PALETTES.items():
        assert len(pal.get("chart_series", [])) >= 3, key


def test_load_custom_palette(tmp_path):
    import json
    from palettes import PALETTES, load_custom_palettes
    pdir = tmp_path / "palettes"
    pdir.mkdir()
    (pdir / "acme-brand.json").write_text(json.dumps({
        "bg": "101820", "bg_deep": "0A0F14", "surface": "1E2A33",
        "accent1": "FEE715", "accent2": "8DA9C4", "accent3": "5C946E",
        "text": "F4F4F4", "text_muted": "9DB2BF", "dark": True}))
    loaded = load_custom_palettes(pdir)
    try:
        assert "acme-brand" in loaded and "acme-brand" in PALETTES
        pal = PALETTES["acme-brand"]
        assert pal["font_title"]                      # font defaults filled
        assert len(pal["chart_series"]) >= 3          # token derived
    finally:
        PALETTES.pop("acme-brand", None)              # don't leak into other tests


def test_invalid_custom_palette_rejected(tmp_path):
    import json
    from palettes import PALETTES, load_custom_palettes
    pdir = tmp_path / "palettes"
    pdir.mkdir()
    (pdir / "broken.json").write_text(json.dumps({"bg": "101820"}))
    loaded = load_custom_palettes(pdir)
    assert "broken" not in loaded and "broken" not in PALETTES


def test_custom_palette_with_non_hex_colors_rejected(tmp_path):
    import json
    from palettes import PALETTES, load_custom_palettes
    pdir = tmp_path / "palettes"
    pdir.mkdir()
    (pdir / "named.json").write_text(json.dumps({
        "bg": "red", "bg_deep": "0A0F14", "surface": "1E2A33",
        "accent1": "FEE715", "accent2": "8DA9C4", "accent3": "5C946E",
        "text": "F4F4F4", "text_muted": "9DB2BF", "dark": True}))
    assert load_custom_palettes(pdir) == []
    assert "named" not in PALETTES


def test_custom_palette_rejects_short_chart_series(tmp_path):
    import json
    from palettes import PALETTES, load_custom_palettes
    pdir = tmp_path / "palettes"
    pdir.mkdir()
    base = {
        "bg": "101820", "bg_deep": "0A0F14", "surface": "1E2A33",
        "accent1": "FEE715", "accent2": "8DA9C4", "accent3": "5C946E",
        "text": "F4F4F4", "text_muted": "9DB2BF", "dark": True,
        "chart_series": ["FEE715", "8DA9C4"],
    }
    (pdir / "short-series.json").write_text(json.dumps(base))
    assert load_custom_palettes(pdir) == []
    assert "short-series" not in PALETTES


def test_sticker_and_kicker_parsed():
    md = ('## Slide 1: Revenue doubled on pricing discipline\n'
          '- Sticker: Illustrative\n'
          '- Kicker: "Lock in pricing gains before renewals"\n')
    _, slides = parse_outline(md)
    assert slides[0]["sticker"] == "Illustrative"
    assert slides[0]["kicker"] == "Lock in pricing gains before renewals"


def test_kicker_restating_title_warned():
    md = ('## Slide 1: Revenue doubled across all regions this year\n'
          '**Layout:** bullet-list\n'
          '- Kicker: "Revenue doubled across all regions"\n'
          '- Point 1: "Growth was broad-based"\n'
          '- Notes: "n"\n')
    meta, slides = parse_outline(md)
    _, warnings = validate(slides, CTX, meta)
    assert any("kicker restates the title" in w for w in warnings), warnings


def test_kicker_advancing_argument_not_warned():
    md = ('## Slide 1: Revenue doubled across all regions this year\n'
          '**Layout:** bullet-list\n'
          '- Kicker: "Approve the follow-on investment before quarter close"\n'
          '- Point 1: "Growth was broad-based"\n'
          '- Notes: "n"\n')
    meta, slides = parse_outline(md)
    _, warnings = validate(slides, CTX, meta)
    assert not any("kicker restates" in w for w in warnings), warnings


def test_kicker_on_no_footer_layout_warned():
    for layout in ("title", "closing", "section-divider", "full-image"):
        md = (f'## Slide 1: Opening\n'
              f'**Layout:** {layout}\n'
              '- Kicker: "Act now before the window closes"\n'
              '- Notes: "n"\n')
        meta, slides = parse_outline(md)
        _, warnings = validate(slides, CTX, meta)
        assert any("Kicker is not rendered" in w and layout in w
                   for w in warnings), (layout, warnings)


def test_heading_over_15_words_warned():
    heading = " ".join(["word"] * 16)
    md = (f'## Slide 1: {heading}\n'
          '**Layout:** bullet-list\n- Point 1: "x"\n- Notes: "n"\n')
    meta, slides = parse_outline(md)
    _, warnings = validate(slides, CTX, meta)
    assert any("exceeds 15 words" in w for w in warnings), warnings


def test_exhibit_heading_without_number_warned():
    md = ('## Slide 1: Margins expanded sharply on pricing discipline\n'
          '**Layout:** waterfall\n'
          '**Data:**\n- FY24: 42\n- Pricing: +9\n- FY25: total\n'
          '- Notes: "n"\n')
    meta, slides = parse_outline(md)
    _, warnings = validate(slides, CTX, meta)
    assert any("no number" in w for w in warnings), warnings


def test_exhibit_heading_with_number_not_warned():
    md = ('## Slide 1: Margins expanded 400bps on pricing discipline\n'
          '**Layout:** waterfall\n'
          '**Data:**\n- FY24: 42\n- Pricing: +9\n- FY25: total\n'
          '- Notes: "n"\n')
    meta, slides = parse_outline(md)
    _, warnings = validate(slides, CTX, meta)
    assert not any("no number" in w for w in warnings), warnings


def test_heading_with_and_warned():
    md = ('## Slide 1: We doubled revenue and we cut operating costs\n'
          '**Layout:** bullet-list\n- Point 1: "x"\n- Notes: "n"\n')
    meta, slides = parse_outline(md)
    _, warnings = validate(slides, CTX, meta)
    assert any("joins two messages" in w for w in warnings), warnings


def test_sticker_and_kicker_rendered(tmp_path):
    from pptx import Presentation
    md = ('## Slide 1: Margins expanded 400bps on pricing discipline\n'
          '**Layout:** bullet-list\n'
          '- Sticker: Illustrative\n'
          '- Kicker: "Lock in pricing gains before renewals"\n'
          '- Pricing actions held through renewals\n'
          '- Notes: "n"\n')
    outline = tmp_path / "o.md"
    outline.write_text(md)
    out = tmp_path / "deck.pptx"
    assert build(outline, out)
    prs = Presentation(str(out))
    texts = {sh.text_frame.text for sh in prs.slides[0].shapes
             if sh.has_text_frame}
    assert "ILLUSTRATIVE" in texts, texts
    assert "Lock in pricing gains before renewals" in texts, texts
    # sticker tag sits top-right (no section label -> y ~0.14)
    tag = [sh for sh in prs.slides[0].shapes
           if sh.has_text_frame and sh.text_frame.text == "ILLUSTRATIVE"][0]
    assert abs(tag.top / 914400 - 0.16) < 0.05, tag.top
    # kicker box top ~6.18
    kick = [sh for sh in prs.slides[0].shapes
            if sh.has_text_frame
            and sh.text_frame.text == "Lock in pricing gains before renewals"][0]
    assert 6.1 < kick.top / 914400 < 6.5, kick.top


def test_sticker_drops_below_section_label(tmp_path):
    from pptx import Presentation
    md = ('## Slide 1: Growth\n**Layout:** section-divider\n'
          '## Slide 2: Margins expanded 400bps on pricing discipline\n'
          '**Layout:** bullet-list\n'
          '- Sticker: Preliminary\n'
          '- Pricing actions held through renewals\n'
          '- Notes: "n"\n')
    outline = tmp_path / "o.md"
    outline.write_text(md)
    out = tmp_path / "deck.pptx"
    assert build(outline, out)
    prs = Presentation(str(out))
    tag = [sh for sh in prs.slides[1].shapes
           if sh.has_text_frame and sh.text_frame.text == "PRELIMINARY"][0]
    assert abs(tag.top / 914400 - 0.52) < 0.05, tag.top


def test_dump_numbers_extracts_per_slide(capsys, tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    from qa_check import dump_numbers
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    tb.text_frame.text = "Revenue grew 12% to $4.2B in FY25"
    p = tmp_path / "n.pptx"
    prs.save(str(p))
    dump_numbers(p)
    out = capsys.readouterr().out
    assert "Slide 1" in out and "12%" in out and "$4.2B" in out and "FY25" in out


def test_num_token_rx_suffix_handling():
    from qa_check import NUM_TOKEN_RX
    assert NUM_TOKEN_RX.findall("spread tightened 200bps") == ["200bps"]
    assert NUM_TOKEN_RX.findall("3 key levers") == ["3"]


# ── T9: image alpha/duotone options ──────────────────────────────────────────
def _make_png(path, color=(120, 80, 40)):
    from PIL import Image
    Image.new("RGB", (320, 200), color).save(path)
    return path


def _add_test_picture(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    img = _make_png(tmp_path / "photo.png")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pic = slide.shapes.add_picture(str(img), Inches(1), Inches(1),
                                   Inches(3), Inches(2))
    return prs, pic


def test_parse_visual_opts_no_opts():
    from helpers import parse_visual_opts
    assert parse_visual_opts("user-image:hero.png") == \
        ("image", "user-image:hero.png", {})
    assert parse_visual_opts("chart:bar") == ("chart", "bar", {})
    assert parse_visual_opts("") == (None, None, {})
    assert parse_visual_opts("none") == (None, None, {})


def test_parse_visual_opts_alpha():
    from helpers import parse_visual_opts
    assert parse_visual_opts("photo.png|alpha=85") == \
        ("image", "photo.png", {"alpha": 85})


def test_parse_visual_opts_duotone():
    from helpers import parse_visual_opts
    assert parse_visual_opts("photo.png|duotone") == \
        ("image", "photo.png", {"duotone": True})


def test_parse_visual_opts_combined():
    from helpers import parse_visual_opts
    assert parse_visual_opts("user-image:p.png|alpha=85|duotone") == \
        ("image", "user-image:p.png", {"alpha": 85, "duotone": True})


def test_parse_visual_opts_malformed_alpha_ignored_with_warning(capsys):
    from helpers import parse_visual_opts
    kind, value, opts = parse_visual_opts("photo.png|alpha=high")
    assert (kind, value, opts) == ("image", "photo.png", {})
    assert "alpha" in capsys.readouterr().err
    kind, value, opts = parse_visual_opts("photo.png|alpha=400")
    assert opts == {}
    assert "alpha" in capsys.readouterr().err


def test_parse_visual_opts_unknown_option_warned(capsys):
    from helpers import parse_visual_opts
    assert parse_visual_opts("photo.png|sepia") == ("image", "photo.png", {})
    assert "sepia" in capsys.readouterr().err


def test_parse_visual_two_tuple_strips_opts(capsys):
    from helpers import parse_visual
    assert parse_visual("photo.png|alpha=85|duotone") == ("image", "photo.png")
    # wrapper is silent even on malformed options (validation hot path)
    assert parse_visual("photo.png|alpha=bad") == ("image", "photo.png")
    assert capsys.readouterr().err == ""


def test_validate_warns_on_malformed_image_option(tmp_path):
    _make_png(tmp_path / "photo.png")
    md = ('## Slide 1: Platform spend concentrates in two clouds today\n'
          '**Layout:** two-column-split\n'
          '**Visual:** photo.png|alpha=bad\n'
          '- Spend concentrates fast\n'
          '- Notes: "n"\n')
    meta, slides = parse_outline(md)
    ctx = {"outline_dir": tmp_path, "assets_dir": tmp_path}
    errors, warnings = validate(slides, ctx, meta)
    assert not errors, errors
    assert any("image option" in w for w in warnings), warnings


def test_set_picture_alpha_xml(tmp_path):
    from helpers import set_picture_alpha
    from pptx.oxml.ns import qn
    _, pic = _add_test_picture(tmp_path)
    set_picture_alpha(pic, 85)
    blip = pic._element.find(qn("p:blipFill") + "/" + qn("a:blip"))
    fixes = blip.findall(qn("a:alphaModFix"))
    assert len(fixes) == 1 and fixes[0].get("amt") == "85000"
    set_picture_alpha(pic, 40)  # re-apply replaces, not stacks
    fixes = blip.findall(qn("a:alphaModFix"))
    assert len(fixes) == 1 and fixes[0].get("amt") == "40000"


def test_set_picture_duotone_xml(tmp_path):
    from helpers import set_picture_duotone
    from pptx.oxml.ns import qn
    _, pic = _add_test_picture(tmp_path)
    set_picture_duotone(pic, "0A0F1E", "C9A84C")
    blip = pic._element.find(qn("p:blipFill") + "/" + qn("a:blip"))
    duos = blip.findall(qn("a:duotone"))
    assert len(duos) == 1
    clrs = [c.get("val") for c in duos[0].findall(qn("a:srgbClr"))]
    assert clrs == ["0A0F1E", "C9A84C"]


def test_image_opts_wired_into_build(tmp_path):
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from palettes import get_palette
    _make_png(tmp_path / "photo.png")
    md = ('## Slide 1: Platform spend concentrates in two clouds today\n'
          '**Layout:** two-column-split\n'
          '**Visual:** photo.png|alpha=85|duotone\n'
          '- Spend concentrates fast\n'
          '- Notes: "n"\n'
          '## Slide 2: Full bleed\n'
          '**Layout:** full-image\n'
          '**Visual:** photo.png|alpha=70\n'
          '- Caption: "Datacenter floor"\n'
          '- Notes: "n"\n')
    outline = tmp_path / "o.md"
    outline.write_text(md)
    out = tmp_path / "deck.pptx"
    assert build(outline, out)
    prs = Presentation(str(out))
    pal = get_palette("")
    pics = [sh for sh in prs.slides[0].shapes
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pics, "two-column-split image not placed"
    blip = pics[0]._element.find(qn("p:blipFill") + "/" + qn("a:blip"))
    assert blip.find(qn("a:alphaModFix")).get("amt") == "85000"
    duo = blip.find(qn("a:duotone"))
    clrs = [c.get("val") for c in duo.findall(qn("a:srgbClr"))]
    assert clrs == [pal["bg"], pal["accent1"]]
    pics2 = [sh for sh in prs.slides[1].shapes
             if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert pics2, "full-image picture not placed"
    blip2 = pics2[0]._element.find(qn("p:blipFill") + "/" + qn("a:blip"))
    assert blip2.find(qn("a:alphaModFix")).get("amt") == "70000"
    assert blip2.find(qn("a:duotone")) is None


# ── T9: PowerPoint sections (p14:sectionLst) ─────────────────────────────────
P14_NS = "{http://schemas.microsoft.com/office/powerpoint/2010/main}"

SECTIONED_MD = ('## Slide 1: T\n**Layout:** title\n- Title: "Deck"\n'
                '## Slide 2: Growth\n**Layout:** section-divider\n'
                '## Slide 3: Revenue doubled on pricing discipline this year\n'
                '**Layout:** bullet-list\n- Point 1: "x"\n- Notes: "n"\n'
                '## Slide 4: Costs\n**Layout:** section-divider\n'
                '## Slide 5: Unit costs fell 20% on platform consolidation\n'
                '**Layout:** bullet-list\n- Point 1: "y"\n- Notes: "n"\n'
                '## Appendix\n'
                '## Slide 6: Detailed cost model assumptions\n'
                '**Layout:** bullet-list\n- Point 1: "z"\n')


def _sections_of(pptx_path):
    from pptx import Presentation
    prs = Presentation(str(pptx_path))
    lst = prs.element.find(".//" + P14_NS + "sectionLst")
    if lst is None:
        return prs, None
    return prs, [(s.get("name"), s.get("id"),
                  [int(e.get("id")) for e in s.findall(
                      P14_NS + "sldIdLst/" + P14_NS + "sldId")])
                 for s in lst.findall(P14_NS + "section")]


def test_sections_injected_with_two_dividers(tmp_path):
    outline = tmp_path / "o.md"
    outline.write_text(SECTIONED_MD)
    out = tmp_path / "deck.pptx"
    assert build(outline, out)
    prs, sections = _sections_of(out)
    assert sections is not None
    assert [s[0] for s in sections] == ["Opening", "Growth", "Costs", "Backup"]
    assert [len(s[2]) for s in sections] == [1, 2, 2, 1]
    deck_ids = [int(s.get("id")) for s in prs.slides._sldIdLst]
    assert [i for s in sections for i in s[2]] == deck_ids
    for _, guid, _ in sections:
        assert re.fullmatch(r"\{[0-9A-F-]{36}\}", guid), guid


def test_sections_guids_deterministic(tmp_path):
    outline = tmp_path / "o.md"
    outline.write_text(SECTIONED_MD)
    out1, out2 = tmp_path / "d1.pptx", tmp_path / "d2.pptx"
    assert build(outline, out1) and build(outline, out2)
    _, s1 = _sections_of(out1)
    _, s2 = _sections_of(out2)
    assert [s[:2] for s in s1] == [s[:2] for s in s2]


def test_sections_meta_off_disables(tmp_path):
    outline = tmp_path / "o.md"
    outline.write_text("**Sections:** off\n\n" + SECTIONED_MD)
    out = tmp_path / "deck.pptx"
    assert build(outline, out)
    _, sections = _sections_of(out)
    assert sections is None


def test_sections_single_divider_not_injected(tmp_path):
    md = ('## Slide 1: T\n**Layout:** title\n- Title: "Deck"\n'
          '## Slide 2: Growth\n**Layout:** section-divider\n'
          '## Slide 3: Revenue doubled on pricing discipline this year\n'
          '**Layout:** bullet-list\n- Point 1: "x"\n- Notes: "n"\n')
    outline = tmp_path / "o.md"
    outline.write_text(md)
    out = tmp_path / "deck.pptx"
    assert build(outline, out)
    _, sections = _sections_of(out)
    assert sections is None


def test_sections_deck_reopens_and_passes_qa(tmp_path):
    from qa_check import check_deck
    outline = tmp_path / "o.md"
    outline.write_text(SECTIONED_MD)
    out = tmp_path / "deck.pptx"
    assert build(outline, out)
    issues = check_deck(out)
    assert not issues["error"], issues["error"]


# ── T9: native table banding attributes ──────────────────────────────────────
def test_table_marked_first_row_and_banding(tmp_path):
    from pptx import Presentation
    md = ('## Slide 1: Pricing tiers converge on the enterprise plan\n'
          '**Layout:** table\n'
          '| Tier | Price |\n| --- | --- |\n| Pro | $99 |\n| Ent | $499 |\n'
          '- Notes: "n"\n')
    outline = tmp_path / "o.md"
    outline.write_text(md)
    out = tmp_path / "deck.pptx"
    assert build(outline, out)
    prs = Presentation(str(out))
    tables = [sh.table for sh in prs.slides[0].shapes
              if getattr(sh, "has_table", False)]
    assert tables, "table not built"
    tblPr = tables[0]._tbl.tblPr
    assert tblPr.get("firstRow") == "1"
    assert tblPr.get("bandRow") == "1"
    # the accessibility checker's header-row rule passes
    from qa_check import check_table_a11y
    issues = {"error": [], "warn": []}
    shape = [sh for sh in prs.slides[0].shapes
             if getattr(sh, "has_table", False)][0]
    check_table_a11y(shape, issues, "Slide 1")
    assert not issues["error"], issues["error"]
