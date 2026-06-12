"""Tests for outline parsing, validation, and golden-path build."""
import sys
from pathlib import Path

import pytest

ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
sys.path.insert(0, str(SCRIPTS))

from build_deck import parse_outline, validate, _parse_data_point, _series_count  # noqa: E402
import builders  # noqa: E402


@pytest.fixture(autouse=True)
def _default_density():
    builders.set_density("compact")
    yield
    import palettes
    palettes.set_cjk(False)


CTX = {"outline_dir": ROOT / "assets", "assets_dir": ROOT / "assets"}


def test_data_block_tolerates_blank_lines():
    md = """## Slide 1: Chart
**Layout:** two-column-split
**Visual:** chart:bar
**Data:**
- 2024: 42

- 2025: 67
- Heading: Growth
- bullet one
"""
    _, slides = parse_outline(md)
    assert slides[0]["data"] == [("2024", 42.0), ("2025", 67.0)]
    assert slides[0]["heading"] == "Growth"
    assert slides[0]["bullets"] == ["bullet one"]


def test_multi_series_data_parsing():
    md = """## Slide 1: Chart
**Layout:** two-column-split
**Visual:** chart:line
**Series:** Revenue, Costs
**Data:**
- Q1: 4.2, 3.1
- Q2: 5.1, 3.0
"""
    _, slides = parse_outline(md)
    assert _series_count(slides[0]) == 2
    assert slides[0]["data"] == [("Q1", [4.2, 3.1]), ("Q2", [5.1, 3.0])]
    errors, _ = validate(slides, CTX)
    assert not errors


def test_multi_series_mismatch_errors():
    md = """## Slide 1: Chart
**Layout:** two-column-split
**Visual:** chart:line
**Series:** Revenue, Costs
**Data:**
- Q1: 4.2
"""
    _, slides = parse_outline(md)
    errors, _ = validate(slides, CTX)
    assert any("multi-series" in e for e in errors)


def test_parse_data_point_dollar_amounts():
    assert _parse_data_point("2024: $42B") == ("2024", 42.0)
    assert _parse_data_point("Q1: -5.2") == ("Q1", -5.2)


def test_validate_warns_missing_notes():
    md = """## Slide 1: X
**Layout:** bullet-list
- one
"""
    _, slides = parse_outline(md)
    _, warnings = validate(slides, CTX)
    assert any("Notes" in w for w in warnings)


def test_validate_title_heading_fallback_warning():
    md = """## Slide 1: My Title
**Layout:** title
"""
    _, slides = parse_outline(md)
    _, warnings = validate(slides, CTX)
    assert any("My Title" in w for w in warnings)


def test_validate_rejects_chart_on_comparison_side():
    md = """## Slide 1: Compare
**Layout:** comparison
**Visual-Left:** chart:bar
- Left label: A
- Right label: B
"""
    _, slides = parse_outline(md)
    errors, _ = validate(slides, CTX)
    assert any("visual_left" in e for e in errors)


def test_example_outline_check():
    outline = ROOT / "assets" / "example-outline.md"
    meta, slides = parse_outline(outline.read_text(encoding="utf-8"))
    errors, _ = validate(slides, CTX)
    assert not errors
    assert len(slides) == 6


def test_golden_build_and_qa(tmp_path):
    outline = ROOT / "assets" / "example-outline.md"
    out = tmp_path / "deck.pptx"
    from build_deck import build

    assert build(outline, out, check_only=False)
    assert out.is_file()

    from qa_check import check_deck

    issues = check_deck(out)
    assert not issues["error"]


TRACKED_MD = """**Auto-Agenda:** track

## Slide 1: Acme FY27 Strategy
**Layout:** title
- Title: Acme FY27 Strategy
- Subtitle: Board readout

## Slide 2: Diagnosis
**Layout:** section-divider
- Subtitle: Where we are

## Slide 3: Costs have outgrown revenue for three years
- Cost CAGR 12% vs revenue 4%
- Notes: Set up the problem.

## Slide 4: Plan
**Layout:** section-divider
- Subtitle: What we will do

## Slide 5: Three levers close the gap by FY27
- Tiering, exit, discount discipline
- Notes: The plan.
"""


def test_auto_agenda_track_inserts_tracker_slides():
    from build_deck import parse_outline, apply_auto_agenda
    meta, slides = parse_outline(TRACKED_MD)
    out = apply_auto_agenda(meta, slides)
    layouts = [s["layout"] for s in out]
    # title, agenda, divider, tracker, content, divider, tracker, content
    assert layouts == ["title", "agenda", "section-divider", "agenda",
                       "bullet-list", "section-divider", "agenda",
                       "bullet-list"], layouts
    trackers = [s for s in out if s["layout"] == "agenda" and s.get("current")]
    assert [t["current"] for t in trackers] == ["Diagnosis", "Plan"]
    assert out[1]["bullets"] == ["Diagnosis", "Plan"]


def test_auto_agenda_on_inserts_only_overview():
    from build_deck import parse_outline, apply_auto_agenda
    meta, slides = parse_outline(TRACKED_MD.replace("track", "on"))
    out = apply_auto_agenda(meta, slides)
    agendas = [s for s in out if s["layout"] == "agenda"]
    assert len(agendas) == 1 and not agendas[0].get("current")


def test_auto_agenda_off_is_identity():
    from build_deck import parse_outline, apply_auto_agenda
    meta, slides = parse_outline(TRACKED_MD.replace("**Auto-Agenda:** track\n\n", ""))
    assert apply_auto_agenda(meta, slides) == slides


def test_ghost_mode_renders_layout_labels(tmp_path):
    from build_deck import build
    from pptx import Presentation
    md = tmp_path / "o.md"
    md.write_text("""## Slide 1: Three levers bridge run-rate down to target
**Layout:** waterfall
**Data:**
- FY25: 46
- Tiering: -8
- FY27: total
- Notes: n.
""")
    out = tmp_path / "g.pptx"
    assert build(str(md), str(out), ghost=True)
    prs = Presentation(str(out))
    texts = " ".join(sh.text_frame.text for sh in prs.slides[0].shapes
                     if getattr(sh, "has_text_frame", False))
    assert "waterfall" in texts            # layout label shown
    assert "Three levers" in texts          # real action title kept
    assert "46" not in texts                # data NOT rendered


def test_stamp_appears_on_every_slide(tmp_path):
    from build_deck import build
    from pptx import Presentation
    md = tmp_path / "o.md"
    md.write_text("""**Stamp:** DRAFT

## Slide 1: Costs have outgrown revenue for three years
- Cost CAGR 12% vs revenue 4%
- Notes: n.

## Slide 2: Three levers close the gap by FY27
- Tiering, exit, discounts
- Notes: n.
""")
    out = tmp_path / "d.pptx"
    assert build(str(md), str(out))
    prs = Presentation(str(out))
    for slide in prs.slides:
        texts = [sh.text_frame.text for sh in slide.shapes
                 if getattr(sh, "has_text_frame", False)]
        assert "DRAFT" in texts, texts


def test_auto_agenda_excludes_appendix_sections():
    from build_deck import parse_outline, apply_auto_agenda
    md = TRACKED_MD + """
## Appendix

## Slide 6: Backup detail
**Layout:** section-divider
- Subtitle: Backup
"""
    meta, slides = parse_outline(md)
    out = apply_auto_agenda(meta, slides)
    overview = next(s for s in out if s["layout"] == "agenda")
    assert overview["bullets"] == ["Diagnosis", "Plan"], overview["bullets"]


def test_auto_agenda_divider_first_gets_no_overview():
    from build_deck import parse_outline, apply_auto_agenda
    md = TRACKED_MD.split("## Slide 2", 1)[1]
    md = "**Auto-Agenda:** track\n\n## Slide 2" + md
    meta, slides = parse_outline(md)
    out = apply_auto_agenda(meta, slides)
    overviews = [s for s in out if s["layout"] == "agenda" and not s.get("current")]
    assert not overviews, [s.get("heading") for s in out]


def test_heading_attributes_set_layout_and_palette():
    from build_deck import parse_outline
    _, slides = parse_outline(
        "## Slide 1: Margin bridge tells the story {layout=waterfall palette=aurora}\n"
        "**Data:**\n- FY25: 46\n- FY27: total\n")
    assert slides[0]["layout"] == "waterfall"
    assert slides[0]["palette"] == "aurora"
    assert slides[0]["heading"] == "Margin bridge tells the story"


def test_explicit_layout_line_overrides_heading_attr():
    from build_deck import parse_outline
    _, slides = parse_outline(
        "## Slide 1: T {layout=waterfall}\n**Layout:** funnel\n"
        "**Data:**\n- A: 10\n- B: 5\n")
    assert slides[0]["layout"] == "funnel"


def test_validation_errors_carry_line_numbers():
    from build_deck import parse_outline, validate
    md = ("## Slide 1: Costs have outgrown revenue for years\n"
          "- a bullet\n"
          "- Notes: n.\n"
          "\n"
          "## Slide 2: Bad chart slide misses its data block\n"
          "**Layout:** waterfall\n"
          "- Notes: n.\n")
    _, slides = parse_outline(md)
    assert slides[0]["_line"] == 1 and slides[1]["_line"] == 5
    errors, _ = validate(slides, {"outline_dir": Path("."),
                                  "assets_dir": Path("assets")})
    assert any("Slide 2 (line 5)" in e for e in errors), errors


def test_unknown_cli_palette_warns(tmp_path, capsys):
    from build_deck import build
    md = tmp_path / "o.md"
    md.write_text("## Slide 1: Costs have outgrown revenue for years\n"
                  "- a bullet\n- Notes: n.\n")
    assert build(str(md), str(tmp_path / "d.pptx"), palette_key="auroa")
    assert "unknown --palette" in capsys.readouterr().err


def test_unknown_heading_attr_warns(capsys):
    from build_deck import parse_outline
    parse_outline("## Slide 1: T {laoyut=waterfall}\n- b\n")
    assert "unknown heading attribute" in capsys.readouterr().err


def test_quoted_heading_attr_warns(capsys):
    from build_deck import parse_outline
    parse_outline('## Slide 1: T {layout="waterfall"}\n- b\n')
    assert "did not parse" in capsys.readouterr().err


# ── Tracker tabs (T6) ────────────────────────────────────────────────────────
TRACKER_MD = """**Tracker:** tabs

## Slide 1: Acme FY27 Strategy
**Layout:** title
- Title: Acme FY27 Strategy

## Slide 2: Diagnosis
**Layout:** section-divider

## Slide 3: Costs have outgrown revenue for three years
- Cost CAGR 12% vs revenue 4%
- Notes: n.

## Slide 4: Plan
**Layout:** section-divider

## Slide 5: Three levers close the gap by FY27
- Tiering, exit, discounts
- Notes: n.
"""


def _round_chips(slide):
    from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
    return [sh for sh in slide.shapes
            if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and sh.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE]


def _slide_texts(slide):
    return [sh.text_frame.text for sh in slide.shapes
            if getattr(sh, "has_text_frame", False)]


def test_tracker_tabs_render_chips(tmp_path):
    from build_deck import build
    from palettes import get_palette
    from pptx import Presentation
    from pptx.enum.dml import MSO_FILL
    from pptx.dml.color import RGBColor
    md = tmp_path / "o.md"
    md.write_text(TRACKER_MD)
    out = tmp_path / "d.pptx"
    assert build(str(md), str(out))
    prs = Presentation(str(out))
    chips = _round_chips(prs.slides[2])  # first content slide (Diagnosis)
    assert len(chips) == 2, "one chip per section"
    filled = [c for c in chips if c.fill.type == MSO_FILL.SOLID]
    assert len(filled) == 1, "exactly one (current) chip is filled"
    accent1 = get_palette("")["accent1"]
    assert filled[0].fill.fore_color.rgb == RGBColor.from_string(accent1)
    texts = _slide_texts(prs.slides[2])
    assert "Diagnosis" in texts and "Plan" in texts


def test_tracker_tabs_current_chip_moves(tmp_path):
    from build_deck import build
    from pptx import Presentation
    from pptx.enum.dml import MSO_FILL
    md = tmp_path / "o.md"
    md.write_text(TRACKER_MD)
    out = tmp_path / "d.pptx"
    assert build(str(md), str(out))
    prs = Presentation(str(out))
    chips = _round_chips(prs.slides[4])  # content slide in section "Plan"
    filled = [c for c in chips if c.fill.type == MSO_FILL.SOLID]
    assert len(chips) == 2 and len(filled) == 1
    # the filled chip is the right-most one (Plan is the last section)
    assert filled[0].left == max(c.left for c in chips)


def test_tracker_tabs_truncate_long_names(tmp_path):
    from build_deck import build
    from pptx import Presentation
    md = tmp_path / "o.md"
    md.write_text(TRACKER_MD.replace(
        "## Slide 4: Plan", "## Slide 4: Commercial Excellence Program"))
    out = tmp_path / "d.pptx"
    assert build(str(md), str(out))
    prs = Presentation(str(out))
    texts = _slide_texts(prs.slides[2])
    assert "Commercial Exc…" in texts, texts


def test_tracker_tabs_many_sections_fall_back(tmp_path, capsys):
    from build_deck import build
    from pptx import Presentation
    parts = ["**Tracker:** tabs\n"]
    n = 1
    for s in range(1, 7):
        parts.append(f"## Slide {n}: Section {s}\n**Layout:** section-divider\n")
        n += 1
        parts.append(f"## Slide {n}: Content slide for section {s} here\n"
                     "- a bullet\n- Notes: n.\n")
        n += 1
    md = tmp_path / "o.md"
    md.write_text("\n".join(parts))
    out = tmp_path / "d.pptx"
    assert build(str(md), str(out))
    err = capsys.readouterr().err
    assert err.count("exceed tab budget") == 1, "warn exactly once per build"
    prs = Presentation(str(out))
    content = prs.slides[1]  # first content slide
    assert not _round_chips(content)
    assert any("1/6 · Section 1" in t for t in _slide_texts(content))


def test_tracker_tabs_without_dividers_warns():
    from build_deck import parse_outline, validate
    meta, slides = parse_outline(
        "**Tracker:** tabs\n\n"
        "## Slide 1: Costs have outgrown revenue for years\n"
        "- a bullet\n- Notes: n.\n")
    _, warnings = validate(slides, CTX, meta)
    assert any("Tracker" in w for w in warnings), warnings


def test_tracker_tabs_appendix_keeps_backup_label(tmp_path):
    from build_deck import build
    from pptx import Presentation
    md = tmp_path / "o.md"
    md.write_text(TRACKER_MD + """
## Appendix

## Slide 6: Backup cost detail by region here
- detail
""")
    out = tmp_path / "d.pptx"
    assert build(str(md), str(out))
    prs = Presentation(str(out))
    backup = prs.slides[5]
    assert not _round_chips(backup)
    assert "BACKUP" in _slide_texts(backup)


# ── Auto-references slide (T10) ──────────────────────────────────────────────
SOURCED_MD = """**References:** on

## Slide 1: Costs have outgrown revenue for three years
- a bullet
- Source: Gartner 2026
- Notes: n.

## Slide 2: Plain slide carries no exhibit source
- b
- Notes: n.

## Slide 3: Unit costs run 3x the regional baseline
- c
- Source: Company filings
- Notes: n.
"""


def test_references_appends_sources_slide():
    from build_deck import apply_references
    meta, slides = parse_outline(SOURCED_MD)
    out = apply_references(meta, slides)
    assert len(out) == 4
    ref = out[-1]
    assert ref["layout"] == "bullet-list"
    assert ref["heading"] == "Sources"
    assert ref["_appendix"] is True
    assert ref["notes"] == "Auto-generated source register."
    assert ref["bullets"] == ["Slide 1 — Gartner 2026",
                              "Slide 3 — Company filings"]
    errors, _ = validate(out, CTX, meta)
    assert not errors


def test_references_uses_exhibit_numbers_when_exhibits_on():
    from build_deck import apply_references
    meta, slides = parse_outline("**Exhibits:** on\n" + SOURCED_MD)
    out = apply_references(meta, slides)
    assert out[-1]["bullets"] == ["Exhibit 1 — Gartner 2026",
                                  "Exhibit 2 — Company filings"]


def test_references_dedupes_repeated_sources():
    from build_deck import apply_references
    meta, slides = parse_outline(
        SOURCED_MD.replace("- Source: Company filings",
                           "- Source: Gartner 2026"))
    out = apply_references(meta, slides)
    assert out[-1]["bullets"] == ["Slides 1, 3 — Gartner 2026"]


def test_references_off_by_default():
    from build_deck import apply_references
    meta, slides = parse_outline(SOURCED_MD.replace("**References:** on\n", ""))
    assert apply_references(meta, slides) == slides


def test_references_on_without_sources_warns(capsys):
    from build_deck import apply_references
    meta, slides = parse_outline(
        "**References:** on\n\n"
        "## Slide 1: Costs have outgrown revenue for years\n"
        "- a bullet\n- Notes: n.\n")
    out = apply_references(meta, slides)
    assert out == slides
    assert "no '- Source:'" in capsys.readouterr().err


def test_references_numbering_counts_auto_agenda_slides():
    from build_deck import apply_auto_agenda, apply_references
    md = "**References:** on\n**Auto-Agenda:** on\n" + TRACKED_MD.replace(
        "**Auto-Agenda:** track\n", "")
    md += "- Source: Team analysis\n"
    meta, slides = parse_outline(md)
    out = apply_references(meta, apply_auto_agenda(meta, slides))
    # title, agenda, divider, content, divider, content(sourced) -> slide 6
    assert out[-1]["bullets"] == ["Slide 6 — Team analysis"]


# ── CSV chart data (T10) ─────────────────────────────────────────────────────
def _csv_outline(tmp_path, csv_text, csv_name="d.csv", extra=""):
    from build_deck import load_data_files
    (tmp_path / csv_name).write_text(csv_text, encoding="utf-8")
    md = (f"## Slide 1: Market grows 32% YoY through 2027\n"
          f"**Layout:** two-column-split\n"
          f"**Visual:** chart:bar\n"
          f"{extra}"
          f"- Data-File: {csv_name}\n"
          f"- Notes: n.\n")
    _, slides = parse_outline(md)
    ctx = {"outline_dir": tmp_path, "assets_dir": tmp_path / "assets"}
    errors, warnings = load_data_files(slides, ctx)
    return slides, errors, warnings


def test_data_file_two_column_csv(tmp_path):
    slides, errors, warnings = _csv_outline(tmp_path, "2024,$42B\n2025,67\n")
    assert not errors and not warnings
    assert slides[0]["data"] == [("2024", 42.0), ("2025", 67.0)]
    errs, _ = validate(slides, {"outline_dir": tmp_path,
                                "assets_dir": tmp_path / "assets"})
    assert not errs


def test_data_file_two_column_header_skipped(tmp_path):
    slides, errors, _ = _csv_outline(tmp_path, "label,value\n2024,42\n")
    assert not errors
    assert slides[0]["data"] == [("2024", 42.0)]


def test_data_file_multi_column_sets_series(tmp_path):
    slides, errors, _ = _csv_outline(
        tmp_path, "Quarter,Revenue,Costs\nQ1,4.2,3.1\nQ2,5.1,3.0\n")
    assert not errors
    assert slides[0]["series"] == "Revenue, Costs"
    assert slides[0]["data"] == [("Q1", [4.2, 3.1]), ("Q2", [5.1, 3.0])]


def test_data_file_keeps_explicit_series(tmp_path):
    slides, errors, _ = _csv_outline(
        tmp_path, "Quarter,Revenue,Costs\nQ1,4.2,3.1\n",
        extra="**Series:** Rev, Cost\n")
    assert not errors
    assert slides[0]["series"] == "Rev, Cost"


def test_data_file_missing_is_error(tmp_path):
    from build_deck import load_data_files
    _, slides = parse_outline("## Slide 1: T\n- Data-File: nope.csv\n")
    errors, _ = load_data_files(slides, {"outline_dir": tmp_path,
                                         "assets_dir": tmp_path / "assets"})
    assert any("not found" in e for e in errors)


def test_data_file_malformed_row_reports_row_number(tmp_path):
    slides, errors, _ = _csv_outline(tmp_path, "2024,42\n2025,n/a\n")
    assert any("row 2" in e for e in errors), errors


def test_data_file_wins_over_explicit_data(tmp_path):
    slides, errors, warnings = _csv_outline(
        tmp_path, "2024,42\n",
        extra="**Data:**\n- 1999: 7\n")
    assert not errors
    assert any("Data-File wins" in w for w in warnings)
    assert slides[0]["data"] == [("2024", 42.0)]


def test_data_file_resolves_from_assets_dir(tmp_path):
    from build_deck import load_data_files
    assets = tmp_path / "assets"
    assets.mkdir()
    (assets / "q.csv").write_text("A,1\nB,2\n", encoding="utf-8")
    _, slides = parse_outline("## Slide 1: T\n- Data-File: q.csv\n")
    errors, _ = load_data_files(slides, {"outline_dir": tmp_path / "elsewhere",
                                         "assets_dir": assets})
    assert not errors
    assert slides[0]["data"] == [("A", 1.0), ("B", 2.0)]


def test_data_file_supports_waterfall_total(tmp_path):
    slides, errors, _ = _csv_outline(tmp_path, "FY25,46\nTiering,-8\nFY27,total\n")
    assert not errors
    assert slides[0]["data"][-1] == ("FY27", "total")


# ── CJK font stacks (T10) ────────────────────────────────────────────────────
def test_deck_has_cjk_detection():
    from build_deck import deck_has_cjk
    _, slides = parse_outline("## Slide 1: 中文标题\n- 项目一\n")
    assert deck_has_cjk(slides)
    _, slides = parse_outline("## Slide 1: English only heading here\n- bullet\n")
    assert not deck_has_cjk(slides)
    _, slides = parse_outline("## Slide 1: Mixed\n- カタカナ bullet\n")
    assert deck_has_cjk(slides)


def test_set_cjk_overlays_palette_fonts(monkeypatch, capsys):
    import palettes
    monkeypatch.setitem(palettes._CJK, "font", "Test CJK")
    monkeypatch.setitem(palettes._CJK, "probed", True)
    palettes.set_cjk(True)
    pal = palettes.get_palette("")
    assert pal["font_title"] == "Test CJK"
    assert pal["font_body"] == "Test CJK"
    assert pal["font_label"] == "Test CJK"
    assert "Test CJK" in capsys.readouterr().err
    palettes.set_cjk(False)
    assert palettes.get_palette("")["font_body"] == "Calibri"


def test_set_cjk_no_font_keeps_palette_fonts(monkeypatch, capsys):
    import palettes
    monkeypatch.setitem(palettes._CJK, "font", None)
    monkeypatch.setitem(palettes._CJK, "probed", True)
    palettes.set_cjk(True)
    assert palettes.get_palette("")["font_body"] == "Calibri"
    assert "no CJK-safe font" in capsys.readouterr().err


def test_build_swaps_fonts_for_cjk_deck(tmp_path, monkeypatch):
    import palettes
    from build_deck import build
    from pptx import Presentation
    monkeypatch.setitem(palettes._CJK, "font", "Hiragino Sans GB")
    monkeypatch.setitem(palettes._CJK, "probed", True)
    md = tmp_path / "o.md"
    md.write_text("## Slide 1: 三个杠杆在2027财年弥合成本差距\n"
                  "- 分层定价\n- Notes: n.\n", encoding="utf-8")
    out = tmp_path / "d.pptx"
    assert build(str(md), str(out))
    prs = Presentation(str(out))
    fonts = {run.font.name for slide in prs.slides for sh in slide.shapes
             if getattr(sh, "has_text_frame", False)
             for para in sh.text_frame.paragraphs for run in para.runs}
    assert "Hiragino Sans GB" in fonts


def test_build_keeps_fonts_for_latin_deck(tmp_path, monkeypatch):
    import palettes
    from build_deck import build
    from pptx import Presentation
    monkeypatch.setitem(palettes._CJK, "font", "Hiragino Sans GB")
    monkeypatch.setitem(palettes._CJK, "probed", True)
    md = tmp_path / "o.md"
    md.write_text("## Slide 1: Costs have outgrown revenue for years\n"
                  "- a bullet\n- Notes: n.\n", encoding="utf-8")
    out = tmp_path / "d.pptx"
    assert build(str(md), str(out))
    prs = Presentation(str(out))
    fonts = {run.font.name for slide in prs.slides for sh in slide.shapes
             if getattr(sh, "has_text_frame", False)
             for para in sh.text_frame.paragraphs for run in para.runs}
    assert "Hiragino Sans GB" not in fonts


def test_tracker_tabs_skips_auto_agenda_slides(tmp_path):
    from build_deck import build
    from pptx import Presentation
    md = tmp_path / "o.md"
    md.write_text("**Auto-Agenda:** track\n" + TRACKER_MD)
    out = tmp_path / "d.pptx"
    assert build(str(md), str(out))
    prs = Presentation(str(out))
    # slide order: title, agenda, divider, agenda(track), content, ...
    assert not _round_chips(prs.slides[3]), "auto agenda slide gets no tab strip"
    assert len(_round_chips(prs.slides[4])) == 2, "content slide keeps tabs"
