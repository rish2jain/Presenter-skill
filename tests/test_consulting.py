"""Tests for the consulting layout builders (builders_consulting.py)."""
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
sys.path.insert(0, str(SCRIPTS))

from pptx import Presentation  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402

import builders  # noqa: E402
from build_deck import parse_outline, validate  # noqa: E402
from builders_consulting import (build_waterfall_slide, build_harvey_slide,
                                 build_mekko_slide, build_gantt_slide,
                                 _fmt_num)  # noqa: E402
from palettes import get_palette  # noqa: E402

PAL = get_palette("midnight-executive")
CTX = {"outline_dir": Path("."), "assets_dir": Path("assets")}


def _prs():
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
    builders.set_canvas(prs)
    return prs


def _texts(slide):
    return [sh.text_frame.text for sh in slide.shapes
            if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip()]


WATERFALL_MD = """## Slide 1: Three levers bridge run-rate down to target
**Layout:** waterfall
**Data:**
- FY25: 46
- Tiering: -8
- Exit: -6
- Discounts: -4
- FY27: total
"""


def test_waterfall_signed_delta_labels():
    _, slides = parse_outline(WATERFALL_MD)
    slide = build_waterfall_slide(_prs(), slides[0], PAL, CTX)
    texts = _texts(slide)
    # deltas are signed, endpoints are plain (the +8 magnitude bug regression)
    for expected in ("-8", "-6", "-4", "46", "28"):
        assert expected in texts, f"{expected!r} missing from {texts}"
    assert "+8" not in texts and "+46" not in texts


def test_waterfall_total_row_parses_as_sentinel():
    _, slides = parse_outline(WATERFALL_MD)
    assert slides[0]["data"][-1] == ("FY27", "total")


def test_fmt_num_signs():
    assert _fmt_num(46) == "46"
    assert _fmt_num(-8, signed=True) == "-8"
    assert _fmt_num(8, signed=True) == "+8"


HARVEY_MD = """## Slide 1: Vendor B leads on weighted criteria
**Layout:** harvey-scorecard
| Criterion | A | B |
|---|---|---|
| Scale | 2 | 4 |
| Support | 0 | 3 |
"""


def test_harvey_pie_angles_in_ooxml_units():
    """Pie adjustments are degrees*0.6 (OOXML 60000ths) — the x100000 regression."""
    _, slides = parse_outline(HARVEY_MD)
    slide = build_harvey_slide(_prs(), slides[0], PAL, CTX)
    pies = [el for el in slide.shapes._spTree.findall(".//" + qn("a:prstGeom"))
            if el.get("prst") == "pie"]
    assert len(pies) == 2  # scores 2 and 3 (0 = empty ring, 4 = full disc)
    gds = {gd.get("name"): gd.get("fmla") for gd in pies[0].findall(".//" + qn("a:gd"))}
    assert gds["adj1"] == "val -5400000"  # -90 degrees
    # score 2 -> end angle 90 degrees = 5400000
    assert gds["adj2"] == "val 5400000"


def test_appendix_marks_slides_and_relaxes_warnings():
    md = """## Slide 1: Main point stated as a full action title
**Layout:** bullet-list
- one
- Notes: "x"

## Appendix
## Slide 2: Backup
**Layout:** table
| a | b |
|---|---|
| 1 | 2 |
"""
    meta, slides = parse_outline(md)
    assert not slides[0].get("_appendix")
    assert slides[1]["_appendix"] is True
    _, warnings = validate(slides, CTX, meta)
    # appendix slides are exempt from notes + action-title warnings
    assert not any("Slide 2" in w and ("notes" in w or "action title" in w)
                   for w in warnings)


def test_quadrant_keys_only_apply_on_matrix_slides():
    md = """## Slide 1: Quarterly revenue grows every quarter this year
**Layout:** two-column-split
**Visual:** chart:bar
**Data:**
- Q1: 4.2
- Q2: 5.1
"""
    _, slides = parse_outline(md)
    assert slides[0]["data"] == [("Q1", 4.2), ("Q2", 5.1)]
    assert "q1" not in slides[0]

    md2 = """## Slide 1: Matrix
**Layout:** matrix-2x2
- Q1: "Quick wins"
- Item: Name="A" X="0.5" Y="0.5"
"""
    _, slides2 = parse_outline(md2)
    assert slides2[0]["q1"] == "Quick wins"


MEKKO_MD = """## Slide 1: Two platforms capture most of segment value
**Layout:** mekko
**Series:** A, B
**Data:**
- Compute: 18, 22
- Storage: 9, 11
"""


def test_mekko_builds_segments_and_widths():
    _, slides = parse_outline(MEKKO_MD)
    slide = build_mekko_slide(_prs(), slides[0], PAL, CTX)
    texts = _texts(slide)
    assert "40" in texts and "20" in texts  # column totals
    # 4 segments + legend chips + accent bar = rect count sanity
    rects = [el for el in slide.shapes._spTree.findall(".//" + qn("a:prstGeom"))
             if el.get("prst") == "rect"]
    assert len(rects) >= 6


GANTT_MD = """## Slide 1: Three workstreams deliver across four quarters
**Layout:** gantt
**Periods:** Q1, Q2, Q3, Q4
- Bar: Row="Infra" Label="Tiering" Start="1" End="2"
- Bar: Row="Platform" Label="Pilot" Start="2" End="4"
- Milestone: Row="Platform" Label="GA" At="4"
"""


def test_gantt_builds_rows_bars_milestone():
    _, slides = parse_outline(GANTT_MD)
    assert len(slides[0]["bars"]) == 2
    assert slides[0]["periods"] == ["Q1", "Q2", "Q3", "Q4"]
    slide = build_gantt_slide(_prs(), slides[0], PAL, CTX)
    texts = _texts(slide)
    for expected in ("Infra", "Platform", "Tiering", "Pilot", "GA", "Q1", "Q4"):
        assert any(expected in t for t in texts), f"{expected!r} missing"
    diamonds = [el for el in slide.shapes._spTree.findall(".//" + qn("a:prstGeom"))
                if el.get("prst") == "diamond"]
    assert len(diamonds) == 1


def test_gantt_validator_rejects_out_of_range_bars():
    md = GANTT_MD.replace('Start="2" End="4"', 'Start="2" End="9"')
    meta, slides = parse_outline(md)
    errors, _ = validate(slides, CTX, meta)
    assert any("outside 1-4" in e for e in errors)


def test_exhibit_numbering_meta():
    md = """**Exhibits:** on

## Slide 1: Costs rise without intervention this year
**Layout:** bullet-list
- one
- Source: "FinOps"
"""
    meta, slides = parse_outline(md)
    assert meta["exhibits"] == "on"
    assert slides[0]["source"] == "FinOps"


def test_action_title_warning_fires_on_topic_labels():
    md = """## Slide 1: Leadership
**Layout:** bullet-list
- one
- Notes: "x"
"""
    meta, slides = parse_outline(md)
    _, warnings = validate(slides, CTX, meta)
    assert any("action title" in w for w in warnings)


BAR_MEKKO_MD = """## Slide 1: EMEA is the margin outlier despite its size
**Layout:** bar-mekko
- Bar: Label="Americas" Size="55" Value="14"
- Bar: Label="EMEA" Size="30" Value="6"
- Bar: Label="APAC" Size="15" Value="11"
- Notes: Width = revenue share, height = EBITDA margin.
"""


def test_bar_mekko_widths_proportional_to_size():
    from builders_consulting import build_bar_mekko_slide
    _, slides = parse_outline(BAR_MEKKO_MD)
    slide = build_bar_mekko_slide(_prs(), slides[0], PAL, CTX)
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    rects = [s for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.height.inches > 0.5]
    assert len(rects) == 3
    widths = sorted((r.width for r in rects), reverse=True)
    assert widths[0] > widths[1] > widths[2]  # 55 > 30 > 15


def test_bar_mekko_registered_and_validated():
    import builders
    assert "bar-mekko" in builders.LAYOUT_MAP
    meta, slides = parse_outline("## Slide 1: Bad mekko\n**Layout:** bar-mekko\n")
    errors, _ = validate(slides, CTX, meta)
    assert any("bar-mekko" in e for e in errors)

BUBBLE_MD = """## Slide 1: Two bets dominate the portfolio by revenue at stake
**Layout:** matrix-2x2
- X-axis: Relative share
- Y-axis: Market growth
- Item: Name="Stars" X="0.8" Y="0.8" Size="40"
- Item: Name="Dogs" X="0.2" Y="0.2" Size="5"
- Notes: Bubble area = revenue.
"""


def test_matrix_bubble_sizes_scale():
    from builders_consulting import build_matrix_slide
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    _, slides = parse_outline(BUBBLE_MD)
    slide = build_matrix_slide(_prs(), slides[0], PAL, CTX)
    ovals = [s for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             and s.width == s.height and s.width.inches > 0.1]
    assert len(ovals) == 2
    big, small = sorted((o.width.inches for o in ovals), reverse=True)
    assert big > small * 1.8, (big, small)  # d=0.24+0.66*(s/smax)^0.5 → ~1.9x (40 vs 5)


def test_matrix_negative_size_falls_back_to_fixed_dot():
    from builders_consulting import build_matrix_slide
    md = BUBBLE_MD.replace('Size="5"', 'Size="-5"')
    _, slides = parse_outline(md)
    slide = build_matrix_slide(_prs(), slides[0], PAL, CTX)  # must not raise
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    ovals = [s for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.width == s.height]
    assert len(ovals) == 2


def test_bar_mekko_rejects_negative_values_in_validate():
    md = BAR_MEKKO_MD.replace('Size="30"', 'Size="-30"')
    meta, slides = parse_outline(md)
    errors, _ = validate(slides, CTX, meta)
    assert any("Size>0" in e for e in errors), errors

def _fill_of(shape):
    try:
        return str(shape.fill.fore_color.rgb)
    except (AttributeError, TypeError):
        return None


def _fills(slide):
    return [f for f in (_fill_of(s) for s in slide.shapes) if f]


# ── heatmap-table ────────────────────────────────────────────────────────────
HEATMAP_MD = """## Slide 1: EMEA unit costs run 3x the NA baseline
**Layout:** heatmap-table
| Region | Cost | Churn |
|---|---|---|
| NA | 10 | 5 |
| EMEA | 30 | 7 |
| APAC | 20 | 6 |
"""


def test_heatmap_fills_interpolate_per_column():
    from builders_consulting import build_heatmap_slide
    _, slides = parse_outline(HEATMAP_MD)
    slide = build_heatmap_slide(_prs(), slides[0], PAL, CTX)
    fills = _fills(slide)
    assert PAL["bg"] in fills                 # column-min cells
    assert fills.count(PAL["accent1"]) >= 3   # accent bar + 2 column-max cells
    # midpoint cells (20 of 10-30, 6 of 5-7) -> exact 50% bg->accent1 blend
    assert fills.count("6A5C35") == 2, fills


def test_heatmap_rag_scale_uses_rag_palette():
    from builders_consulting import build_heatmap_slide
    md = HEATMAP_MD + "- Scale: rag\n"
    _, slides = parse_outline(md)
    assert slides[0]["scale"] == "rag"
    slide = build_heatmap_slide(_prs(), slides[0], PAL, CTX)
    fills = _fills(slide)
    for key in ("rag_bad", "rag_mid", "rag_good"):
        assert PAL[key] in fills, (key, fills)


def test_heatmap_rag_defaults_on_all_builtin_palettes():
    import re as _re
    from palettes import PALETTES
    for name, pal in PALETTES.items():
        for key in ("rag_bad", "rag_mid", "rag_good"):
            assert _re.fullmatch(r"[0-9A-Fa-f]{6}", pal[key]), (name, key)


def test_rag_colors_contrast_on_surface_all_palettes():
    """rag_bad and rag_good must each achieve >= 3.0:1 contrast on the palette
    surface color (driver-tree boxes render rag delta text at 12pt on surface).
    """
    from palettes import PALETTES

    def _luminance(hex_str):
        h = hex_str.lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        def chan(c):
            c = c / 255
            return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
        return 0.2126 * chan(r) + 0.7152 * chan(g) + 0.0722 * chan(b)

    def _contrast(h1, h2):
        l1, l2 = sorted([_luminance(h1), _luminance(h2)], reverse=True)
        return (l1 + 0.05) / (l2 + 0.05)

    failures = []
    for name, pal in PALETTES.items():
        surf = pal["surface"]
        for key in ("rag_bad", "rag_good"):
            ratio = _contrast(pal[key], surf)
            if ratio < 3.0:
                failures.append(
                    f"{name}: {key}={pal[key]} on surface={surf} → {ratio:.2f}:1 (<3.0)")
    assert not failures, "RAG contrast failures:\n" + "\n".join(failures)


def test_heatmap_validator_requires_numeric_table():
    md = """## Slide 1: Heatmap without any numeric body column must fail
**Layout:** heatmap-table
| A | B |
|---|---|
| x | y |
"""
    meta, slides = parse_outline(md)
    errors, _ = validate(slides, CTX, meta)
    assert any("heatmap-table" in e for e in errors), errors


# ── tornado ──────────────────────────────────────────────────────────────────
TORNADO_MD = """## Slide 1: Margin and WACC swings dominate the 35-point range
**Layout:** tornado
**Series:** Low, High
**Data:**
- WACC: -12, 18
- Growth: -8, 10
- Margin: -15, 20
"""


def _tornado_right_widths(slide):
    """Widths of right-side (accent1) bars, top-to-bottom."""
    bars = sorted((s.top, s.width) for s in slide.shapes
                  if _fill_of(s) == PAL["accent1"]
                  and 0.1 < s.height.inches < 0.6 and s.width.inches < 5)
    return [w for _, w in bars]


def test_tornado_sorts_by_span_and_scales_bars():
    from builders_consulting import build_tornado_slide
    _, slides = parse_outline(TORNADO_MD)
    assert slides[0]["data"][0] == ("WACC", [-12.0, 18.0])
    slide = build_tornado_slide(_prs(), slides[0], PAL, CTX)
    widths = _tornado_right_widths(slide)
    assert len(widths) == 3
    # sorted by |low|+|high| desc: Margin (35), WACC (30), Growth (18)
    assert widths[0] > widths[1] > widths[2]
    assert abs(widths[0] / widths[2] - 2.0) < 0.05  # right values 20 vs 10
    texts = _texts(slide)
    for expected in ("-12", "+18", "Margin", "Low", "High"):
        assert any(expected in t for t in texts), (expected, texts)


def test_tornado_sort_off_preserves_input_order():
    from builders_consulting import build_tornado_slide
    md = TORNADO_MD + "- Sort: off\n"
    _, slides = parse_outline(md)
    slide = build_tornado_slide(_prs(), slides[0], PAL, CTX)
    widths = _tornado_right_widths(slide)
    # input order: WACC (18), Growth (10), Margin (20)
    assert widths[2] > widths[0] > widths[1]


def test_tornado_validator_requires_two_pair_rows():
    md = """## Slide 1: Tornado without a two-name Series declaration fails
**Layout:** tornado
**Data:**
- WACC: -12
"""
    meta, slides = parse_outline(md)
    errors, _ = validate(slides, CTX, meta)
    assert any("tornado" in e for e in errors), errors


# ── football-field ───────────────────────────────────────────────────────────
FOOTBALL_MD = """## Slide 1: Valuation methods converge on the 45-52 band
**Layout:** football-field
**Series:** Low, High
**Data:**
- DCF: 42, 58
- Comparables: 45, 52
- Precedents: 40, 50
- Marker: Current price, 47
"""


def test_football_field_bars_and_marker():
    from builders_consulting import build_football_field_slide
    _, slides = parse_outline(FOOTBALL_MD)
    assert slides[0]["marker"] == "Current price, 47"
    slide = build_football_field_slide(_prs(), slides[0], PAL, CTX)
    rounds = [el for el in slide.shapes._spTree.findall(".//" + qn("a:prstGeom"))
              if el.get("prst") == "roundRect"]
    assert len(rounds) == 3
    assert len(slide.shapes._spTree.findall(qn("p:cxnSp"))) == 1  # marker line
    texts = _texts(slide)
    for expected in ("DCF", "42", "58", "Current price"):
        assert any(expected in t for t in texts), (expected, texts)


def test_football_field_bar_widths_proportional_to_range():
    from builders_consulting import build_football_field_slide
    from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
    _, slides = parse_outline(FOOTBALL_MD)
    slide = build_football_field_slide(_prs(), slides[0], PAL, CTX)
    bars = sorted((s.top, s.width) for s in slide.shapes
                  if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                  and s.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE)
    widths = [w for _, w in bars]  # row order: DCF 16, Comparables 7, Precedents 10
    assert abs(widths[0] / widths[1] - 16 / 7) < 0.05
    assert abs(widths[0] / widths[2] - 16 / 10) < 0.05


def test_football_marker_outside_range_warns_and_skips(capsys):
    md = FOOTBALL_MD.replace("Current price, 47", "Current price, 99")
    meta, slides = parse_outline(md)
    _, warnings = validate(slides, CTX, meta)
    assert any("Marker" in w and "outside" in w for w in warnings), warnings
    from builders_consulting import build_football_field_slide
    slide = build_football_field_slide(_prs(), slides[0], PAL, CTX)
    assert "Marker" in capsys.readouterr().err
    assert not slide.shapes._spTree.findall(qn("p:cxnSp"))


def test_football_validator_rejects_inverted_range():
    md = FOOTBALL_MD.replace("DCF: 42, 58", "DCF: 58, 42")
    meta, slides = parse_outline(md)
    errors, _ = validate(slides, CTX, meta)
    assert any("low < high" in e for e in errors), errors


def test_new_layouts_registered_and_action_titled():
    from build_deck import ACTION_TITLE_LAYOUTS
    for name in ("heatmap-table", "tornado", "football-field"):
        assert name in builders.LAYOUT_MAP
        assert name in ACTION_TITLE_LAYOUTS
        assert name not in builders.GHOST_KEEP_REAL  # ghosts like other exhibits


def test_matrix_bubble_clamped_inside_plot():
    from builders_consulting import build_matrix_slide
    md = BUBBLE_MD.replace('X="0.8" Y="0.8"', 'X="1.0" Y="1.0"')
    _, slides = parse_outline(md)
    slide = build_matrix_slide(_prs(), slides[0], PAL, CTX)
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    EMU_IN = 914400
    plot_right = (2.0 + 9.6) * EMU_IN  # L + W from the builder
    for s in slide.shapes:
        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and s.width == s.height \
                and s.width / EMU_IN > 0.1:
            assert s.left + s.width <= plot_right + 100, "bubble escapes plot"


# ── I1: football-field bar_h floor at 18 rows ────────────────────────────────
def _make_football_md(n_rows):
    rows = "\n".join(f"- Method{i}: {30 + i}, {50 + i}" for i in range(n_rows))
    return (
        f"## Slide 1: Valuation range from {n_rows} methods\n"
        "**Layout:** football-field\n"
        "**Series:** Low, High\n"
        "**Data:**\n"
        f"{rows}\n"
    )


def test_football_field_18_rows_positive_bar_heights():
    """bar_h must not go negative when row_h < 0.22 (18+ rows)."""
    from builders_consulting import build_football_field_slide
    from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
    _, slides = parse_outline(_make_football_md(18))
    slide = build_football_field_slide(_prs(), slides[0], PAL, CTX)
    bars = [s for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and s.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE]
    assert len(bars) == 18
    for bar in bars:
        assert bar.height.inches > 0, f"bar height {bar.height.inches} <= 0"


# ── I3: football-field marker label clamped when near global max ──────────────
def test_football_marker_label_right_edge_clamped():
    """Label right edge must stay <= 12.63in when marker is at the global max."""
    from builders_consulting import build_football_field_slide
    # marker at 58 = global max of the default FOOTBALL_MD data
    md = FOOTBALL_MD.replace("Current price, 47", "Current price, 58")
    _, slides = parse_outline(md)
    slide = build_football_field_slide(_prs(), slides[0], PAL, CTX)
    EMU_IN = 914400
    max_right = 12.63 * EMU_IN
    label_tbs = [s for s in slide.shapes
                 if getattr(s, "has_text_frame", False)
                 and "Current price" in s.text_frame.text]
    assert label_tbs, "marker label text box not found"
    for tb in label_tbs:
        right_edge = (tb.left + tb.width) / EMU_IN
        assert right_edge <= 12.63 + 0.01, (
            f"label right edge {right_edge:.3f}in exceeds 12.63in")


# ── I2: _cell_value treats malformed comma groups as non-numeric ──────────────
def test_cell_value_rejects_bad_comma_grouping():
    """'10,0' must NOT parse as numeric — surface fill should be used."""
    from builders_consulting import _cell_value
    assert _cell_value("10,0") is None, "'10,0' should be non-numeric"
    assert _cell_value("1,000") == 1000.0
    assert _cell_value("1,234,567") == 1234567.0
    assert _cell_value("$1,200") == 1200.0
    assert _cell_value("42%") == 42.0


def test_heatmap_bad_comma_cell_gets_surface_fill():
    """A cell like '10,0' is non-numeric and should receive the surface fill."""
    from builders_consulting import build_heatmap_slide
    md = """## Slide 1: Mixed numeric and bad comma cell
**Layout:** heatmap-table
| Category | Score |
|---|---|
| A | 10,0 |
| B | 20 |
| C | 30 |
"""
    _, slides = parse_outline(md)
    slide = build_heatmap_slide(_prs(), slides[0], PAL, CTX)
    fills = _fills(slide)
    # surface fill must be present (the non-numeric '10,0' cell)
    assert PAL["surface"] in fills, f"expected surface fill for '10,0' cell; fills={fills}"


# ── I4: heatmap with 15 rows builds and warns ────────────────────────────────
def _make_heatmap_md(n_body_rows):
    header = "| Metric | Score |\n|---|---|\n"
    rows = "\n".join(f"| Item{i} | {i * 3} |" for i in range(1, n_body_rows + 1))
    return (
        f"## Slide 1: Dense heatmap with {n_body_rows} rows\n"
        "**Layout:** heatmap-table\n"
        + header + rows + "\n"
    )


def test_heatmap_15_rows_builds_and_warns(capsys):
    """15-row heatmap must build without error and emit a density warning."""
    from builders_consulting import build_heatmap_slide
    _, slides = parse_outline(_make_heatmap_md(15))
    slide = build_heatmap_slide(_prs(), slides[0], PAL, CTX)
    # must produce at least 15 body rects
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    rects = [s for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             and s.height.inches > 0]
    assert len(rects) >= 15
    # warning must have been emitted
    err = capsys.readouterr().err
    assert "15" in err and "heatmap-table" in err, (
        f"expected density warning in stderr; got: {err!r}")


# ── driver-tree ──────────────────────────────────────────────────────────────
DRIVER_MD = """## Slide 1: Pricing drives 60% of the planned revenue uplift
**Layout:** driver-tree
- Node: Id="rev" Label="Revenue" Value="$120M" Delta="+8%" Parent=""
- Node: Id="price" Label="Price" Value="$70M" Delta="+12%" Parent="rev"
- Node: Id="vol" Label="Volume" Value="$50M" Delta="-2%" Parent="rev"
- Node: Id="mix" Label="Mix" Value="$30M" Parent="price"
"""


def test_driver_tree_parses_nodes():
    _, slides = parse_outline(DRIVER_MD)
    nodes = slides[0]["nodes"]
    assert len(nodes) == 4
    assert nodes[0] == {"id": "rev", "label": "Revenue", "value": "$120M",
                        "delta": "+8%", "parent": ""}


def _label_left(slide, label):
    lefts = [s.left for s in slide.shapes
             if getattr(s, "has_text_frame", False)
             and s.text_frame.text == label]
    assert lefts, f"label {label!r} not found"
    return min(lefts)


def test_driver_tree_columns_advance_with_depth():
    from builders_consulting import build_driver_tree_slide
    _, slides = parse_outline(DRIVER_MD)
    slide = build_driver_tree_slide(_prs(), slides[0], PAL, CTX)
    assert _label_left(slide, "Revenue") < _label_left(slide, "Price") \
        < _label_left(slide, "Mix")
    assert _label_left(slide, "Price") == _label_left(slide, "Volume")
    # 3 parent->child edges = 3 elbow connectors
    conns = slide.shapes._spTree.findall(qn("p:cxnSp"))
    assert len(conns) == 3


def _run_colors_of(slide, text):
    out = []
    for s in slide.shapes:
        if getattr(s, "has_text_frame", False) and s.text_frame.text == text:
            for para in s.text_frame.paragraphs:
                out.extend(str(r.font.color.rgb) for r in para.runs)
    return out


def test_driver_tree_delta_colors_by_sign():
    from builders_consulting import build_driver_tree_slide
    _, slides = parse_outline(DRIVER_MD)
    slide = build_driver_tree_slide(_prs(), slides[0], PAL, CTX)
    assert PAL["rag_good"] in _run_colors_of(slide, "+12%")
    assert PAL["rag_bad"] in _run_colors_of(slide, "-2%")


def test_driver_tree_12_nodes_stays_in_canvas():
    rows = ['- Node: Id="r" Label="Root" Value="100" Parent=""']
    rows += [f'- Node: Id="n{i}" Label="Driver {i}" Value="{i}" Parent="r"'
             for i in range(11)]
    md = ("## Slide 1: Eleven drivers decompose the full cost base\n"
          "**Layout:** driver-tree\n" + "\n".join(rows) + "\n")
    from builders_consulting import build_driver_tree_slide
    meta, slides = parse_outline(md)
    errors, _ = validate(slides, CTX, meta)
    assert not errors, errors
    slide = build_driver_tree_slide(_prs(), slides[0], PAL, CTX)
    EMU_IN = 914400
    for s in slide.shapes:
        assert s.left >= 0
        assert (s.left + s.width) / EMU_IN <= 13.33 + 0.01
        assert (s.top + s.height) / EMU_IN <= 7.5 + 0.01


def test_driver_tree_depth3_centering():
    """Depth-3 tree (root -> mid -> leaves): mid boxes are horizontally between
    root and leaves; root center-y ≈ mean of its children's center-y; all boxes
    stay within the 13.33 x 7.5 canvas."""
    from builders_consulting import build_driver_tree_slide
    md = """## Slide 1: Three-level cost decomposition shows labour as the driver
**Layout:** driver-tree
- Node: Id="root" Label="Total Cost" Value="$200M" Parent=""
- Node: Id="mid1" Label="Labour" Value="$120M" Parent="root"
- Node: Id="mid2" Label="Infra" Value="$80M" Parent="root"
- Node: Id="leaf1" Label="Headcount" Value="$90M" Parent="mid1"
- Node: Id="leaf2" Label="Contractors" Value="$30M" Parent="mid1"
- Node: Id="leaf3" Label="Cloud" Value="$50M" Parent="mid2"
- Node: Id="leaf4" Label="Facilities" Value="$30M" Parent="mid2"
"""
    meta, slides = parse_outline(md)
    errors, _ = validate(slides, CTX, meta)
    assert not errors, errors
    slide = build_driver_tree_slide(_prs(), slides[0], PAL, CTX)

    EMU_IN = 914400

    def _boxes(slide):
        """Return {label_text: (left_in, center_y_in)} for all labeled boxes."""
        result = {}
        for s in slide.shapes:
            if not getattr(s, "has_text_frame", False):
                continue
            txt = s.text_frame.text
            left_in = s.left / EMU_IN
            cy_in = (s.top + s.height / 2) / EMU_IN
            result[txt] = (left_in, cy_in)
        return result

    boxes = _boxes(slide)

    # Gather label-box lefts for representative nodes at each depth
    root_left = min(l for (l, _) in [boxes[t] for t in boxes if "Total Cost" in t])
    mid_lefts = [l for (l, _) in [boxes[t] for t in boxes
                                  if t in ("Labour", "Infra")]]
    leaf_lefts = [l for (l, _) in [boxes[t] for t in boxes
                                   if t in ("Headcount", "Contractors",
                                            "Cloud", "Facilities")]]
    assert mid_lefts, "mid-node boxes not found"
    assert leaf_lefts, "leaf boxes not found"
    mid_left = min(mid_lefts)
    leaf_left = min(leaf_lefts)

    # Mid column is between root and leaves (horizontally)
    assert root_left < mid_left < leaf_left, (
        f"mid ({mid_left:.3f}) not between root ({root_left:.3f}) "
        f"and leaf ({leaf_left:.3f})")

    # Root center-y ≈ mean of its direct children's center-y (within 0.15in)
    root_cy = next(cy for (_, cy) in [boxes[t] for t in boxes
                                      if "Total Cost" in t])
    mid_cys = [cy for (_, cy) in [boxes[t] for t in boxes
                                   if t in ("Labour", "Infra")]]
    mean_mid_cy = sum(mid_cys) / len(mid_cys)
    assert abs(root_cy - mean_mid_cy) < 0.15, (
        f"root center_y {root_cy:.3f} not near mean of children {mean_mid_cy:.3f}")

    # All boxes within canvas
    for s in slide.shapes:
        assert s.left >= 0
        assert (s.left + s.width) / EMU_IN <= 13.33 + 0.01
        assert (s.top + s.height) / EMU_IN <= 7.5 + 0.01


def test_driver_tree_validator_rules():
    # two roots
    md = DRIVER_MD.replace('Parent="rev"', 'Parent=""', 1)
    meta, slides = parse_outline(md)
    errors, _ = validate(slides, CTX, meta)
    assert any("exactly one root" in e for e in errors), errors
    # unknown parent
    md = DRIVER_MD.replace('Parent="price"', 'Parent="ghost"')
    meta, slides = parse_outline(md)
    errors, _ = validate(slides, CTX, meta)
    assert any("ghost" in e for e in errors), errors
    # duplicate ids
    md = DRIVER_MD.replace('Id="mix"', 'Id="vol"')
    meta, slides = parse_outline(md)
    errors, _ = validate(slides, CTX, meta)
    assert any("unique" in e for e in errors), errors
    # depth > 3
    md = DRIVER_MD + '- Node: Id="deep" Label="Too deep" Parent="mix"\n'
    meta, slides = parse_outline(md)
    errors, _ = validate(slides, CTX, meta)
    assert any("depth" in e for e in errors), errors
    # > 12 nodes
    extra = "\n".join(f'- Node: Id="x{i}" Label="X{i}" Parent="rev"'
                      for i in range(9))
    meta, slides = parse_outline(DRIVER_MD + extra + "\n")
    errors, _ = validate(slides, CTX, meta)
    assert any("12" in e for e in errors), errors


# ── stakeholder-map ──────────────────────────────────────────────────────────
STAKE_MD = """## Slide 1: Move the CFO from skeptic to sponsor this quarter
**Layout:** stakeholder-map
- Item: Name="CFO" X="0.2" Y="0.9" TargetX="0.8" TargetY="0.9"
- Item: Name="COO" X="0.7" Y="0.6"
- Item: Name="CISO" X="0.4" Y="0.3"
"""


def test_stakeholder_map_default_axes_and_arrow():
    from builders_consulting import build_stakeholder_map_slide
    from pptx.enum.dml import MSO_FILL
    _, slides = parse_outline(STAKE_MD)
    slide = build_stakeholder_map_slide(_prs(), slides[0], PAL, CTX)
    texts = _texts(slide)
    assert any("Support" in t for t in texts)
    assert any("Influence" in t for t in texts)
    conns = slide.shapes._spTree.findall(qn("p:cxnSp"))
    assert len(conns) == 1  # one item has targets
    # arrowhead appended (version-guarded helper; present on this python-pptx)
    assert conns[0].findall(".//" + qn("a:tailEnd"))
    # hollow target circle: oval with background (no) fill
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    hollow = [s for s in slide.shapes
              if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
              and s.width == s.height
              and s.fill.type == MSO_FILL.BACKGROUND]
    assert len(hollow) == 1


def test_stakeholder_map_axis_override():
    from builders_consulting import build_stakeholder_map_slide
    md = STAKE_MD + '- X-axis: "Alignment →"\n'
    _, slides = parse_outline(md)
    slide = build_stakeholder_map_slide(_prs(), slides[0], PAL, CTX)
    texts = _texts(slide)
    assert any("Alignment" in t for t in texts)
    assert not any("Support" in t for t in texts)


def test_stakeholder_map_validator():
    md = """## Slide 1: One stakeholder is not a map of the landscape
**Layout:** stakeholder-map
- Item: Name="CFO" X="0.2" Y="0.9"
"""
    meta, slides = parse_outline(md)
    errors, _ = validate(slides, CTX, meta)
    assert any("stakeholder-map" in e for e in errors), errors
    md2 = STAKE_MD.replace('TargetX="0.8"', 'TargetX="1.4"')
    meta, slides = parse_outline(md2)
    errors, _ = validate(slides, CTX, meta)
    assert any("TargetX" in e for e in errors), errors


def test_stakeholder_map_single_target_key_warns():
    md = STAKE_MD.replace(' TargetY="0.9"', "")
    meta, slides = parse_outline(md)
    _, warnings = validate(slides, CTX, meta)
    assert any("TargetX/TargetY" in w for w in warnings), warnings


# ── raci ─────────────────────────────────────────────────────────────────────
RACI_MD = """## Slide 1: Every workstream has a single accountable owner
**Layout:** raci
| Activity | Alice | Bob | Cara |
|---|---|---|---|
| Design | A | R | C |
| Build | I | A | R |
"""


def test_raci_chips_colored_by_letter():
    from builders_consulting import build_raci_slide
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    _, slides = parse_outline(RACI_MD)
    slide = build_raci_slide(_prs(), slides[0], PAL, CTX)
    ovals = [s for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             and s.width == s.height]
    fills = [_fill_of(o) for o in ovals]
    assert fills.count(PAL["accent1"]) == 2     # R chips
    assert fills.count(PAL["accent2"]) == 2     # A chips
    assert fills.count(PAL["accent3"]) == 1     # C chip
    assert fills.count(PAL["text_muted"]) == 1  # I chip
    texts = _texts(slide)
    for expected in ("Activity", "Alice", "Design", "Responsible"):
        assert any(expected in t for t in texts), (expected, texts)


def test_raci_blank_cell_renders_no_chip():
    from builders_consulting import build_raci_slide
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    md = RACI_MD.replace("| Design | A | R | C |", "| Design | A | R |  |")
    _, slides = parse_outline(md)
    slide = build_raci_slide(_prs(), slides[0], PAL, CTX)
    ovals = [s for s in slide.shapes
             if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
             and s.width == s.height]
    assert len(ovals) == 5  # 6 cells minus the blank one


def test_raci_validator_letters_and_accountable():
    md = RACI_MD.replace("| Design | A | R | C |", "| Design | A | X | C |")
    meta, slides = parse_outline(md)
    errors, _ = validate(slides, CTX, meta)
    assert any("R/A/C/I" in e for e in errors), errors
    # 'A,R' combined cells are invalid too
    md = RACI_MD.replace("| Design | A | R | C |", "| Design | A,R |  | C |")
    meta, slides = parse_outline(md)
    errors, _ = validate(slides, CTX, meta)
    assert any("R/A/C/I" in e for e in errors), errors
    # zero or 2+ Accountable cells warn (not error)
    md = RACI_MD.replace("| Build | I | A | R |", "| Build | I | R | R |")
    meta, slides = parse_outline(md)
    errors, warnings = validate(slides, CTX, meta)
    assert not errors, errors
    assert any("Accountable" in w for w in warnings), warnings
    md = RACI_MD.replace("| Build | I | A | R |", "| Build | A | A | R |")
    meta, slides = parse_outline(md)
    _, warnings = validate(slides, CTX, meta)
    assert any("Accountable" in w for w in warnings), warnings


def test_t5_layouts_registered_and_action_titled():
    from build_deck import ACTION_TITLE_LAYOUTS
    for name in ("driver-tree", "stakeholder-map", "raci"):
        assert name in builders.LAYOUT_MAP
        assert name in ACTION_TITLE_LAYOUTS
        assert name not in builders.GHOST_KEEP_REAL
