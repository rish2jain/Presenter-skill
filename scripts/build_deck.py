#!/usr/bin/env python3
"""
build_deck.py — Build a .pptx from a markdown outline.

Usage:
    python scripts/build_deck.py outline.md --output deck.pptx
        [--palette midnight-executive] [--template template.pptx]
        [--assets-dir assets] [--check]

--check validates the outline (layouts, images, chart data) without building.
Outline format: references/generation-guide.md.
"""
import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from builders import LAYOUT_MAP, ctx_blank_layout
from helpers import add_speaker_notes, parse_visual, resolve_image_path, warn
from narrative import check_unsourced_stats, validate_narrative
from palettes import PALETTES, apply_variant, get_palette, load_custom_palettes
from smart_layout import auto_layout
from template_helpers import (find_best_layout, label_placeholder_names,
                              load_template_config, palette_from_theme,
                              validate_template_mapping, LAYOUT_CATEGORIES)

KV_RE = re.compile(r'(\w+)="([^"]*)"')
NUM_RE = re.compile(r"-?\d[\d,]*\.?\d*")

# "- Key: value" prefixes mapped to slide-dict fields
FIELD_KEYS = {
    "Title": "title", "Subtitle": "subtitle", "Heading": "heading",
    "Headline": "heading", "Caption": "caption", "Contact": "contact",
    "Notes": "notes", "Series": "series", "Source": "source",
    "Left label": "left_label", "Right label": "right_label",
    # consulting layouts
    "Quote": "quote", "Attribution": "attribution", "Context": "context",
    "Value": "value", "Label": "label", "Callout": "callout",
    "Current": "current", "Situation": "situation",
    "Recommendation": "recommendation",
    "Bracket": "bracket", "CAGR": "cagr", "Axis-Max": "axis_max",
    "Value-Line": "value_line", "Labels": "labels_mode",
    "X-axis": "x_axis", "Y-axis": "y_axis",
    "Q1": "q1", "Q2": "q2", "Q3": "q3", "Q4": "q4",
    "Sticker": "sticker", "Kicker": "kicker",
    "Scale": "scale", "Sort": "sort", "Marker": "marker",
}


def _parse_kv(line):
    return {k.lower(): v for k, v in KV_RE.findall(line)}


def _series_count(slide):
    """Number of chart series declared via **Series:** (comma-separated names)."""
    series = slide.get("series", "")
    if not series or "," not in series:
        return 1
    return len([n for n in series.split(",") if n.strip()])


def _parse_data_point(text, n_series=1):
    """'2024: $42B' -> ('2024', 42.0); multi: '2024: 42, 38' -> ('2024', [42.0, 38.0])."""
    label, sep, value = text.partition(":")
    if not sep:
        return None
    if n_series > 1:
        nums = []
        for part in value.split(","):
            m = NUM_RE.search(part)
            if not m:
                return None
            nums.append(float(m.group().replace(",", "")))
        if len(nums) != n_series:
            return None
        return label.strip().strip('"'), nums
    if value.strip().lower() in ("total", "end"):  # waterfall closing bar
        return label.strip().strip('"'), "total"
    m = NUM_RE.search(value)
    if not m:
        return None
    return label.strip().strip('"'), float(m.group().replace(",", ""))


META_KEYS = {"Palette": "palette", "Footer": "footer",
              "Page-Numbers": "page_numbers", "Size": "size",
              "Density": "density", "Purpose": "purpose",
              "Variant": "variant", "Motif": "motif", "Takeaway": "takeaway",
              "Exhibits": "exhibits", "Auto-Agenda": "auto_agenda",
              "Stamp": "stamp", "Scale-Group": "scale_group",
              "Tracker": "tracker"}


def parse_outline(md_text):
    """Parse a markdown outline into (deck_meta, list of slide dicts).

    Deck-level front-matter: `**Palette:** aurora`, `**Footer:** "Confidential"`,
    `**Page-Numbers:** on`, `**Size:** 16:9|4:3` before the first `## Slide`.
    """
    meta, slides, current, in_data = {}, [], None, False

    for lineno, raw in enumerate(md_text.splitlines(), 1):
        line = raw.strip()
        if current is None and line.startswith("**"):
            key, _, value = line.lstrip("*").partition(":**")
            if key.strip() in META_KEYS:
                meta[META_KEYS[key.strip()]] = value.strip().strip('"')
                continue
        if line.startswith("## Appendix"):
            if current:
                slides.append(current)
                current = None
            meta["_appendix_from"] = len(slides)  # subsequent slides are backup
            continue
        if line.startswith("## Slide "):
            if current:
                slides.append(current)
            current = {"bullets": [], "stats": [], "cards": [], "items": [],
                       "data": [], "table_rows": [], "steps": [], "tiles": [],
                       "matrix_items": [], "bars": [], "milestones": [],
                       "nodes": [], "left_bullets": [], "right_bullets": []}
            current["_line"] = lineno
            if "_appendix_from" in meta:
                current["_appendix"] = True
            in_data = False
            m = re.match(r"## Slide \d+:\s*(.*)", line)
            if m:
                heading = m.group(1).strip()
                m_attr = re.search(r"\{([^{}]*)\}\s*$", heading)
                if m_attr:
                    pairs = re.findall(r"(\w[\w-]*)=([\w./#-]+)",
                                       m_attr.group(1))
                    if not pairs:
                        warn(f"line {lineno}: heading attributes "
                             f"{m_attr.group(0)!r} did not parse — values "
                             "must be unquoted (layout=waterfall)")
                    for k, v in pairs:
                        if k == "layout":
                            current["layout"] = v.lower()
                        elif k == "palette":
                            current["palette"] = v.lower()
                        else:
                            warn(f"line {lineno}: unknown heading attribute "
                                 f"'{k}' ignored (supported: layout, palette)")
                    heading = heading[:m_attr.start()].strip()
                current["heading"] = heading
            continue
        if current is None:
            continue
        if not line:
            continue  # blank lines inside **Data:** blocks are ignored

        if line.startswith("**"):
            in_data = False
            key, _, value = line.lstrip("*").partition(":**")
            key, value = key.strip(), value.strip()
            if key == "Layout":
                current["layout"] = value.lower()
            elif key == "Visual":
                current["visual"] = value
            elif key == "Visual-Left":
                current["visual_left"] = value
            elif key == "Visual-Right":
                current["visual_right"] = value
            elif key == "Palette":
                current["palette"] = value.lower()
            elif key == "Series":
                current["series"] = value
            elif key == "Benchmark":
                current["benchmark"] = value
            elif key == "Periods":
                current["periods"] = [s.strip() for s in value.split(",")
                                      if s.strip()]
            elif key == "Data":
                in_data = True
            continue

        if line.startswith("|"):
            cells = [c.strip() for c in line.strip("|").split("|")]
            if not all(set(c) <= {"-", ":", " "} for c in cells):  # skip ---- row
                current["table_rows"].append(cells)
            continue

        if not line.startswith("- "):
            continue
        item = line[2:].strip()

        # Known field keys take precedence over (and terminate) a Data block,
        # so '- Heading: "32% YoY"' after **Data:** isn't read as a data point.
        matched = False
        for prefix, field in FIELD_KEYS.items():
            # Q1-Q4 are quadrant labels only on matrix slides; elsewhere they
            # are ordinary labels (e.g. '- Q1: 4.2, 3.1' chart data rows)
            if field in ("q1", "q2", "q3", "q4") and \
                    current.get("layout") != "matrix-2x2":
                continue
            if item.startswith(prefix + ":"):
                current[field] = item[len(prefix) + 1:].strip().strip('"')
                matched = True
                in_data = False
                break
        if matched:
            continue

        if in_data:
            point = _parse_data_point(item, _series_count(current))
            if point:
                current["data"].append(point)
                continue
            in_data = False  # fall through: not a data line after all

        if item.startswith("Left:"):
            current["left_bullets"].append(item[5:].strip().strip('"'))
        elif item.startswith("Right:"):
            current["right_bullets"].append(item[6:].strip().strip('"'))
        elif re.match(r"Card \d+:", item):
            current["cards"].append(_parse_kv(item))
        elif re.match(r"Step \d+:", item):
            current["steps"].append(_parse_kv(item))
        elif re.match(r"Tile \d+:", item):
            current["tiles"].append(_parse_kv(item))
        elif item.startswith("Item:"):
            current["matrix_items"].append(_parse_kv(item))
        elif item.startswith("Bar:"):
            current["bars"].append(_parse_kv(item))
        elif item.startswith("Milestone:"):
            current["milestones"].append(_parse_kv(item))
        elif item.startswith("Node:"):
            current["nodes"].append(_parse_kv(item))
        elif item.startswith("Value=") or item.startswith("Stat"):
            kv = _parse_kv(item)
            if kv:
                current["stats"].append(kv)
        elif "date=" in item.lower() and KV_RE.search(item):
            current["items"].append(_parse_kv(item))
        elif re.match(r"(Point|Finding) \d+:", item):
            current["bullets"].append(item.split(":", 1)[1].strip().strip('"'))
        else:
            current["bullets"].append(item.strip('"'))

    if current:
        slides.append(current)
    for slide in slides:
        if not slide.get("layout"):
            slide["layout"] = auto_layout(slide)
    return meta, slides


def _agenda_slide(sections, current=None):
    slide = {"bullets": list(sections), "stats": [], "cards": [], "items": [],
             "data": [], "table_rows": [], "steps": [], "tiles": [],
             "matrix_items": [], "bars": [], "milestones": [],
             "nodes": [], "left_bullets": [], "right_bullets": [],
             "layout": "agenda", "heading": "Agenda",
             "notes": "Auto-generated agenda tracker.", "_auto": True}
    if current:
        slide["current"] = current
        slide["heading"] = "Where we are"
    return slide


def apply_auto_agenda(meta, slides):
    """**Auto-Agenda:** on  -> overview agenda after the title slide.
                       track -> + current-highlighted agenda after each divider.
    Sections come from section-divider headings (appendix dividers excluded);
    no dividers -> no-op.  The overview is only inserted when slide 1 is a
    title slide."""
    mode = meta.get("auto_agenda", "").lower()
    if mode not in ("on", "track"):
        return slides
    sections = [s.get("heading") or s.get("title", "")
                for s in slides
                if s.get("layout") == "section-divider" and not s.get("_appendix")]
    if not sections:
        return slides
    out = []
    for i, s in enumerate(slides):
        out.append(s)
        if i == 0 and s.get("layout") == "title" and not s.get("_appendix"):
            out.append(_agenda_slide(sections))
        if mode == "track" and s.get("layout") == "section-divider" \
                and not s.get("_appendix"):
            out.append(_agenda_slide(
                sections, current=s.get("heading") or s.get("title", "")))
    return out


# Chart kinds where sharing one value-axis maximum is meaningful (a value
# axis exists and isn't a fixed 0-100% scale).
SCALE_GROUP_KINDS = ("bar", "column", "hbar", "line", "area")


def apply_scale_groups(meta, slides):
    """**Scale-Group:** auto — slides whose chart visuals share a kind get one
    common axis maximum (honest visual comparison). Slides with an explicit
    Axis-Max are left alone; groups need 2+ slides. Default: off."""
    if meta.get("scale_group", "").lower() != "auto":
        return slides
    groups = {}  # chart kind -> [(slide_no, slide_dict, data_max)]
    for i, p in enumerate(slides, 1):
        kind, value = parse_visual(p.get("visual", ""))
        if kind != "chart" or value not in SCALE_GROUP_KINDS \
                or p.get("axis_max"):
            continue
        nums = []
        for _, v in p.get("data", []):
            if isinstance(v, list):
                nums.extend(x for x in v if isinstance(x, (int, float)))
            elif isinstance(v, (int, float)):
                nums.append(v)
        if not nums:
            continue
        groups.setdefault(value, []).append((i, p, max(nums)))
    from charts import _nice_ceil
    for chart_kind, members in groups.items():
        if len(members) < 2:
            continue
        axis_max = _nice_ceil(max(m for _, _, m in members))
        for _, p, _ in members:
            p["axis_max"] = axis_max
        where = ", ".join(str(i) for i, _, _ in members)
        print(f"  Scale-group: {chart_kind} charts on slides {where} "
              f"share axis max {axis_max:g}")
    return slides


# ── Validation ───────────────────────────────────────────────────────────────
NO_FOOTER_LAYOUTS = {"title", "closing", "section-divider", "full-image"}

# Layouts whose headings should pass the consulting "titles test"
ACTION_TITLE_LAYOUTS = {
    "bullet-list", "two-column-split", "exec-summary", "exec-summary-scqa",
    "comparison", "table", "stat-callout", "timeline", "waterfall",
    "matrix-2x2", "chart-callout", "dashboard", "funnel", "harvey-scorecard",
    "process-flow", "mekko", "gantt", "bar-mekko",
    "heatmap-table", "tornado", "football-field",
    "driver-tree", "stakeholder-map", "raci",
}


def _validate_driver_tree(p, where, errors):
    """Driver-tree: one root, parents exist, unique ids, depth<=3, <=12 nodes."""
    nodes = p.get("nodes", [])
    if not nodes:
        errors.append(f"{where}: driver-tree needs '- Node: Id=\"..\" "
                      'Label=".." Value=".." Parent=".."\' rows')
        return
    ids = [n.get("id", "").strip() for n in nodes]
    if len(nodes) > 12:
        errors.append(f"{where}: driver-tree supports at most 12 nodes "
                      f"(got {len(nodes)})")
    if not all(ids):
        errors.append(f"{where}: every driver-tree Node needs an Id")
        return
    if len(set(ids)) != len(ids):
        errors.append(f"{where}: driver-tree Node Ids must be unique")
        return
    roots = [n for n in nodes if not n.get("parent", "").strip()]
    if len(roots) != 1:
        errors.append(f"{where}: driver-tree needs exactly one root node "
                      f"(empty Parent) — got {len(roots)}")
    known = set(ids)
    bad = sorted({n["parent"] for n in nodes
                  if n.get("parent", "").strip() and n["parent"] not in known})
    if bad:
        errors.append(f"{where}: driver-tree Parent ids not found: "
                      f"{', '.join(bad)}")
    if len(roots) != 1 or bad:
        return
    children = {}
    for n in nodes:
        if n.get("parent", "").strip():
            children.setdefault(n["parent"], []).append(n["id"])
    depth, queue = {roots[0]["id"]: 1}, [roots[0]["id"]]
    while queue:
        nid = queue.pop(0)
        for ch in children.get(nid, []):
            depth[ch] = depth[nid] + 1
            queue.append(ch)
    if len(depth) != len(nodes):
        missing = sorted(set(ids) - set(depth))
        errors.append(f"{where}: driver-tree nodes not reachable from the "
                      f"root (cycle?): {', '.join(missing)}")
    elif max(depth.values()) > 3:
        errors.append(f"{where}: driver-tree depth exceeds 3 levels "
                      "(root + 2 — split the decomposition)")


_TARGET_KEYS = {"targetx": "TargetX", "targety": "TargetY"}


def _validate_stakeholder_map(p, where, errors, warnings):
    """Stakeholder-map: 2+ valid 0-1 items; targets, when given, 0-1 pairs."""
    valid = 0
    for it in p.get("matrix_items", []):
        try:
            x, y = float(it.get("x", "")), float(it.get("y", ""))
            if 0 <= x <= 1 and 0 <= y <= 1:
                valid += 1
        except ValueError:
            pass
        given = [k for k in _TARGET_KEYS if it.get(k)]
        if len(given) == 1:
            warnings.append(f"{where}: Item {it.get('name', '')!r} has only "
                            "one of TargetX/TargetY — arrow will be skipped")
        for k in given:
            try:
                if not 0 <= float(it[k]) <= 1:
                    raise ValueError
            except ValueError:
                errors.append(f"{where}: Item {it.get('name', '')!r} "
                              f"{_TARGET_KEYS[k]} must be a 0-1 float "
                              f"(got {it[k]!r})")
    if valid < 2:
        errors.append(f"{where}: stakeholder-map needs 2+ '- Item: "
                      'Name=".." X="0-1" Y="0-1"\' rows')


def _validate_raci(p, where, errors, warnings):
    """RACI: header + body table; cells R/A/C/I or blank; one 'A' per row."""
    t_rows = p.get("table_rows", [])
    if len(t_rows) < 2 or max((len(r) for r in t_rows), default=0) < 2:
        errors.append(f"{where}: raci needs a markdown table "
                      "(| Activity | Person… | header plus body rows)")
        return
    for row in t_rows[1:]:
        cells = [c.strip().upper() for c in row[1:]]
        for c in cells:
            if c and c not in ("R", "A", "C", "I"):
                errors.append(f"{where}: raci cell {c!r} in row {row[0]!r} "
                              "must be a single R/A/C/I letter or blank")
        n_a = cells.count("A")
        if n_a != 1:
            warnings.append(f"{where}: activity {row[0]!r} has {n_a} 'A' "
                            "cells — exactly one Accountable per activity")


def validate(slides, ctx, meta=None):
    """Fail-fast outline checks. Returns (errors, warnings)."""
    meta = meta or {}
    errors, warnings = [], []
    if not slides:
        return ["Outline contains no '## Slide N:' sections"], []

    deck_pal = meta.get("palette")
    if deck_pal and deck_pal not in PALETTES:
        warnings.append(f"unknown deck palette '{deck_pal}' — default will be used")

    if meta.get("tracker", "").lower() == "tabs" and not any(
            p.get("layout") == "section-divider" and not p.get("_appendix")
            for p in slides):
        warnings.append("**Tracker:** tabs needs section-divider slides — "
                        "falling back to the plain section label")

    for n, p in enumerate(slides, 1):
        where = f"Slide {n} (line {p['_line']})" if p.get("_line") else f"Slide {n}"
        layout = p.get("layout", "bullet-list")
        if layout not in LAYOUT_MAP:
            errors.append(f"{where}: unknown layout '{layout}' "
                          f"(valid: {', '.join(sorted(LAYOUT_MAP))})")
            continue
        if "palette" in p and p["palette"] not in PALETTES:
            warnings.append(f"{where}: unknown palette '{p['palette']}', "
                            "default will be used")

        for vis_key in ("visual", "visual_left", "visual_right"):
            kind, value = parse_visual(p.get(vis_key, ""))
            if not kind:
                continue
            if kind == "image" and not resolve_image_path(value, ctx):
                errors.append(f"{where}: image not found for {vis_key}: '{value}'")
            elif kind == "chart":
                if value not in __import__("charts").CHART_TYPES:
                    errors.append(f"{where}: unknown chart type '{value}'")
                if vis_key == "visual" and not p.get("data"):
                    errors.append(f"{where}: 'chart:{value}' visual but no "
                                  "**Data:** block (or no parsable '- label: value' lines)")
                elif vis_key != "visual":
                    errors.append(f"{where}: charts on {vis_key} are not supported "
                                  "(use **Visual:** on two-column-split)")

        if not p.get("notes") and not p.get("_appendix") \
                and layout != "section-divider":
            warnings.append(f"{where}: missing '- Notes:' speaker notes")

        # Titles test: content-slide headings should be action titles —
        # full-sentence takeaways, not topic labels ("Leadership" fails).
        if (layout in ACTION_TITLE_LAYOUTS and not p.get("_appendix")
                and len(p.get("heading", "").split()) < 5):
            warnings.append(
                f"{where}: heading {p.get('heading', '')!r} is a topic label, "
                "not an action title — state the takeaway as a sentence "
                "(run --titles for the titles test)")

        if layout in ACTION_TITLE_LAYOUTS and not p.get("_appendix"):
            heading = p.get("heading", "")
            if len(heading.split()) > 15:
                warnings.append(f"{where}: heading exceeds 15 words — "
                                "tighten the takeaway")
            chart_ish = parse_visual(p.get("visual", ""))[0] == "chart" or \
                layout in ("waterfall", "mekko", "bar-mekko", "chart-callout")
            if chart_ish and heading and not any(c.isdigit() for c in heading):
                warnings.append(f"{where}: exhibit slide heading has no "
                                "number — quantify the takeaway")
            if " and " in heading.lower():
                warnings.append(f"{where}: heading joins two messages with "
                                "'and' — consider splitting the slide")

        if p.get("kicker"):
            kick_words = {w for w in re.findall(r"[a-z0-9']+",
                                                p["kicker"].lower())
                          if len(w) > 3}
            head_words = {w for w in re.findall(r"[a-z0-9']+",
                                                p.get("heading", "").lower())
                          if len(w) > 3}
            if kick_words and len(kick_words & head_words) / len(kick_words) > 0.6:
                warnings.append(f"{where}: kicker restates the title — "
                                "make it advance the argument")
        if p.get("kicker") and layout in NO_FOOTER_LAYOUTS:
            warnings.append(f"{where}: Kicker is not rendered on '{layout}' "
                            "layouts — move it to a content slide")

        if layout in ("title", "closing") and not p.get("title") and p.get("heading"):
            warnings.append(f"{where}: no '- Title:' — using '## Slide N:' heading "
                            f"({p['heading']!r}) as title text")

        n_series = _series_count(p)
        if n_series > 1 and parse_visual(p.get("visual", ""))[0] == "chart":
            if not p.get("data"):
                errors.append(
                    f"{where}: multi-series chart needs **Data:** rows with "
                    f"{n_series} comma-separated values per line "
                    f"(**Series:** {p.get('series')})")
            else:
                for label, vals in p["data"]:
                    if not isinstance(vals, list) or len(vals) != n_series:
                        errors.append(
                            f"{where}: multi-series row {label!r} needs {n_series} "
                            f"comma-separated values (got **Series:** {p.get('series')})")
                        break

        import builders
        cap = builders.density()["bullet_max"]
        if len(p.get("bullets", [])) > cap:
            warnings.append(f"{where}: {len(p['bullets'])} bullets — only the "
                            f"first {cap} render; split the slide")

        # Overflow guard: estimated text height vs the layout's vertical
        # budget (catches long bullets; the bullet-count cap catches many).
        if layout in ("bullet-list", "bullets", "two-column-split"):
            import textfit
            geo = (builders.twocol_geometry() if layout == "two-column-split"
                   else builders.bullet_geometry())
            shown = p.get("bullets", [])[:geo["max"]]
            if shown:
                ratio, _, overflowing = textfit.bullets_fit(
                    shown, geo["font_pt"], geo["col_w"], geo["avail_h"],
                    cols=geo["cols"], step_in=geo["step"])
                if ratio > 1.4:
                    errors.append(
                        f"{where}: text overflows layout capacity "
                        f"(~{round(ratio * 100)}%) — split the slide or cut copy")
                elif ratio > 1.0:
                    warnings.append(
                        f"{where}: slide ~{round(ratio * 100)}% of text "
                        "capacity — autofit will shrink text")
                # Per-slot collision guard: non-terminal bullets that overflow
                # their fixed y-step slot will render on top of the next bullet.
                if overflowing and geo["step"] > 0:
                    line_h = geo["font_pt"] * 1.2 / 72.0
                    for idx in overflowing:
                        plain, has_icon = textfit.strip_markup(shown[idx])
                        w = geo["col_w"] - (textfit.ICON_INDENT if has_icon
                                            else textfit.MARKER_INDENT)
                        lines = max(
                            textfit.estimate_lines(plain, geo["font_pt"], w), 1)
                        sev = lines * line_h / geo["step"]
                        msg = (
                            f"{where}: bullet {idx + 1} wraps to ~{lines} lines "
                            "and will collide with the next bullet — "
                            "shorten it or split the slide")
                        if sev > 1.8:
                            errors.append(msg)
                        else:
                            warnings.append(msg)
            head = textfit.strip_markup(p.get("heading", ""))[0]
            if head:
                h_w, h_pt = ((6.0, 30) if layout == "two-column-split"
                             else (11.9, 32))
                h_lines = textfit.estimate_lines(head, h_pt, h_w, bold=True)
                if h_lines > 2:
                    warnings.append(f"{where}: title wraps to ~{h_lines} "
                                    "lines — shorten the heading")
        if layout == "stat-callout" and not p.get("stats"):
            errors.append(f"{where}: stat-callout requires "
                          '\'- Value="..." Label="..." Sublabel="..."\' lines')
        if layout == "table" and not p.get("table_rows"):
            errors.append(f"{where}: table layout requires markdown table rows")
        if layout in ("cards-3", "cards-4"):
            want = 3 if layout == "cards-3" else 4
            if len(p.get("cards", [])) != want:
                warnings.append(f"{where}: {layout} has {len(p.get('cards', []))} "
                                f"cards (expected {want})")
        if layout == "timeline" and not p.get("items") and not p.get("bullets"):
            errors.append(f"{where}: timeline needs '- Date=\"..\" Title=\"..\"' "
                          "items or 'Date: text' bullets")
        if layout == "full-image" and parse_visual(p.get("visual", ""))[0] != "image":
            warnings.append(f"{where}: full-image without an image visual "
                            "renders a plain background")

        # consulting layout requirements
        numeric_data = [d for d in p.get("data", [])
                        if isinstance(d[1], (int, float, list))]
        if parse_visual(p.get("visual", ""))[0] == "chart" and \
                any(d[1] == "total" for d in p.get("data", [])):
            errors.append(f"{where}: 'total' rows are only valid on waterfall "
                          "layouts, not charts")
        if layout == "waterfall" and len(p.get("data", [])) < 2:
            errors.append(f"{where}: waterfall needs a **Data:** block — start "
                          "value, +/- deltas, optional '- End label: total' row")
        if layout == "funnel" and len(numeric_data) < 2:
            errors.append(f"{where}: funnel needs a **Data:** block with 2+ "
                          "numeric '- Stage: value' rows")
        if layout == "matrix-2x2" and not p.get("matrix_items"):
            errors.append(f"{where}: matrix-2x2 needs '- Item: Name=\"..\" "
                          "X=\"0-1\" Y=\"0-1\"' rows")
        if layout == "harvey-scorecard" and len(p.get("table_rows", [])) < 2:
            errors.append(f"{where}: harvey-scorecard needs a markdown table "
                          "(header row + criteria rows, cells 0-4)")
        if layout in ("process-flow", "next-steps") and not p.get("steps") \
                and not p.get("bullets"):
            errors.append(f"{where}: {layout} needs '- Step N: ...' rows")
        if layout == "big-number" and not p.get("value"):
            errors.append(f"{where}: big-number needs '- Value: \"$4.2B\"'")
        if layout == "chart-callout":
            if not p.get("callout"):
                errors.append(f"{where}: chart-callout needs '- Callout: \"...\"'")
            if not p.get("data") and parse_visual(p.get("visual", ""))[0] != "image":
                errors.append(f"{where}: chart-callout needs a chart **Data:** "
                              "block or an image visual")
        if layout == "dashboard" and not p.get("tiles"):
            errors.append(f"{where}: dashboard needs '- Tile N: Value=\"..\" "
                          "Label=\"..\" Delta=\"..\" Trend=\"up|down\"' rows")
        if layout == "quote-evidence" and not p.get("quote"):
            errors.append(f"{where}: quote-evidence needs '- Quote: \"...\"'")
        if layout == "exec-summary-scqa" and not (
                p.get("situation") and p.get("recommendation")):
            errors.append(f"{where}: exec-summary-scqa needs '- Situation:' and "
                          "'- Recommendation:' (findings via '- Finding N:')")
        if layout == "agenda" and not p.get("bullets"):
            errors.append(f"{where}: agenda needs section names as bullets")
        if layout == "mekko":
            if _series_count(p) < 2 or not any(
                    isinstance(d[1], list) for d in p.get("data", [])):
                errors.append(f"{where}: mekko needs **Series:** with 2+ names "
                              "and multi-value **Data:** rows")
        if layout == "gantt":
            if not p.get("periods"):
                errors.append(f"{where}: gantt needs '**Periods:** Q1, Q2, ...'")
            if not p.get("bars"):
                errors.append(f"{where}: gantt needs '- Bar: Row=\"..\" "
                              "Label=\"..\" Start=\"1\" End=\"2\"' rows")
            for bar in p.get("bars", []):
                try:
                    s, e = float(bar.get("start", 1)), float(bar.get("end", 1))
                    if not (1 <= s <= e <= len(p.get("periods", [])) ):
                        errors.append(f"{where}: Bar {bar.get('label', '')!r} "
                                      f"Start/End outside 1-{len(p['periods'])}")
                except ValueError:
                    errors.append(f"{where}: Bar {bar.get('label', '')!r} has "
                                  "non-numeric Start/End")
        if layout == "bar-mekko":
            ok_bars = [b for b in p.get("bars", [])
                       if "size" in b and "value" in b]
            if len(ok_bars) < 2:
                errors.append(f"{where}: bar-mekko needs 2+ '- Bar: "
                              'Label=".." Size=".." Value=".."\' rows')
            for bar in ok_bars:
                try:
                    size, value = float(bar["size"]), float(bar["value"])
                    if size <= 0 or value < 0:
                        errors.append(f"{where}: Bar {bar.get('label', '')!r} "
                                      "needs Size>0 and Value>=0")
                except ValueError:
                    errors.append(f"{where}: Bar {bar.get('label', '')!r} has "
                                  "non-numeric Size/Value")
        if layout == "heatmap-table":
            from builders_consulting import _cell_value
            t_rows = p.get("table_rows", [])
            n_cols = max((len(r) for r in t_rows), default=0)
            numeric = any(_cell_value(c) is not None
                          for r in t_rows[1:] for c in r)
            if len(t_rows) < 2 or n_cols < 2 or not numeric:
                errors.append(f"{where}: heatmap-table needs a markdown table "
                              "(header + body rows, 2+ columns) with at least "
                              "one numeric body column")
        if layout in ("tornado", "football-field"):
            pairs = [d for d in p.get("data", [])
                     if isinstance(d[1], list) and len(d[1]) == 2]
            if len(pairs) < 2:
                errors.append(f"{where}: {layout} needs '**Series:** Low, "
                              "High' and 2+ '- Name: low, high' **Data:** rows")
            elif layout == "football-field":
                for lbl, (lo, hi) in pairs:
                    if not lo < hi:
                        errors.append(f"{where}: football-field row {lbl!r} "
                                      f"needs low < high (got {lo:g}, {hi:g})")
                if p.get("marker"):
                    from builders_consulting import _parse_marker
                    parsed = _parse_marker(p["marker"])
                    gmin = min(v[0] for _, v in pairs)
                    gmax = max(v[1] for _, v in pairs)
                    if not parsed:
                        warnings.append(
                            f"{where}: Marker needs 'label, value' (got "
                            f"{p['marker']!r}) — it will be skipped")
                    elif not gmin <= parsed[1] <= gmax:
                        warnings.append(
                            f"{where}: Marker value {parsed[1]:g} outside "
                            f"[{gmin:g}, {gmax:g}] — it will be skipped")
        if layout == "driver-tree":
            _validate_driver_tree(p, where, errors)
        if layout == "stakeholder-map":
            _validate_stakeholder_map(p, where, errors, warnings)
        if layout == "raci":
            _validate_raci(p, where, errors, warnings)

    # Narrative + data integrity
    warnings.extend(validate_narrative(meta, slides))
    warnings.extend(check_unsourced_stats(slides))

    template_path = ctx.get("template_path")
    if template_path and Path(template_path).is_file():
        from pptx import Presentation
        config = load_template_config(template_path)
        prs = Presentation(str(template_path))
        warnings.extend(validate_template_mapping(prs, slides, config))

    # Layout variety: 3+ identical consecutive layouts reads as monotonous
    # (appendix/backup slides are deliberately dense and repetitive — exempt)
    layouts = [p.get("layout", "bullet-list") for p in slides
               if not p.get("_appendix")]
    run_start = 0
    for i in range(1, len(layouts) + 1):
        if i == len(layouts) or layouts[i] != layouts[run_start]:
            if i - run_start >= 3 and layouts[run_start] not in ("full-image",):
                warnings.append(
                    f"Slides {run_start + 1}-{i}: '{layouts[run_start]}' used "
                    f"{i - run_start}x consecutively — vary layouts (cards, "
                    "stat-callout, two-column-split...)")
            run_start = i
    return errors, warnings


# ── Template-aware build ─────────────────────────────────────────────────────
TEMPLATE_WANT = LAYOUT_CATEGORIES


def _fill_placeholders(slide, p):
    """Fill title/body/subtitle placeholders from the slide dict."""
    from pptx.enum.shapes import PP_PLACEHOLDER
    title_text = p.get("title") or p.get("heading", "")
    body_lines = p.get("bullets", [])
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        ph_type = ph.placeholder_format.type
        if idx == 0 or ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            ph.text = title_text
        elif ph_type == PP_PLACEHOLDER.SUBTITLE or (idx == 1 and not body_lines):
            ph.text = p.get("subtitle", "")
        elif ph.has_text_frame and body_lines:
            tf = ph.text_frame
            tf.clear()
            for i, line in enumerate(body_lines):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = line
            body_lines = []  # only fill the first body placeholder


def build_template_slide(prs, p, layout_name, ctx, config=None):
    """Try a placeholder-based slide from the user's template."""
    from pptx.enum.shapes import PP_PLACEHOLDER as PH

    kind, value = parse_visual(p.get("visual", ""))
    if layout_name == "two-column-split" and kind == "image":
        layout, _ = find_best_layout(prs, "picture", config)
        if layout is not None:
            path = resolve_image_path(value, ctx)
            if path:
                slide = prs.slides.add_slide(layout)
                label_placeholder_names(slide, layout)
                _fill_placeholders(slide, p)
                for ph in slide.placeholders:
                    if ph.placeholder_format.type == PH.PICTURE:
                        ph.insert_picture(str(path))
                        break
                return slide

    want = TEMPLATE_WANT.get(layout_name)
    if not want:
        return None
    layout, score = find_best_layout(prs, want, config)
    if layout is None or score < 1:
        return None
    slide = prs.slides.add_slide(layout)
    label_placeholder_names(slide, layout)
    _fill_placeholders(slide, p)
    return slide


# ── Main ─────────────────────────────────────────────────────────────────────
SIZES = {"16:9": (13.33, 7.5), "4:3": (10.0, 7.5)}


def create_presentation(template_path=None, size="16:9"):
    if template_path:
        return Presentation(template_path)
    prs = Presentation()
    w, h = SIZES.get(size, SIZES["16:9"])
    prs.slide_width = Inches(w)
    prs.slide_height = Inches(h)
    return prs


def _add_footer(slide, i, pal, footer, page_numbers, appendix=False):
    import builders
    if footer:
        builders.add_tb(slide, footer, 0.7, 7.08, 6.0, 0.35, size=11,
                        color=pal["text_muted"], font=pal["font_label"])
    if page_numbers:
        from pptx.enum.text import PP_ALIGN
        label = f"B·{i + 1}" if appendix else str(i + 1)
        builders.add_tb(slide, label, 11.9, 7.08, 0.8, 0.35, size=11,
                        color=pal["text_muted"], font=pal["font_label"],
                        align=PP_ALIGN.RIGHT)


def _add_section_tracker(slide, pal, section):
    """Section tracker: small current-section label top-right (consulting
    navigation convention for long decks)."""
    if not section:
        return
    import builders
    from pptx.enum.text import PP_ALIGN
    builders.add_tb(slide, section.upper(), 7.6, 0.12, 5.0, 0.32, size=11,
                    bold=True, color=pal["text_muted"], font=pal["font_label"],
                    align=PP_ALIGN.RIGHT)


# Tab strip geometry (design-space inches; builders rescale for 4:3 etc.)
TAB_RIGHT_EDGE = 12.63
TAB_Y, TAB_H, TAB_GAP = 0.12, 0.30, 0.08
TAB_STRIP_MAX_W = 5.5   # over this (or 6+ sections) -> compact text form
TAB_LABEL_CHARS = 14


def _tab_label(name):
    return name[:TAB_LABEL_CHARS] + "…" if len(name) > TAB_LABEL_CHARS else name


def _add_tracker_tabs(slide, pal, sections, current, state):
    """**Tracker:** tabs — chip-per-section strip top-right; current section
    chip filled accent1, others outlined muted. Too many sections -> compact
    'n/N · Section' text at the plain-label position (warn once per build)."""
    import builders
    from pptx.enum.text import PP_ALIGN
    labels = [_tab_label(s) for s in sections]
    widths = [max(0.7, 0.09 * len(lb) + 0.25) for lb in labels]
    total = sum(widths) + TAB_GAP * (len(widths) - 1)
    if len(sections) >= 6 or total > TAB_STRIP_MAX_W:
        if not state.get("warned"):
            warn(f"Tracker tabs: {len(sections)} sections exceed tab budget "
                 "— using compact text form")
            state["warned"] = True
        n = sections.index(current) + 1 if current in sections else 0
        builders.add_tb(slide, f"{n}/{len(sections)} · {current}",
                        7.6, 0.12, 5.0, 0.32, size=11, bold=True,
                        color=pal["text_muted"], font=pal["font_label"],
                        align=PP_ALIGN.RIGHT)
        return
    x = TAB_RIGHT_EDGE
    for label, w, sec in reversed(list(zip(labels, widths, sections))):
        x -= w
        if sec == current:
            builders.add_round_rect(slide, x, TAB_Y, w, TAB_H, pal["accent1"])
            builders.add_tb(slide, label, x, TAB_Y + 0.02, w, TAB_H - 0.02,
                            size=10, bold=True, color=pal["bg"],
                            font=pal["font_label"], align=PP_ALIGN.CENTER)
        else:
            chip = builders.add_round_rect(slide, x, TAB_Y, w, TAB_H,
                                           pal["bg"],
                                           line_hex=pal["text_muted"],
                                           line_pt=0.75)
            chip.fill.background()  # outline only — works on any background
            builders.add_tb(slide, label, x, TAB_Y + 0.02, w, TAB_H - 0.02,
                            size=10, color=pal["text_muted"],
                            font=pal["font_label"], align=PP_ALIGN.CENTER)
        x -= TAB_GAP


def _add_sticker(slide, pal, text, below_section=False):
    """Per-slide bordered status tag (ILLUSTRATIVE, PRELIMINARY...) top-right;
    drops below the section tracker label when both are present."""
    import builders
    from pptx.enum.text import PP_ALIGN
    label = text.upper()
    w = min(3.0, max(1.2, 0.12 * len(label)))
    left = 12.63 - w
    top = 0.50 if below_section else 0.14
    box = builders.add_rect(slide, left, top, w, 0.32, pal["bg"],
                            line_hex=pal["accent3"], line_pt=1.0)
    box.fill.background()  # outline only — works on photos and gradients
    builders.add_tb(slide, label, left, top + 0.02, w, 0.28, size=11,
                    bold=True, color=pal["text_muted"], font=pal["font_label"],
                    align=PP_ALIGN.CENTER)


def _add_kicker_box(slide, pal, text):
    """Takeaway box: full-width accent-edged band above the source/footer."""
    import builders
    builders.add_rect(slide, 0.7, 6.18, 11.93, 0.52, pal["surface"])
    builders.add_rect(slide, 0.7, 6.18, 0.06, 0.52, pal["accent1"])
    builders.add_tb(slide, text, 0.95, 6.25, 11.4, 0.4, size=14, bold=True,
                    color=pal["text"], font=pal["font_body"])


def _add_source(slide, pal, source, exhibit_no=None):
    """Exhibit attribution footnote (every consulting exhibit cites a source).
    With **Exhibits:** on, sourced slides are numbered 'Exhibit N · Source: ...'."""
    import builders
    text = source if source.lower().startswith("source") else f"Source: {source}"
    if exhibit_no:
        text = f"Exhibit {exhibit_no} · {text}"
    builders.add_tb(slide, text, 0.7, 6.80, 9.5, 0.3, size=11,
                    color=pal["text_muted"], font=pal["font_label"])


def _add_stamp(slide, pal, text):
    """Bordered status tag (DRAFT, CONFIDENTIAL) top-left, above the heading."""
    import builders
    from pptx.enum.text import PP_ALIGN
    box = builders.add_rect(slide, 0.7, 0.14, 1.5, 0.32, pal["bg"],
                            line_hex=pal["accent3"], line_pt=1.0)
    box.fill.background()  # outline only — works on photos and gradients
    builders.add_tb(slide, text.upper(), 0.7, 0.16, 1.5, 0.28, size=11,
                    bold=True, color=pal["text_muted"], font=pal["font_label"],
                    align=PP_ALIGN.CENTER)


def build(outline_path, output_path, palette_key=None,
          template_path=None, assets_dir=None, check_only=False, size=None,
          density=None, variant=None, ghost=False):
    outline_path = Path(outline_path)
    ctx = {
        "outline_dir": outline_path.parent,
        "assets_dir": Path(assets_dir) if assets_dir else Path("assets"),
        "template_path": template_path,
    }
    custom = load_custom_palettes(ctx["assets_dir"] / "palettes")
    if custom:
        print(f"  Loaded custom palettes: {', '.join(custom)}")
    if palette_key and palette_key not in PALETTES:
        print(f"  [WARN] unknown --palette '{palette_key}' — using default",
              file=sys.stderr)
    meta, slides_data = parse_outline(outline_path.read_text(encoding="utf-8"))
    slides_data = apply_auto_agenda(meta, slides_data)
    slides_data = apply_scale_groups(meta, slides_data)

    import builders
    variant_key = variant or meta.get("variant", "")
    palette_key, density = apply_variant(variant_key, palette_key, density)
    builders.set_density(density or meta.get("density", ""))
    errors, warnings = validate(slides_data, ctx, meta)
    for w in warnings:
        print(f"  [WARN] {w}", file=sys.stderr)
    if errors:
        for e in errors:
            print(f"  [ERROR] {e}", file=sys.stderr)
        print(f"\nOutline validation failed: {len(errors)} error(s).", file=sys.stderr)
        return False
    if check_only:
        print(f"Outline OK: {len(slides_data)} slides, "
              f"{len(warnings)} warning(s).")
        return True

    # precedence: CLI flag > outline front-matter > variant preset > default
    pal_default = get_palette(palette_key or meta.get("palette", ""))
    template_config = None
    if template_path:
        template_config = load_template_config(template_path)
        from profile_template import extract_theme_colors
        from pptx import Presentation as _Prs
        theme = extract_theme_colors(_Prs(str(template_path)))
        theme_pal = palette_from_theme(theme)
        if theme_pal:
            for k, v in theme_pal.items():
                if not k.startswith("_"):
                    pal_default[k] = v
    if meta.get("motif"):
        pal_default = {**pal_default, "motif": meta["motif"]}
    size = size or meta.get("size", "16:9")
    footer = meta.get("footer", "")
    page_numbers = meta.get("page_numbers", "").lower() in ("on", "true", "yes")
    exhibits_on = meta.get("exhibits", "").lower() in ("on", "true", "yes")
    exhibit_no = 0

    prs = create_presentation(template_path, size)
    import builders
    builders.set_canvas(prs)  # rescale styled builders to the actual slide size

    failures = 0
    built_slides = []
    current_section = None
    sections = [s.get("heading") or s.get("title", "") for s in slides_data
                if s.get("layout") == "section-divider"
                and not s.get("_appendix")]
    tracker_tabs = meta.get("tracker", "").lower() == "tabs" and bool(sections)
    tab_state = {}  # warn-once flag for the tab-budget fallback
    for i, p in enumerate(slides_data):
        pal = get_palette(p["palette"]) if "palette" in p else pal_default
        layout_name = p.get("layout", "bullet-list")
        if layout_name == "section-divider":
            current_section = p.get("heading") or p.get("title", "")
        print(f"  Building slide {i + 1}: "
              f"{p.get('heading', p.get('title', '?'))} [{layout_name}]")
        try:
            slide = None
            templated = False
            if template_path:
                slide = build_template_slide(prs, p, layout_name, ctx, template_config)
                templated = slide is not None
            if slide is None:
                if ghost and layout_name not in builders.GHOST_KEEP_REAL:
                    slide = builders.build_ghost_slide(prs, p, pal, ctx)
                else:
                    slide = LAYOUT_MAP[layout_name](prs, p, pal, ctx)
            if p.get("notes"):
                add_speaker_notes(slide, p["notes"])
            section_label = "BACKUP" if p.get("_appendix") else current_section
            if not templated and layout_name not in NO_FOOTER_LAYOUTS:
                _add_footer(slide, i, pal, footer, page_numbers,
                            appendix=p.get("_appendix", False))
                if tracker_tabs and section_label \
                        and not p.get("_appendix") and not p.get("_auto"):
                    _add_tracker_tabs(slide, pal, sections, section_label,
                                      tab_state)
                else:
                    _add_section_tracker(slide, pal, section_label)
                if p.get("kicker"):
                    _add_kicker_box(slide, pal, p["kicker"])
            if not templated and p.get("sticker"):
                _add_sticker(slide, pal, p["sticker"],
                             below_section=bool(section_label)
                             and layout_name not in NO_FOOTER_LAYOUTS)
            if not templated and p.get("source"):
                exhibit_no += 1
                _add_source(slide, pal, p["source"],
                            exhibit_no if exhibits_on else None)
            if not templated and meta.get("stamp"):
                _add_stamp(slide, pal, meta["stamp"])
            built_slides.append(slide)
        except Exception as e:
            failures += 1
            print(f"  [ERROR] Slide {i + 1} failed: {e}", file=sys.stderr)
            import traceback
            traceback.print_exc()

    if failures:
        print(f"\nBuild aborted: {failures} slide(s) failed — no file written.",
              file=sys.stderr)
        return False

    prs.save(output_path)
    print(f"\nSaved: {output_path} ({len(built_slides)} slides)")
    return True


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Build a .pptx from a markdown outline.")
    parser.add_argument("outline", help="Path to outline.md")
    parser.add_argument("--output", default="deck.pptx", help="Output .pptx path")
    parser.add_argument("--palette", default=None,
                        help="Color palette: built-ins "
                             f"({', '.join(sorted(PALETTES))}) or a custom "
                             "palette JSON name from <assets>/palettes/")
    parser.add_argument("--template", default=None, help="Optional .pptx template")
    parser.add_argument("--assets-dir", default=None,
                        help="Root dir for user-images/ and auto/ (default: ./assets)")
    parser.add_argument("--size", default=None, choices=sorted(SIZES),
                        help="Slide aspect (default 16:9)")
    parser.add_argument("--density", default=None,
                        choices=("compact", "comfortable"),
                        help="Content density (default compact)")
    parser.add_argument("--variant", default=None,
                        choices=("a", "b", "c", "consulting"),
                        help="Palette+density preset (a=default, b=aurora/comfortable)")
    parser.add_argument("--ghost", action="store_true",
                        help="Build a skeleton deck: real titles, placeholder exhibits (storyline alignment)")
    parser.add_argument("--check", action="store_true",
                        help="Validate the outline only; do not build")
    parser.add_argument("--titles", action="store_true",
                        help="Print slide titles in order (the consulting "
                             "'titles test' — they should read as an argument)")
    args = parser.parse_args()

    if args.titles:
        _, _slides = parse_outline(Path(args.outline).read_text(encoding="utf-8"))
        print("Titles test — read top to bottom; it should work as an argument:\n")
        for n, s in enumerate(_slides, 1):
            tag = " [BACKUP]" if s.get("_appendix") else ""
            print(f"  {n:2d}. {s.get('title') or s.get('heading', '?')}{tag}")
        sys.exit(0)

    ok = build(args.outline, args.output, args.palette,
               args.template, args.assets_dir, args.check, args.size,
               args.density, args.variant, ghost=args.ghost)
    sys.exit(0 if ok else 1)
