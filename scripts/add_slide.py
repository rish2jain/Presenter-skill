#!/usr/bin/env python3
"""Add a blank slide from a template layout index to an existing .pptx."""
import argparse
import sys
from pathlib import Path

from pptx import Presentation


def add_slide(pptx_in, pptx_out, layout_index):
    prs = Presentation(str(pptx_in))
    if layout_index < 0 or layout_index >= len(prs.slide_layouts):
        print(f"ERROR: layout index {layout_index} out of range "
              f"(0-{len(prs.slide_layouts) - 1})", file=sys.stderr)
        return False
    layout = prs.slide_layouts[layout_index]
    prs.slides.add_slide(layout)
    prs.save(str(pptx_out))
    print(f"Added slide using layout {layout_index} ({layout.name!r}) "
          f"→ {pptx_out} ({len(prs.slides)} slides)")
    return True


def main():
    parser = argparse.ArgumentParser(description="Add slide from layout index.")
    parser.add_argument("input", help="Input .pptx")
    parser.add_argument("output", help="Output .pptx")
    parser.add_argument("--layout", type=int, default=1,
                        help="Slide layout index (default 1 = title+content)")
    args = parser.parse_args()
    ok = add_slide(args.input, args.output, args.layout)
    sys.exit(0 if ok else 1)


if __name__ == "__main__":
    main()
