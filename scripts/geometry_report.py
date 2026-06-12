#!/usr/bin/env python3
"""
geometry_report.py — deterministic per-slide layout metrics for LLM
self-verification *before* any render (cheap; no LibreOffice needed).

Per slide:
  (a) pairwise overlaps of visible shapes, z-order aware (containment =
      card/background pattern, ignored; labels drawn over bars/cards and
      solid-on-solid layering = intentional, suppressed; text-on-text and
      text hidden under a solid shape = reported whenever the estimated
      glyph areas — not just the oversized text boxes — actually cross)
  (b) gap consistency: uneven spacing within shape rows/columns (runs are
      split at large gaps — separate visual groups judged independently)
  (c) near-misses: edges almost aligned (off by < 0.08in), between
      structural (non-text) shapes only — a text frame's visible edge is
      its inset glyph edge, not its box edge; pairs sharing the opposite
      edge of the same axis are data-driven sizes (bar lengths), skipped
  (d) whitespace ratio + left/right and top/bottom visual-mass imbalance
  (e) word count (> 90 words = text overload)

Usage:
    python3 scripts/geometry_report.py deck.pptx [--json] [--slides 2,5]

Human output lists only slides with findings; --json emits full metrics
for every slide.
"""
import argparse
import json
import math
import statistics
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR

from qa_check import iter_shapes

EMU_IN = 914400
OVERLAP_MIN_SQIN = 0.02     # report overlaps above this area
SHAPE_OVERLAP_SQIN = 0.3    # textless pairs need this much to be reported
CONTAIN_FRAC = 0.95         # >=95% of the smaller shape inside = containment
CLUSTER_TOL_IN = 0.15       # row/column clustering tolerance
SIZE_RATIO = 1.3            # gap runs require similar cross-axis size
GAP_SPREAD_IN = 0.05        # stdev of gaps above this = uneven spacing
MIN_GAP_IN = 0.02           # smaller gaps = touching run, not a spaced row
GAP_SPLIT_IN = 1.5          # a gap this large splits a run into groups
NEAR_MISS_MIN_IN = 0.04     # below this = invisible at render / builder inset
NEAR_MISS_MAX_IN = 0.08     # edges off by less than this = near-miss
CO_ALIGN_TOL_IN = 0.005     # edges this close = exactly co-aligned
NEAR_MISS_CAP = 4           # max near-miss lines reported per slide
INSET_LR_IN = 0.1           # default text-frame left/right inset
INSET_TB_IN = 0.05          # default text-frame top/bottom inset
DEFAULT_PT = 18.0           # font size assumed when no run specifies one
CHAR_W_EM = 0.55            # mean glyph width as a fraction of font size
LINE_SPACING = 1.3          # line height as a fraction of font size
GRID_STEP_IN = 0.1          # rasterization resolution for coverage
MAX_WORDS = 90              # words per slide before "text overload"
CROWDED_WS = 0.20           # whitespace ratio below this = crowded
IMBALANCE_PP = 60.0         # half-vs-half coverage gap (percentage points)


def _is_solid(shape):
    """Explicit solid (or gradient) fill — occludes whatever is below."""
    try:
        return shape.fill.type in (MSO_FILL.SOLID, MSO_FILL.GRADIENT)
    except Exception:
        return False


def _margin_in(value, default):
    return value / EMU_IN if value is not None else default


def _glyph_box(shape, l, t, w, h):
    """Estimated box the rendered text actually occupies, in inches.
    Text frames carry insets and rarely fill their box, so classifying
    overlaps on raw bounds reads breathing room as collisions.  Height =
    per-paragraph wrapped-line estimate from run font sizes; honors the
    vertical anchor.  Any surprise falls back to the full box."""
    try:
        tf = shape.text_frame
        ml = _margin_in(tf.margin_left, INSET_LR_IN)
        mr = _margin_in(tf.margin_right, INSET_LR_IN)
        mt = _margin_in(tf.margin_top, INSET_TB_IN)
        mb = _margin_in(tf.margin_bottom, INSET_TB_IN)
        avail_w = max(w - ml - mr, 0.1)
        text_h = 0.0
        for para in tf.paragraphs:
            sizes = [r.font.size.pt for r in para.runs if r.font.size]
            pt = max(sizes) if sizes else DEFAULT_PT
            line_w = max(len(para.text), 1) * CHAR_W_EM * pt / 72.0
            lines = max(int(math.ceil(line_w / avail_w)), 1)
            text_h += lines * pt / 72.0 * LINE_SPACING
        box_h = max(h - mt - mb, 0.0)
        text_h = min(text_h, box_h)
        anchor = tf.vertical_anchor
        if anchor == MSO_ANCHOR.MIDDLE:
            gt = t + mt + (box_h - text_h) / 2
        elif anchor == MSO_ANCHOR.BOTTOM:
            gt = t + mt + box_h - text_h
        else:
            gt = t + mt
        return {"l": l + ml, "t": gt, "w": avail_w, "h": text_h}
    except Exception:
        return {"l": l, "t": t, "w": w, "h": h}


def _boxes(slide):
    """Visible leaf shapes as dicts of inches (+ name, words, txt, solid),
    in spTree document order — i.e. z-order, lowest first."""
    out = []
    for shape in iter_shapes(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            continue
        try:
            l, t, w, h = shape.left, shape.top, shape.width, shape.height
        except (AttributeError, TypeError):
            continue
        if None in (l, t, w, h) or w <= 0 or h <= 0:
            continue
        text = shape.text_frame.text \
            if getattr(shape, "has_text_frame", False) else ""
        box = {"name": shape.name or str(shape.shape_type),
               "l": l / EMU_IN, "t": t / EMU_IN,
               "w": w / EMU_IN, "h": h / EMU_IN,
               "words": len(text.split()), "txt": bool(text.strip()),
               "solid": _is_solid(shape)}
        box["glyph"] = _glyph_box(shape, box["l"], box["t"],
                                  box["w"], box["h"]) if box["txt"] else None
        out.append(box)
    return out


def _inter_area(a, b):
    w = min(a["l"] + a["w"], b["l"] + b["w"]) - max(a["l"], b["l"])
    h = min(a["t"] + a["h"], b["t"] + b["h"]) - max(a["t"], b["t"])
    return w * h if w > 0 and h > 0 else 0.0


def _contained(a, b):
    """a is (near-)fully inside b."""
    area_a = a["w"] * a["h"]
    return area_a > 0 and _inter_area(a, b) >= CONTAIN_FRAC * area_a


def find_overlaps(boxes):
    """Z-order-aware overlap classification (boxes arrive lowest-z first).
    Text-on-text and text hidden under a solid shape are the defects —
    judged on estimated glyph areas, not the oversized text boxes; a
    label drawn over a bar/card/band and solid-on-solid layering are
    the standard intentional patterns and are suppressed."""
    out = []
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            lo, hi = boxes[i], boxes[j]  # hi drawn on top of lo
            area = _inter_area(lo, hi)
            if area <= OVERLAP_MIN_SQIN:
                continue
            if _contained(lo, hi) or _contained(hi, lo):
                continue  # card/background pattern — intentional
            if lo["txt"] and hi["txt"]:
                if _inter_area(lo["glyph"], hi["glyph"]) <= OVERLAP_MIN_SQIN:
                    continue  # boxes touch but rendered text cannot
            elif lo["txt"]:
                if not hi["solid"]:
                    continue  # unfilled frame over text hides nothing
                if _inter_area(lo["glyph"], hi) <= OVERLAP_MIN_SQIN:
                    continue  # solid shape misses the rendered text
            elif hi["txt"]:
                continue  # label over bar/card/band — intentional
            elif hi["solid"] or area <= SHAPE_OVERLAP_SQIN:
                continue  # deliberate layering / negligible decoration
            out.append({"a": lo["name"], "b": hi["name"],
                        "area_sqin": round(area, 2)})
    return out


def _top_level(boxes):
    """Boxes not contained inside any other box (card children excluded)."""
    return [a for a in boxes
            if not any(a is not b and _contained(a, b) for b in boxes)]


def _cluster(boxes, key):
    clusters = []
    for box in sorted(boxes, key=lambda b: b[key]):
        if clusters and box[key] - clusters[-1][-1][key] <= CLUSTER_TOL_IN:
            clusters[-1].append(box)
        else:
            clusters.append([box])
    return clusters


def _size_runs(cl, size):
    """Split a cluster into runs of similar cross-axis size: card grids and
    stacks are uniform; title/body/footer chains are not (and aren't rows)."""
    runs, run = [], [cl[0]]
    for box in cl[1:]:
        lo, hi = sorted((box[size], run[-1][size]))
        if lo > 0 and hi / lo <= SIZE_RATIO:
            run.append(box)
        else:
            runs.append(run)
            run = [box]
    runs.append(run)
    return runs


def _gap_groups(run, pos, size):
    """Split a run at any gap > GAP_SPLIT_IN: a large jump means separate
    visual groups, judged for even spacing independently."""
    groups, group = [], [run[0]]
    for box in run[1:]:
        if box[pos] - (group[-1][pos] + group[-1][size]) > GAP_SPLIT_IN:
            groups.append(group)
            group = [box]
        else:
            group.append(box)
    groups.append(group)
    return groups


def _gap_findings(clusters, pos, size, label):
    out = []
    for cl in clusters:
        if len(cl) < 3:
            continue
        for run in _size_runs(sorted(cl, key=lambda b: b[pos]), size):
            for group in _gap_groups(run, pos, size):
                if len(group) < 3:
                    continue
                gaps = [round(group[i + 1][pos]
                              - (group[i][pos] + group[i][size]), 3)
                        for i in range(len(group) - 1)]
                if any(g < MIN_GAP_IN for g in gaps):
                    continue  # touching/overlapping run, not a spaced row
                if statistics.pstdev(gaps) > GAP_SPREAD_IN:
                    seq = "/".join(f"{g:.2f}" for g in gaps)
                    out.append({"axis": label, "gaps": gaps,
                                "text":
                                f"uneven {label} spacing: gaps {seq}in"})
    return out


def find_uneven_gaps(boxes):
    top = _top_level(boxes)
    return (_gap_findings(_cluster(top, "t"), "l", "w", "row")
            + _gap_findings(_cluster(top, "l"), "t", "h", "column"))


_EDGES = (("left", lambda b: b["l"]),
          ("right", lambda b: b["l"] + b["w"]),
          ("top", lambda b: b["t"]),
          ("bottom", lambda b: b["t"] + b["h"]))
_EDGE_FN = dict(_EDGES)
_OPPOSITE = {"left": "right", "right": "left",
             "top": "bottom", "bottom": "top"}


def find_near_misses(boxes):
    """Edge near-misses between structural (non-text) shapes.  Text
    frames align by their inset glyph edge, not their box edge, so
    box-edge deltas this small are invisible for them.  Pairs exactly
    co-aligned on the opposite edge of the same axis are skipped too:
    there the difference is a data-driven width or height (bars off a
    shared spine/baseline), not a misalignment."""
    top = _top_level(boxes)
    out = []
    for i in range(len(top)):
        for j in range(i + 1, len(top)):
            a, b = top[i], top[j]
            if a["txt"] or b["txt"]:
                continue  # glyph edge, not box edge, is what aligns
            for edge, f in _EDGES:
                d = round(abs(f(a) - f(b)), 3)
                if not NEAR_MISS_MIN_IN < d < NEAR_MISS_MAX_IN:
                    continue
                opp = _EDGE_FN[_OPPOSITE[edge]]
                if abs(opp(a) - opp(b)) <= CO_ALIGN_TOL_IN:
                    continue  # shared opposite edge = data-sized pair
                out.append({"a": a["name"], "b": b["name"], "edge": edge,
                            "off_in": round(d, 2)})
    return out


def coverage_metrics(boxes, sw_in, sh_in):
    """Whitespace ratio and half-vs-half visual-mass imbalance, via a
    ~0.1in pixel grid so overlaps aren't double-counted."""
    nx = max(int(round(sw_in / GRID_STEP_IN)), 1)
    ny = max(int(round(sh_in / GRID_STEP_IN)), 1)
    grid = [[False] * nx for _ in range(ny)]
    for b in boxes:
        x0 = max(int(b["l"] / GRID_STEP_IN), 0)
        x1 = min(int((b["l"] + b["w"]) / GRID_STEP_IN) + 1, nx)
        y0 = max(int(b["t"] / GRID_STEP_IN), 0)
        y1 = min(int((b["t"] + b["h"]) / GRID_STEP_IN) + 1, ny)
        for y in range(y0, y1):
            row = grid[y]
            for x in range(x0, x1):
                row[x] = True
    covered = sum(sum(row) for row in grid)
    half_x, half_y = nx // 2, ny // 2

    def pct(cells, total):
        return 100.0 * cells / total if total else 0.0

    left = sum(sum(row[:half_x]) for row in grid)
    right = sum(sum(row[half_x:]) for row in grid)
    top = sum(sum(row) for row in grid[:half_y])
    bottom = sum(sum(row) for row in grid[half_y:])
    return {
        "whitespace_ratio": round(1 - covered / (nx * ny), 3),
        "lr_imbalance_pp": round(abs(
            pct(left, half_x * ny) - pct(right, (nx - half_x) * ny)), 1),
        "tb_imbalance_pp": round(abs(
            pct(top, half_y * nx) - pct(bottom, (ny - half_y) * nx)), 1),
    }


def analyze_slide(slide, sw_in, sh_in):
    boxes = _boxes(slide)
    metrics = {
        "shapes": len(boxes),
        "words": sum(b["words"] for b in boxes),
        "overlaps": find_overlaps(boxes),
        "uneven_gaps": find_uneven_gaps(boxes),
        "near_misses": find_near_misses(boxes),
    }
    metrics.update(coverage_metrics(boxes, sw_in, sh_in))
    return metrics


def analyze_deck(prs, only=None):
    """slide_no -> metrics dict; only = set of slide numbers or None."""
    sw_in, sh_in = prs.slide_width / EMU_IN, prs.slide_height / EMU_IN
    return {n: analyze_slide(slide, sw_in, sh_in)
            for n, slide in enumerate(prs.slides, 1)
            if not only or n in only}


def findings(metrics):
    """Human-readable finding lines for one slide's metrics."""
    out = []
    for o in metrics["overlaps"]:
        out.append(f"overlap: '{o['a']}' x '{o['b']}' "
                   f"by {o['area_sqin']:.2f} sq in")
    for g in metrics["uneven_gaps"]:
        out.append(g["text"])
    misses = sorted(metrics["near_misses"], key=lambda m: m["off_in"])
    for m in misses[:NEAR_MISS_CAP]:
        out.append(f"almost aligned: '{m['a']}' vs '{m['b']}' {m['edge']} "
                   f"edges off by {m['off_in']:.2f}in")
    if len(misses) > NEAR_MISS_CAP:
        out.append(f"(+{len(misses) - NEAR_MISS_CAP} more near-miss "
                   f"alignments suppressed)")
    if metrics["words"] > MAX_WORDS:
        out.append(f"text overload: {metrics['words']} words (>{MAX_WORDS})")
    if metrics["shapes"] and metrics["whitespace_ratio"] < CROWDED_WS:
        out.append(f"crowded: whitespace ratio "
                   f"{metrics['whitespace_ratio']:.2f} (<{CROWDED_WS})")
    if metrics["lr_imbalance_pp"] > IMBALANCE_PP:
        out.append(f"visual mass imbalance: left/right coverage differs by "
                   f"{metrics['lr_imbalance_pp']:.0f}pp")
    if metrics["tb_imbalance_pp"] > IMBALANCE_PP:
        out.append(f"visual mass imbalance: top/bottom coverage differs by "
                   f"{metrics['tb_imbalance_pp']:.0f}pp")
    return out


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx")
    parser.add_argument("--json", action="store_true",
                        help="Emit full metrics for all slides as JSON")
    parser.add_argument("--slides", default=None,
                        help="Comma-separated slide numbers (e.g. 2,5)")
    args = parser.parse_args()
    only = None
    if args.slides:
        try:
            only = {int(s) for s in args.slides.split(",") if s.strip()}
        except ValueError:
            print(f"invalid --slides value: {args.slides!r}")
            sys.exit(2)
    prs = Presentation(args.pptx)
    report = analyze_deck(prs, only)
    if args.json:
        print(json.dumps(report, indent=2))
        return
    total = 0
    for n, metrics in report.items():
        lines = findings(metrics)
        if not lines:
            continue
        total += len(lines)
        print(f"Slide {n}:")
        for line in lines:
            print(f"  - {line}")
    print(f"\n{Path(args.pptx).name}: {len(report)} slide(s) analyzed, "
          f"{total} finding(s)")


if __name__ == "__main__":
    main()
