#!/usr/bin/env python3
"""
diff_deck.py — Compare outline content to built .pptx text (content QA).

Usage:
    python3 scripts/diff_deck.py outline.md deck.pptx
    python3 scripts/diff_deck.py outline.md deck.pptx --strict  # missing text = error

Uses markitdown when installed; falls back to python-pptx text extraction.
"""
import argparse
import re
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent
sys.path.insert(0, str(ROOT))

from build_deck import parse_outline  # noqa: E402

PLACEHOLDER_RX = re.compile(
    r"lorem|ipsum|xxxx|click to add|\[ add visual \]|this page layout", re.I)


def _normalize(text):
    return re.sub(r"\s+", " ", (text or "").lower()).strip()


def _slide_expected_text(slide):
    parts = []
    layout = slide.get("layout", "")
    for key in ("title", "subtitle", "caption", "contact"):
        if slide.get(key):
            parts.append(slide[key])
    # "## Slide N: Title" section labels are not deck copy — skip on title/closing.
    if slide.get("heading") and layout not in ("title", "closing"):
        parts.append(slide["heading"])
    parts.extend(slide.get("bullets", []))
    for stat in slide.get("stats", []):
        parts.extend(str(v) for v in stat.values())
    for card in slide.get("cards", []):
        parts.extend(str(v) for v in card.values())
    for row in slide.get("table_rows", []):
        parts.extend(row)
    return [_normalize(p) for p in parts if p and len(_normalize(p)) > 3]


def extract_pptx_texts(pptx_path):
    """Return list of per-slide text blobs."""
    try:
        from markitdown import MarkItDown
        md = MarkItDown()
        result = md.convert(str(pptx_path), extract_pages=True)
        if getattr(result, "pages", None):
            return [_normalize(p.content) for p in result.pages]
    except ImportError:
        pass
    except Exception as e:
        print(f"  [WARN] markitdown failed: {e}", file=sys.stderr)

    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    def iter_shapes(shapes):
        for shape in shapes:
            yield shape
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from iter_shapes(shape.shapes)

    prs = Presentation(str(pptx_path))
    texts = []
    for slide in prs.slides:
        chunks = []
        for shape in iter_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                t = shape.text_frame.text.strip()
                if t:
                    chunks.append(t)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                chunks.append(notes)
        texts.append(_normalize(" ".join(chunks)))
    return texts


def diff_outline(outline_path, pptx_path, strict=False):
    _, slides = parse_outline(Path(outline_path).read_text(encoding="utf-8"))
    built = extract_pptx_texts(pptx_path)
    errors, warnings = [], []

    if len(built) != len(slides):
        warnings.append(
            f"Slide count mismatch: outline {len(slides)} vs deck {len(built)}")

    for n, slide in enumerate(slides, 1):
        where = f"Slide {n}"
        expected = _slide_expected_text(slide)
        blob = built[n - 1] if n <= len(built) else ""
        if PLACEHOLDER_RX.search(blob):
            errors.append(f"{where}: leftover placeholder text in deck")
        for phrase in expected[:8]:  # key phrases only
            if phrase and phrase not in blob:
                msg = f"{where}: expected text not found in deck: {phrase[:50]!r}"
                (errors if strict else warnings).append(msg)

    return errors, warnings


def main():
    parser = argparse.ArgumentParser(description="Diff outline vs built deck text.")
    parser.add_argument("outline")
    parser.add_argument("pptx")
    parser.add_argument("--strict", action="store_true",
                        help="Treat missing expected text as errors")
    args = parser.parse_args()

    errors, warnings = diff_outline(args.outline, args.pptx, args.strict)
    for w in warnings:
        print(f"  [WARN]  {w}")
    for e in errors:
        print(f"  [ERROR] {e}")
    print(f"\n{Path(args.pptx).name}: {len(errors)} error(s), "
          f"{len(warnings)} warning(s)")
    sys.exit(1 if errors else 0)


if __name__ == "__main__":
    main()
