"""Shared slide-building primitives: text, shapes, images, overlays, notes."""
import sys
from pathlib import Path

from lxml import etree
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

from palettes import hex_rgb

EMU_PER_INCH = 914400


# ── Text & shapes ────────────────────────────────────────────────────────────
import re as _re

_RICH_RX = _re.compile(r"(\*\*.+?\*\*|\{accent\}.+?\{/\})")


def parse_rich_segments(text):
    """Split '**bold** and {accent}colored{/}' into (text, bold, accented) runs."""
    segments = []
    for part in _RICH_RX.split(text):
        if not part:
            continue
        if part.startswith("**") and part.endswith("**"):
            segments.append((part[2:-2], True, False))
        elif part.startswith("{accent}") and part.endswith("{/}"):
            segments.append((part[8:-3], False, True))
        else:
            segments.append((part, False, False))
    return segments or [(text, False, False)]


def add_tb(slide, text, left, top, w, h, size=16, bold=False,
           color="F1F5F9", align=PP_ALIGN.LEFT, font="Calibri", accent=None):
    """Text box with inline rich-text support: **bold** and {accent}...{/}
    (accent runs use the `accent` hex color; ignored if accent is None)."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for seg_text, seg_bold, seg_accent in parse_rich_segments(text):
        run = p.add_run()
        run.text = seg_text
        run.font.size = Pt(size)
        run.font.bold = bold or seg_bold
        run.font.color.rgb = hex_rgb(accent if (seg_accent and accent) else color)
        run.font.name = font
    return tb


def _add_filled_shape(slide, shape_type, left, top, w, h, fill_hex,
                      line_hex=None, line_pt=0):
    shape = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = hex_rgb(line_hex)
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_rect(slide, left, top, w, h, fill_hex, line_hex=None, line_pt=0):
    return _add_filled_shape(slide, 1, left, top, w, h, fill_hex, line_hex, line_pt)


def add_circle(slide, left, top, d, fill_hex, line_hex=None, line_pt=0):
    return _add_filled_shape(slide, 9, left, top, d, d, fill_hex, line_hex, line_pt)


_SHADOW_XML = (
    '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    '<a:outerShdw blurRad="76200" dist="38100" dir="5400000" rotWithShape="0">'
    '<a:srgbClr val="000000"><a:alpha val="32000"/></a:srgbClr>'
    "</a:outerShdw></a:effectLst>"
)


def add_soft_shadow(shape):
    """Subtle downward outer shadow for cards (depth without clutter)."""
    spPr = shape._element.spPr
    for old in spPr.findall(qn("a:effectLst")):
        spPr.remove(old)
    spPr.append(etree.fromstring(_SHADOW_XML))


def set_slide_bg_gradient(slide, hex_from, hex_to, angle_deg=115):
    """Two-stop linear gradient slide background (hero slides on dark themes)."""
    fill = slide.background.fill
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].color.rgb = hex_rgb(hex_from)
    stops[-1].color.rgb = hex_rgb(hex_to)
    try:
        fill.gradient_angle = angle_deg
    except ValueError:
        pass  # non-linear default gradient; colors still applied


def set_fill_alpha(shape, alpha_pct):
    """Make a solid-filled shape semi-transparent (alpha_pct = opacity 0-100)."""
    for srgb in shape._element.findall(".//" + qn("a:solidFill") + "/" + qn("a:srgbClr")):
        for existing in srgb.findall(qn("a:alpha")):
            srgb.remove(existing)
        alpha = etree.SubElement(srgb, qn("a:alpha"))
        alpha.set("val", str(int(alpha_pct * 1000)))


def add_overlay(slide, left, top, w, h, fill_hex="000000", alpha_pct=50):
    """Semi-transparent overlay rectangle (for caption readability on images)."""
    shape = add_rect(slide, left, top, w, h, fill_hex)
    set_fill_alpha(shape, alpha_pct)
    return shape


# ── Images ───────────────────────────────────────────────────────────────────
def _image_size(img_path):
    from PIL import Image
    with Image.open(img_path) as im:
        return im.size


def set_alt_text(picture, description):
    picture._element._nvXxPr.cNvPr.set("descr", description)


def send_to_back(slide, shape):
    sp_tree = slide.shapes._spTree
    sp_tree.remove(shape._element)
    sp_tree.insert(2, shape._element)


def add_picture_contain(slide, img_path, left, top, w, h, alt=None):
    """Place image inside the box (inches), preserving aspect ratio, centered."""
    ow, oh = _image_size(img_path)
    ratio = ow / oh
    tw, th = w, h
    if tw / th > ratio:
        tw = th * ratio
    else:
        th = tw / ratio
    pic = slide.shapes.add_picture(
        str(img_path),
        Inches(left + (w - tw) / 2), Inches(top + (h - th) / 2),
        Inches(tw), Inches(th),
    )
    if alt:
        set_alt_text(pic, alt)
    return pic


def add_picture_cover(slide, prs, img_path, alt=None):
    """Full-bleed background image: fills the slide, crops overflow (no distortion)."""
    ow, oh = _image_size(img_path)
    img_ratio = ow / oh
    slide_ratio = prs.slide_width / prs.slide_height

    pic = slide.shapes.add_picture(str(img_path), 0, 0,
                                   prs.slide_width, prs.slide_height)
    if img_ratio > slide_ratio:
        crop = 1 - slide_ratio / img_ratio
        pic.crop_left = crop / 2
        pic.crop_right = crop / 2
    elif img_ratio < slide_ratio:
        crop = 1 - img_ratio / slide_ratio
        pic.crop_top = crop / 2
        pic.crop_bottom = crop / 2
    if alt:
        set_alt_text(pic, alt)
    send_to_back(slide, pic)
    return pic


def _picture_blip(picture):
    return picture._element.find(qn("p:blipFill") + "/" + qn("a:blip"))


def set_picture_alpha(picture, pct):
    """Make a picture semi-transparent (pct = opacity 0-100) via a:alphaModFix."""
    blip = _picture_blip(picture)
    if blip is None:
        return
    for existing in blip.findall(qn("a:alphaModFix")):
        blip.remove(existing)
    fix = etree.SubElement(blip, qn("a:alphaModFix"))
    fix.set("amt", str(int(pct * 1000)))


def set_picture_duotone(picture, hex_dark, hex_light):
    """Brand-tint a photo: a:duotone mapping shadows->hex_dark, lights->hex_light."""
    blip = _picture_blip(picture)
    if blip is None:
        return
    for existing in blip.findall(qn("a:duotone")):
        blip.remove(existing)
    duo = etree.SubElement(blip, qn("a:duotone"))
    for hex_color in (hex_dark, hex_light):
        clr = etree.SubElement(duo, qn("a:srgbClr"))
        clr.set("val", hex_color)


# ── Visual spec resolution ───────────────────────────────────────────────────
VISUAL_SUBDIRS = {
    "user-image": "user-images",
    "auto-image": "auto",
    "hero-image": "",
}


def parse_visual(spec):
    """Split a **Visual:** value into (kind, value), options stripped.

    kinds: 'chart' (value = chart type), 'image' (value = path spec), None.
    Silent on malformed options — parse_visual_opts warns; this wrapper is
    called repeatedly during validation and must not spam stderr.
    """
    kind, value, _ = parse_visual_opts(spec, warn_opts=False)
    return kind, value


def parse_visual_opts(spec, warn_opts=True):
    """Split a **Visual:** value into (kind, value, opts).

    Image specs take pipe-separated options after the path:
    'user-image:photo.png|alpha=85|duotone'.  Supported opts:
      alpha=N   -> {"alpha": N}  (opacity percent, 0-100)
      duotone   -> {"duotone": True}  (brand-tint via palette bg/accent1)
    Malformed or unknown options are ignored with a warning.
    """
    spec = (spec or "").strip()
    if not spec or spec.lower() == "none":
        return None, None, {}
    if spec.startswith("chart:"):
        return "chart", spec.split(":", 1)[1].strip().lower(), {}
    value, *raw_opts = (part.strip() for part in spec.split("|"))
    opts = {}
    for opt in raw_opts:
        low = opt.lower()
        if low == "duotone":
            opts["duotone"] = True
        elif low.startswith("alpha="):
            try:
                pct = int(low.split("=", 1)[1])
                if not 0 <= pct <= 100:
                    raise ValueError
                opts["alpha"] = pct
            except ValueError:
                if warn_opts:
                    warn(f"image option {opt!r} ignored — alpha must be "
                         "0-100 (e.g. photo.png|alpha=85)")
        elif opt and warn_opts:
            warn(f"unknown image option {opt!r} ignored "
                 "(supported: alpha=N, duotone)")
    return "image", value, opts


def resolve_image_path(spec, ctx):
    """Resolve 'user-image:foo.png' / 'auto-image:x.png' / plain paths to a file.

    Tries (in order): literal path, relative to outline dir, relative to
    assets dir (with source-specific subdir), relative to CWD assets/.
    Returns Path or None.
    """
    spec = spec.strip()
    if ":" in spec and spec.split(":", 1)[0] in VISUAL_SUBDIRS:
        kind, name = spec.split(":", 1)
        subdir = VISUAL_SUBDIRS[kind]
    else:
        kind, name, subdir = None, spec, ""
    name = name.strip()

    assets_dir = Path(ctx.get("assets_dir", "assets"))
    outline_dir = Path(ctx.get("outline_dir", "."))
    candidates = [
        Path(name),
        outline_dir / name,
        assets_dir / subdir / name if subdir else assets_dir / name,
        assets_dir / name,
        outline_dir / "assets" / subdir / name if subdir else outline_dir / "assets" / name,
    ]
    for c in candidates:
        if c.is_file():
            return c
    return None


def add_speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def warn(msg):
    print(f"  [WARN] {msg}", file=sys.stderr)
