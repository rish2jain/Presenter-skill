"""Consulting layout builders: waterfall, 2x2 matrix, harvey scorecard,
process flow, big number, chart callout, dashboard, quote+evidence, funnel,
next-steps, agenda, SCQA executive summary.

Registered into builders.LAYOUT_MAP at import time (see bottom of builders.py).
All geometry is in design inches on the 13.33 x 7.5 canvas; builders' wrappers
rescale to the actual slide size.
"""
import re

from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import builders as B
from charts import round_to_sum
from helpers import warn

NEG_HEX = "D9655B"   # waterfall decreases / down-trends (readable on all palettes)
POS_HEX = None       # increases use pal accent2


# ── waterfall ────────────────────────────────────────────────────────────────
def build_waterfall_slide(prs, p, pal, ctx):
    """Value bridge: first bar = start, middle bars = +/- deltas, a 'total'
    row (or the computed end) closes the bridge."""
    slide = B._blank_slide(prs, pal, pal["bg"])
    B._heading(slide, p, pal)
    data = [(lbl, val) for lbl, val in p.get("data", [])]
    if len(data) < 2:
        warn("waterfall has <2 data rows; rendering bullet-list")
        return B.build_bullet_slide(prs, p, pal, ctx)

    has_total = isinstance(data[-1][1], str)
    deltas = data[1:-1] if has_total else data[1:]
    start_val = data[0][1]
    end_val = start_val + sum(v for _, v in deltas)
    # bar = (label, lo, hi, color, display_value); endpoints show their plain
    # value, deltas show the SIGNED change (-8, not the +8 bar magnitude)
    bars = ([(data[0][0], 0, start_val, pal["accent1"], start_val)]
            + _delta_bars(start_val, deltas, pal)
            + [((data[-1][0] if has_total else "Total"), 0, end_val,
                pal["accent1"], end_val)])

    top, bottom = 1.9, 6.15
    peak = max(max(b[1], b[2]) for b in bars) or 1
    scale = (bottom - top) / peak
    n = len(bars)
    gap = 0.18
    bar_w = (11.9 - gap * (n - 1)) / n

    # connector between bar i and i+1 sits at the RUNNING TOTAL after bar i
    # (bottom of a decrease bar, top of an increase bar — never floats)
    cum_after = [start_val]
    for _, v in deltas:
        cum_after.append(cum_after[-1] + v)
    cum_after.append(end_val)

    for i, (label, lo, hi, color, shown) in enumerate(bars):
        x = 0.7 + i * (bar_w + gap)
        y0 = bottom - hi * scale
        h = max((hi - lo) * scale, 0.04)
        B.add_rect(slide, x, y0, bar_w, h, color)
        is_delta = 0 < i < n - 1
        B.add_tb(slide, _fmt_num(shown, signed=is_delta),
                 x - 0.1, y0 - 0.34, bar_w + 0.2, 0.3, size=12, bold=True,
                 color=pal["text"], align=PP_ALIGN.CENTER, font=pal["font_body"])
        B.add_tb(slide, label, x - 0.12, bottom + 0.10, bar_w + 0.24, 0.75,
                 size=11, color=pal["text_muted"], align=PP_ALIGN.CENTER,
                 font=pal["font_label"])
        if i > 0:
            level_y = bottom - cum_after[i - 1] * scale
            B.add_rect(slide, x - gap, level_y - 0.009, gap, 0.018,
                       pal["text_muted"])
    if p.get("bracket"):
        _bracket_on_waterfall(slide, pal, p, bars, bottom, scale, bar_w, gap)
    B.add_rect(slide, 0.7, bottom, 11.9, 0.02, pal["surface"])
    return slide


def _delta_bars(start, deltas, pal):
    bars, cum = [], start
    for label, v in deltas:
        lo, hi = (cum, cum + v) if v >= 0 else (cum + v, cum)
        bars.append((label, lo, hi,
                     pal["accent2"] if v >= 0 else NEG_HEX, v))
        cum += v
    return bars


def _fmt_num(v, signed=False):
    sign = "+" if signed else ""
    if abs(v) >= 1000:
        return f"{v:{sign},.0f}"
    if v != int(v):
        return f"{v:{sign}.1f}"
    return f"{int(v):{sign}d}"


def _bracket_on_waterfall(slide, pal, p, bars, bottom, scale, bar_w, gap):
    """Difference bracket between two named bars: ┌────┐ + centered label."""
    spec = p.get("bracket", "")
    parts = [s.strip().strip('"') for s in spec.split(",") if s.strip()]
    if len(parts) < 2:
        warn(f"Bracket needs two bar labels, got: {spec!r}")
        return
    labels = [b[0].lower() for b in bars]
    try:
        ia, ib = labels.index(parts[0].lower()), labels.index(parts[1].lower())
    except ValueError:
        warn(f"Bracket labels not found among bars: {spec!r}")
        return
    (_, _, ha, _, va), (_, _, hb, _, vb) = bars[ia], bars[ib]
    xa = 0.7 + ia * (bar_w + gap) + bar_w / 2
    xb = 0.7 + ib * (bar_w + gap) + bar_w / 2
    y = max(min(bottom - ha * scale, bottom - hb * scale) - 0.66, 1.55)
    label = parts[2] if len(parts) > 2 else (
        f"{(vb - va) / va:+.0%}" if va else _fmt_num(vb - va, signed=True))
    B.add_rect(slide, xa, y, xb - xa, 0.018, pal["text_muted"])      # beam
    B.add_rect(slide, xa, y, 0.018, 0.14, pal["text_muted"])          # left tick
    B.add_rect(slide, xb - 0.018, y, 0.018, 0.14, pal["text_muted"])  # right tick
    B.add_tb(slide, label, xa, y - 0.36, xb - xa, 0.32, size=13, bold=True,
             color=pal["text"], align=PP_ALIGN.CENTER, font=pal["font_body"])


# ── 2x2 matrix ───────────────────────────────────────────────────────────────
def build_matrix_slide(prs, p, pal, ctx):
    slide = B._blank_slide(prs, pal, pal["bg"])
    B._heading(slide, p, pal)
    L, T, W, H = 2.0, 1.85, 9.6, 4.6  # plot area
    # quadrant tints + frame
    B.add_rect(slide, L, T, W, H, pal["surface"])
    B.add_rect(slide, L, T + H / 2, W, 0.02, pal["text_muted"])
    B.add_rect(slide, L + W / 2, T, 0.02, H, pal["text_muted"])
    # axes labels
    B.add_tb(slide, p.get("x_axis", "X →"), L, T + H + 0.12, W, 0.4, size=13,
             bold=True, color=pal["text_muted"], align=PP_ALIGN.CENTER,
             font=pal["font_body"])
    y_label = p.get("y_axis", "Y ↑").replace("→", "↑")  # vertical axis points up
    B.add_tb(slide, y_label, 0.4, T + H / 2 - 0.2, 1.5, 0.8,
             size=13, bold=True, color=pal["text_muted"], font=pal["font_body"])
    # quadrant labels: Q1 TL, Q2 TR, Q3 BL, Q4 BR
    quads = [("q1", L, T), ("q2", L + W / 2, T),
             ("q3", L, T + H / 2), ("q4", L + W / 2, T + H / 2)]
    for key, qx, qy in quads:
        if p.get(key):
            B.add_tb(slide, p[key], qx + 0.15, qy + 0.1, W / 2 - 0.3, 0.4,
                     size=12, bold=True, color=pal["text_muted"],
                     font=pal["font_body"])
    items = p.get("matrix_items", [])[:12]
    sizes = []
    for it in items:
        try:
            s = float(it["size"])
            sizes.append(s if s > 0 else None)
        except (KeyError, ValueError):
            sizes.append(None)
    smax = max((s for s in sizes if s), default=0)
    from helpers import set_fill_alpha
    for item, s in zip(items, sizes):
        try:
            fx, fy = float(item.get("x", 0.5)), float(item.get("y", 0.5))
        except ValueError:
            continue
        cx = L + fx * W
        cy = T + (1 - fy) * H
        d = 0.18 if not (s and smax) else 0.24 + 0.66 * (s / smax) ** 0.5
        cx = max(L + d / 2, min(L + W - d / 2, cx))
        cy = max(T + d / 2, min(T + H - d / 2, cy))
        dot = B.add_circle(slide, cx - d / 2, cy - d / 2, d, pal["accent1"])
        if s and smax:
            set_fill_alpha(dot, 80)  # bubbles overlap; keep grid visible
        off = d / 2 + 0.06
        if cx + off + 2.2 > L + W:  # label would cross the right border
            B.add_tb(slide, item.get("name", ""), cx - off - 2.2, cy - 0.16,
                     2.2, 0.35, size=11, color=pal["text"],
                     align=PP_ALIGN.RIGHT, font=pal["font_body"])
        else:
            B.add_tb(slide, item.get("name", ""), cx + off, cy - 0.16, 2.2,
                     0.35, size=11, color=pal["text"], font=pal["font_body"])
    return slide


# ── harvey-ball scorecard ────────────────────────────────────────────────────
def _harvey(slide, pal, cx, cy, d, quarters):
    """0-4 quarters filled. Outline circle + pie wedge (or full disc)."""
    ring = B.add_circle(slide, cx, cy, d, pal["bg_deep"] if pal["dark"] else "FFFFFF",
                        line_hex=pal["accent1"], line_pt=1.2)
    if quarters >= 4:
        B.add_circle(slide, cx + 0.03, cy + 0.03, d - 0.06, pal["accent1"])
    elif quarters > 0:
        l, t, w, _ = B._sc(cx + 0.03, cy + 0.03, d - 0.06, d - 0.06)
        pie = slide.shapes.add_shape(142, Inches(l), Inches(t), Inches(w), Inches(w))
        pie.fill.solid()
        from palettes import hex_rgb
        pie.fill.fore_color.rgb = hex_rgb(pal["accent1"])
        pie.line.fill.background()
        # OOXML pie angles are 60000ths of a degree; python-pptx adjustments
        # write value*100000 — so multiply degrees by 0.6
        pie.adjustments[0] = -90.0 * 0.6
        pie.adjustments[1] = (-90.0 + 90.0 * quarters) * 0.6
    return ring


def build_harvey_slide(prs, p, pal, ctx):
    rows = p.get("table_rows", [])
    if len(rows) < 2:
        warn("harvey-scorecard needs a markdown table; rendering as table")
        return B.build_table_slide(prs, p, pal, ctx)
    slide = B._blank_slide(prs, pal, pal["bg"])
    B._heading(slide, p, pal)
    n_rows, n_cols = len(rows), max(len(r) for r in rows)
    first_col_w = 3.6
    col_w = (11.9 - first_col_w) / max(n_cols - 1, 1)
    row_h = min(0.62, 4.9 / n_rows)
    top = 1.9
    for r, row in enumerate(rows):
        y = top + r * (row_h + 0.08)
        if r > 0:
            B.add_rect(slide, 0.7, y, 11.9, row_h,
                       pal["surface"] if r % 2 else pal["bg_deep"])
        for c in range(n_cols):
            cell = row[c] if c < len(row) else ""
            x = 0.7 + (first_col_w + (c - 1) * col_w if c else 0)
            w = col_w if c else first_col_w
            if r == 0:  # header
                B.add_tb(slide, cell, x, y + 0.05, w, row_h, size=14, bold=True,
                         color=pal["accent1"],
                         align=PP_ALIGN.CENTER if c else PP_ALIGN.LEFT,
                         font=pal["font_body"])
            elif c == 0:
                B.add_tb(slide, cell, x + 0.15, y + row_h / 2 - 0.16, w - 0.2, 0.35,
                         size=13, color=pal["text"], font=pal["font_body"])
            elif cell.strip().isdigit():
                d = min(row_h - 0.18, 0.4)
                _harvey(slide, pal, x + w / 2 - d / 2, y + row_h / 2 - d / 2, d,
                        max(0, min(4, int(cell))))
            else:
                B.add_tb(slide, cell, x, y + row_h / 2 - 0.16, w, 0.35, size=12,
                         color=pal["text_muted"], align=PP_ALIGN.CENTER,
                         font=pal["font_body"])
    B.add_tb(slide, "Empty = not met · quarter fills = partial · solid = fully met",
             0.7, top + n_rows * (row_h + 0.08) + 0.1, 11.9, 0.3, size=11,
             color=pal["text_muted"], font=pal["font_label"])
    return slide


# ── process flow ─────────────────────────────────────────────────────────────
def build_process_flow_slide(prs, p, pal, ctx):
    steps = p.get("steps", [])[:6]
    if not steps:
        steps = [{"title": B.split_icon(b)[1]}
                 for b in p.get("bullets", [])[:6]]
    slide = B._blank_slide(prs, pal, pal["bg"])
    B._heading(slide, p, pal)
    n = max(len(steps), 1)
    gap = 0.12
    w = (11.9 - gap * (n - 1)) / n
    for i, step in enumerate(steps):
        x = 0.7 + i * (w + gap)
        l, t, ww, hh = B._sc(x, 2.3, w, 1.15)
        chev = slide.shapes.add_shape(52, Inches(l), Inches(t),
                                      Inches(ww), Inches(hh))  # 52 = CHEVRON
        chev.fill.solid()
        from palettes import hex_rgb
        chev.fill.fore_color.rgb = hex_rgb(pal["accent1"] if i == 0 else pal["surface"])
        chev.line.fill.background()
        B.add_tb(slide, step.get("title", ""), x + 0.45, 2.62, w - 0.55, 0.6,
                 size=14, bold=True,
                 color=pal["bg_deep"] if i == 0 else pal["text"],
                 font=pal["font_body"])
        if step.get("desc"):
            B.add_tb(slide, step["desc"], x + 0.05, 3.65, w - 0.1, 2.8, size=12,
                     color=pal["text_muted"], font=pal["font_body"],
                     accent=pal["accent1"])
    return slide


# ── big number ───────────────────────────────────────────────────────────────
def build_big_number_slide(prs, p, pal, ctx):
    slide = B._blank_slide(prs, pal, pal["bg_deep"], accent_bar=False, gradient=True)
    B.add_tb(slide, p.get("value", ""), 0.7, 2.0, 11.9, 2.2, size=110, bold=True,
             color=pal["accent1"], align=PP_ALIGN.CENTER, font=pal["font_title"])
    B.add_tb(slide, p.get("label", p.get("heading", "")), 1.5, 4.5, 10.33, 0.9,
             size=24, bold=True, color=pal["text"], align=PP_ALIGN.CENTER,
             font=pal["font_body"])
    if p.get("context"):
        B.add_tb(slide, p["context"], 1.5, 5.45, 10.33, 0.7, size=16,
                 color=pal["text_muted"], align=PP_ALIGN.CENTER,
                 font=pal["font_label"])
    return slide


# ── chart + insight callout ──────────────────────────────────────────────────
def build_chart_callout_slide(prs, p, pal, ctx):
    slide = B._blank_slide(prs, pal, pal["bg"])
    B._heading(slide, p, pal)
    B._place_visual(slide, prs, p, pal, ctx, 0.7, 1.85, 8.4, 4.9)
    box = B.add_rect(slide, 9.4, 2.5, 3.45, 2.8, pal["surface"],
                     line_hex=pal["accent1"], line_pt=1.5)
    from helpers import add_soft_shadow
    add_soft_shadow(box)
    l, t, w, h = B._sc(9.12, 3.6, 0.32, 0.26)
    arrow = slide.shapes.add_shape(34, Inches(l), Inches(t), Inches(w), Inches(h))
    arrow.fill.solid()
    from palettes import hex_rgb
    arrow.fill.fore_color.rgb = hex_rgb(pal["accent1"])
    arrow.line.fill.background()
    B.add_tb(slide, "KEY INSIGHT", 9.65, 2.7, 3.0, 0.35, size=14, bold=True,
             color=pal["accent1"], font=pal["font_label"])
    B.add_tb(slide, p.get("callout", ""), 9.65, 3.1, 3.0, 2.05, size=14,
             color=pal["text"], font=pal["font_body"], accent=pal["accent1"])
    return slide


# ── KPI dashboard ────────────────────────────────────────────────────────────
def build_dashboard_slide(prs, p, pal, ctx):
    tiles = p.get("tiles", [])[:6]
    slide = B._blank_slide(prs, pal, pal["bg"])
    B._heading(slide, p, pal)
    cols = 3 if len(tiles) > 4 else 2
    rows = -(-len(tiles) // cols)
    gap = 0.3
    tw = (11.9 - gap * (cols - 1)) / cols
    th = min(2.15, (4.7 - gap * (rows - 1)) / max(rows, 1))
    from helpers import add_soft_shadow
    for i, tile in enumerate(tiles):
        r, c = divmod(i, cols)
        x = 0.7 + c * (tw + gap)
        y = 1.95 + r * (th + gap)
        card = B.add_rect(slide, x, y, tw, th, pal["surface"])
        add_soft_shadow(card)
        B.add_tb(slide, tile.get("value", ""), x + 0.25, y + 0.18, tw - 0.5, 0.85,
                 size=32, bold=True, color=pal["accent1"], font=pal["font_title"])
        B.add_tb(slide, tile.get("label", ""), x + 0.25, y + 0.98, tw - 0.5,
                 0.45, size=13, color=pal["text"], font=pal["font_body"])
        if tile.get("delta"):
            up = tile.get("trend", "up").lower() != "down"
            # 14pt bold = WCAG large text (3:1), so these hold on white tiles
            B.add_tb(slide, f"{'▲' if up else '▼'} {tile['delta']}",
                     x + 0.25, y + 1.42, tw - 0.5, 0.4, size=14, bold=True,
                     color=pal["accent1"] if up else NEG_HEX,
                     font=pal["font_body"])
    return slide


# ── quote + evidence ─────────────────────────────────────────────────────────
def build_quote_evidence_slide(prs, p, pal, ctx):
    slide = B._blank_slide(prs, pal, pal["bg"])
    B.add_tb(slide, "“", 0.6, 1.0, 1.4, 1.6, size=110, bold=True,
             color=pal["accent1"], font=pal["font_title"])
    quote_tb = B.add_tb(slide, p.get("quote", ""), 1.1, 2.3, 6.9, 3.2, size=22,
                        color=pal["text"], font=pal["font_title"],
                        accent=pal["accent1"])
    from pptx.enum.text import MSO_ANCHOR
    quote_tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    if p.get("attribution"):
        B.add_tb(slide, f"— {p['attribution']}", 1.1, 5.6, 6.9, 0.6, size=14,
                 bold=True, color=pal["text_muted"], font=pal["font_body"])
    stats = p.get("stats", [])[:1]
    if stats:
        from helpers import add_soft_shadow
        card = B.add_rect(slide, 8.7, 2.2, 4.0, 3.2, pal["surface"],
                          line_hex=pal["accent1"], line_pt=1.5)
        add_soft_shadow(card)
        B.add_tb(slide, stats[0].get("value", ""), 8.9, 2.6, 3.6, 1.3, size=48,
                 bold=True, color=pal["accent1"], align=PP_ALIGN.CENTER,
                 font=pal["font_title"])
        B.add_tb(slide, stats[0].get("label", ""), 8.9, 4.0, 3.6, 1.1, size=14,
                 color=pal["text"], align=PP_ALIGN.CENTER, font=pal["font_body"])
    return slide


# ── funnel ───────────────────────────────────────────────────────────────────
def build_funnel_slide(prs, p, pal, ctx):
    data = [(l, v) for l, v in p.get("data", []) if isinstance(v, (int, float))]
    if len(data) < 2:
        warn("funnel needs a **Data:** block with 2+ numeric rows")
        return B.build_bullet_slide(prs, p, pal, ctx)
    slide = B._blank_slide(prs, pal, pal["bg"])
    B._heading(slide, p, pal)
    n = len(data)
    top_v = data[0][1] or 1
    row_h = min(0.85, 4.7 / n)
    label_w = 3.0  # stage labels live left of the bars (never overflow them)
    bar_area = 11.9 - label_w - 1.2
    for i, (label, v) in enumerate(data):
        w = max(1.0, bar_area * (v / top_v))
        x = 0.7 + label_w + 0.2 + (bar_area - w) / 2
        y = 1.95 + i * (row_h + 0.18)
        B.add_tb(slide, label, 0.7, y + row_h / 2 - 0.18, label_w, 0.4,
                 size=13, bold=i == 0, color=pal["text"],
                 align=PP_ALIGN.RIGHT, font=pal["font_body"])
        B.add_rect(slide, x, y, w, row_h,
                   pal["accent1"] if i == 0 else pal["surface"],
                   line_hex=pal["accent1"], line_pt=0 if i == 0 else 1.0)
        B.add_tb(slide, _fmt_plain(v), x, y + row_h / 2 - 0.18, w, 0.4, size=14,
                 bold=True, color=pal["bg_deep"] if i == 0 else pal["accent1"],
                 align=PP_ALIGN.CENTER, font=pal["font_body"])
        if i:
            pct = v / data[i - 1][1] * 100 if data[i - 1][1] else 0
            B.add_tb(slide, f"↳ {pct:.0f}%", x + w + 0.15, y + row_h / 2 - 0.16,
                     1.1, 0.32, size=11, color=pal["text_muted"],
                     font=pal["font_label"])
    return slide


def _fmt_plain(v):
    return f"{v:,.0f}" if v == int(v) else f"{v:,.1f}"


# ── next steps (conclusions that stay up during Q&A) ─────────────────────────
def build_next_steps_slide(prs, p, pal, ctx):
    steps = p.get("steps", [])[:6]
    if not steps:  # plain bullets validate too — render them as step cards
        steps = [{"action": B.split_icon(b)[1]}
                 for b in p.get("bullets", [])[:6]]
    slide = B._blank_slide(prs, pal, pal["bg"])
    B._heading(slide, p, pal)
    from helpers import add_soft_shadow
    row_h = min(0.85, 4.8 / max(len(steps), 1))
    for i, s in enumerate(steps):
        y = 1.95 + i * (row_h + 0.15)
        card = B.add_rect(slide, 0.7, y, 11.9, row_h, pal["surface"])
        add_soft_shadow(card)
        B.add_circle(slide, 0.95, y + row_h / 2 - 0.19, 0.38, pal["accent1"])
        B.add_tb(slide, str(i + 1), 0.95, y + row_h / 2 - 0.21, 0.38, 0.4,
                 size=15, bold=True, color=pal["bg_deep"],
                 align=PP_ALIGN.CENTER, font=pal["font_body"])
        B.add_tb(slide, s.get("action", s.get("title", "")), 1.55,
                 y + row_h / 2 - 0.19, 6.6, 0.42, size=14, bold=True,
                 color=pal["text"], font=pal["font_body"])
        B.add_tb(slide, s.get("owner", ""), 8.3, y + row_h / 2 - 0.19, 2.1, 0.42,
                 size=14, bold=True, color=pal["accent1"], font=pal["font_body"])
        B.add_tb(slide, s.get("when", ""), 10.5, y + row_h / 2 - 0.19, 2.0, 0.42,
                 size=13, color=pal["text_muted"], align=PP_ALIGN.RIGHT,
                 font=pal["font_body"])
    return slide


# ── agenda with current-section highlight ────────────────────────────────────
def build_agenda_slide(prs, p, pal, ctx):
    sections = p.get("bullets", [])[:8]
    current = p.get("current", "").strip().lower()
    slide = B._blank_slide(prs, pal, pal["bg_deep"], gradient=True)
    B.add_tb(slide, p.get("heading", "Agenda"), 0.8, 0.7, 11.7, 1.0, size=34,
             bold=True, color=pal["accent1"], font=pal["font_title"])
    row_h = min(0.75, 5.0 / max(len(sections), 1))
    for i, name in enumerate(sections):
        y = 2.0 + i * (row_h + 0.12)
        is_cur = current and current in name.lower()
        if is_cur:
            B.add_rect(slide, 0.8, y - 0.06, 11.0, row_h, pal["surface"])
            B.add_rect(slide, 0.8, y - 0.06, 0.08, row_h, pal["accent1"])
        B.add_tb(slide, f"{i + 1:02d}", 1.1, y, 0.8, 0.5, size=16, bold=True,
                 color=pal["accent1"] if is_cur else pal["text_muted"],
                 font=pal["font_title"])
        B.add_tb(slide, name, 2.0, y, 9.5, 0.5, size=17, bold=is_cur,
                 color=pal["text"] if is_cur else pal["text_muted"],
                 font=pal["font_body"])
    return slide


# ── SCQA executive summary (Situation | Findings | Recommendation) ───────────
def build_exec_scqa_slide(prs, p, pal, ctx):
    slide = B._blank_slide(prs, pal, pal["bg"])
    B._heading(slide, p, pal)
    from helpers import add_soft_shadow
    cols = [
        ("SITUATION", [p.get("situation", "")], False),
        ("FINDINGS", p.get("bullets", [])[:4], False),
        ("RECOMMENDATION", [p.get("recommendation", "")], True),
    ]
    gap = 0.3
    col_w = (11.9 - 2 * gap) / 3
    for i, (header, items, emphasize) in enumerate(cols):
        x = 0.7 + i * (col_w + gap)
        card = B.add_rect(slide, x, 1.9, col_w, 4.7, pal["surface"],
                          line_hex=pal["accent1"] if emphasize else None,
                          line_pt=1.5 if emphasize else 0)
        add_soft_shadow(card)
        B.add_tb(slide, header, x + 0.25, 2.1, col_w - 0.5, 0.4, size=14,
                 bold=True, color=pal["accent1"], font=pal["font_label"])
        single = len([t for t in items if t]) == 1
        for j, item in enumerate([t for t in items if t]):
            tb = B.add_tb(slide, item, x + 0.25, 2.6 + j * 1.0, col_w - 0.5,
                          0.96 if not single else 3.8,
                          size=13, bold=emphasize, color=pal["text"],
                          font=pal["font_body"], accent=pal["accent1"])
            if single:  # center short text in the tall column
                from pptx.enum.text import MSO_ANCHOR
                tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return slide


LAYOUTS = {
    "waterfall": build_waterfall_slide,
    "matrix-2x2": build_matrix_slide,
    "harvey-scorecard": build_harvey_slide,
    "process-flow": build_process_flow_slide,
    "big-number": build_big_number_slide,
    "chart-callout": build_chart_callout_slide,
    "dashboard": build_dashboard_slide,
    "quote-evidence": build_quote_evidence_slide,
    "funnel": build_funnel_slide,
    "next-steps": build_next_steps_slide,
    "agenda": build_agenda_slide,
    "exec-summary-scqa": build_exec_scqa_slide,
}


# ── mekko (marimekko) ────────────────────────────────────────────────────────
def _rel_lum(hex_color):
    def chan(c):
        c = c / 255
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    r, g, b = (int(hex_color[i:i + 2], 16) for i in (0, 2, 4))
    return 0.2126 * chan(r) + 0.7152 * chan(g) + 0.0722 * chan(b)


def _label_color_on(fill_hex, pal):
    """Pick whichever palette color (text vs bg_deep) contrasts most with the
    fill — works on dark AND light themes (text is dark on light palettes)."""
    def ratio(a, b):
        la, lb = sorted((_rel_lum(a), _rel_lum(b)), reverse=True)
        return (la + 0.05) / (lb + 0.05)
    return max((pal["text"], pal["bg_deep"]),
               key=lambda c: ratio(c, fill_hex))


def build_mekko_slide(prs, p, pal, ctx):
    """Variable-width 100%-stacked columns: column width ∝ column total,
    segment height = share of column. Needs multi-series **Data:** rows."""
    names = [n.strip() for n in p.get("series", "").split(",") if n.strip()]
    data = [(lbl, v) for lbl, v in p.get("data", []) if isinstance(v, list)]
    if not names or not data:
        warn("mekko needs **Series:** names and multi-series **Data:** rows")
        return B.build_bullet_slide(prs, p, pal, ctx)
    slide = B._blank_slide(prs, pal, pal["bg"])
    B._heading(slide, p, pal)

    accents = [pal["accent1"], pal["accent2"], pal["accent3"], NEG_HEX]
    # legend
    lx = 0.75
    for i, name in enumerate(names):
        B.add_rect(slide, lx, 1.72, 0.18, 0.18, accents[i % len(accents)])
        B.add_tb(slide, name, lx + 0.26, 1.64, 1.9, 0.32, size=12,
                 color=pal["text_muted"], font=pal["font_label"])
        lx += 0.32 + min(len(name) * 0.085 + 0.3, 2.0)

    top, bottom = 2.25, 6.05
    height = bottom - top
    gap = 0.06
    totals = [sum(v) for _, v in data]
    grand = sum(totals) or 1
    plot_w = 11.9 - gap * (len(data) - 1)
    x = 0.7
    for (label, vals), total in zip(data, totals):
        col_w = plot_w * (total / grand)
        # largest-remainder so each column's segment labels sum to 100
        pcts = (round_to_sum([v / total * 100 for v in vals])
                if total else [0] * len(vals))
        y = bottom
        for i, v in enumerate(vals):
            seg_h = height * (v / total) if total else 0
            y -= seg_h
            fill = accents[i % len(accents)]
            B.add_rect(slide, x, y, col_w, seg_h, fill,
                       line_hex=pal["bg"], line_pt=0.75)
            if seg_h >= 0.34 and col_w >= 0.9:
                B.add_tb(slide, f"{pcts[i]:.0f}%", x, y + seg_h / 2 - 0.16,
                         col_w, 0.32, size=14, bold=True,
                         color=_label_color_on(fill, pal),
                         align=PP_ALIGN.CENTER, font=pal["font_body"])
        B.add_tb(slide, _fmt_plain(total), x - 0.1, top - 0.36, col_w + 0.2, 0.3,
                 size=12, bold=True, color=pal["text"],
                 align=PP_ALIGN.CENTER, font=pal["font_body"])
        B.add_tb(slide, label, x - 0.12, bottom + 0.08, col_w + 0.24, 0.66,
                 size=11, color=pal["text_muted"], align=PP_ALIGN.CENTER,
                 font=pal["font_label"])
        x += col_w + gap
    return slide


# ── gantt / swimlane roadmap ─────────────────────────────────────────────────
def build_gantt_slide(prs, p, pal, ctx):
    """Workstream rows x period columns with phase bars and milestone diamonds.
    Syntax: **Periods:** Q1, Q2... + '- Bar: Row=".." Label=".." Start="1"
    End="2"' + '- Milestone: Row=".." Label=".." At="2.5"'."""
    periods = p.get("periods", [])
    bars = p.get("bars", [])
    if not periods or not bars:
        warn("gantt needs **Periods:** and '- Bar:' rows")
        return B.build_bullet_slide(prs, p, pal, ctx)
    slide = B._blank_slide(prs, pal, pal["bg"])
    B._heading(slide, p, pal)

    rows = list(dict.fromkeys([b.get("row", "") for b in bars]))
    label_w, grid_x = 2.3, 3.1
    grid_w = 12.6 - grid_x
    col_w = grid_w / len(periods)
    top = 2.35
    row_h = min(0.85, 4.3 / max(len(rows), 1))
    grid_h = row_h * len(rows)

    for i, period in enumerate(periods):
        B.add_tb(slide, period, grid_x + i * col_w, top - 0.42, col_w, 0.34,
                 size=12, bold=True, color=pal["text_muted"],
                 align=PP_ALIGN.CENTER, font=pal["font_label"])
        if i:
            B.add_rect(slide, grid_x + i * col_w, top, 0.015, grid_h,
                       pal["surface"])
    accents = [pal["accent1"], pal["accent2"], pal["accent3"]]
    for r, row_name in enumerate(rows):
        y = top + r * row_h
        if r:
            B.add_rect(slide, 0.7, y, 11.9, 0.015, pal["surface"])
        B.add_tb(slide, row_name, 0.8, y + row_h / 2 - 0.17, label_w, 0.38,
                 size=13, bold=True, color=pal["text"], font=pal["font_body"])
    for bar in bars:
        try:
            start = float(bar.get("start", 1))
            end = float(bar.get("end", start))
        except ValueError:
            continue
        r = rows.index(bar.get("row", ""))
        fill = accents[r % len(accents)]
        bx = grid_x + (start - 1) * col_w + 0.04
        bw = max((end - start + 1) * col_w - 0.08, 0.25)
        y = top + r * row_h + row_h / 2 - 0.21
        B.add_rect(slide, bx, y, bw, 0.42, fill)
        B.add_tb(slide, bar.get("label", ""), bx + 0.08, y + 0.04, bw - 0.16,
                 0.34, size=14, bold=True, color=_label_color_on(fill, pal),
                 font=pal["font_body"])
    for ms in p.get("milestones", []):
        try:
            at = float(ms.get("at", 1))
        except ValueError:
            continue
        if ms.get("row", "") not in rows:
            continue
        r = rows.index(ms.get("row", ""))
        cx = grid_x + (at - 0.5) * col_w  # At=N means the middle of period N
        cy = top + r * row_h + row_h / 2
        # if the milestone falls within a bar's span on its row, lift it above
        # the bar so it doesn't cover the bar label
        overlaps = any(
            b.get("row") == ms.get("row")
            and float(b.get("start", 1)) - 0.5 <= at <= float(b.get("end", 1)) + 0.5
            for b in bars if str(b.get("start", "")).replace(".", "").isdigit())
        if overlaps:
            cy = top + r * row_h + 0.13
        l, t, w, _ = B._sc(cx - 0.13, cy - 0.13, 0.26, 0.26)
        d = slide.shapes.add_shape(4, Inches(l), Inches(t), Inches(w), Inches(w))
        d.fill.solid()
        from palettes import hex_rgb
        d.fill.fore_color.rgb = hex_rgb(pal["text"])
        d.line.fill.background()
        if ms.get("label"):
            # label sits just above the row, clear of bars and the diamond
            B.add_tb(slide, ms["label"], cx - 0.75,
                     top + r * row_h - 0.30, 1.5, 0.3, size=11, bold=True,
                     color=pal["text"], align=PP_ALIGN.CENTER,
                     font=pal["font_label"])
    return slide


LAYOUTS["mekko"] = build_mekko_slide
LAYOUTS["gantt"] = build_gantt_slide


# ── bar mekko (profit pool: width = size, height = value) ───────────────────
def build_bar_mekko_slide(prs, p, pal, ctx):
    slide = B._blank_slide(prs, pal, pal["bg"])
    B._heading(slide, p, pal)
    rows = []
    for b in p.get("bars", []):
        try:
            size, value = float(b["size"]), float(b["value"])
        except (KeyError, ValueError):
            warn(f"bar-mekko row needs numeric Size and Value: {b}")
            continue
        if size <= 0 or value < 0:
            warn(f"bar-mekko row needs Size>0 and Value>=0: {b}")
            continue
        rows.append((b.get("label", "?"), size, value))
    if len(rows) < 2:
        warn("bar-mekko has <2 valid bars; rendering bullet-list")
        return B.build_bullet_slide(prs, p, pal, ctx)

    L, bottom, total_w, top = 0.7, 6.15, 11.9, 2.25
    gap = 0.08
    total_size = sum(s for _, s, _ in rows) or 1
    vmax = max(v for _, _, v in rows) or 1
    scale = (bottom - top) / vmax
    accents = [pal["accent1"], pal["accent2"], pal["accent3"]]

    x = L
    for i, (label, size, value) in enumerate(rows):
        w = (total_w - gap * (len(rows) - 1)) * size / total_size
        h = max(value * scale, 0.04)
        B.add_rect(slide, x, bottom - h, w, h, accents[i % len(accents)])
        B.add_tb(slide, _fmt_num(value), x - 0.1, bottom - h - 0.34, w + 0.2,
                 0.3, size=12, bold=True, color=pal["text"],
                 align=PP_ALIGN.CENTER, font=pal["font_body"])
        B.add_tb(slide, f"{label}\n{_fmt_num(size)}", x - 0.1, bottom + 0.10,
                 w + 0.2, 0.75, size=11, color=pal["text_muted"],
                 align=PP_ALIGN.CENTER, font=pal["font_label"])
        x += w + gap
    B.add_rect(slide, L, bottom, total_w, 0.02, pal["surface"])
    return slide


LAYOUTS["bar-mekko"] = build_bar_mekko_slide


# ── heatmap table ────────────────────────────────────────────────────────────
_CELL_NUM_RX = re.compile(r"-?(?:\d{1,3}(?:,\d{3})*|\d+)(?:\.\d+)?")


def _cell_value(cell):
    """Numeric value of a table cell ('$1,200', '42%', '-3.5') or None."""
    stripped = str(cell).strip().lstrip("$€£").rstrip("%").strip()
    m = _CELL_NUM_RX.fullmatch(stripped)
    return float(m.group().replace(",", "")) if m else None


def _lerp_hex(a, b, t):
    """Linear blend between two hex colors, t in 0..1."""
    return "".join(
        f"{round(int(a[i:i + 2], 16) + (int(b[i:i + 2], 16) - int(a[i:i + 2], 16)) * t):02X}"
        for i in (0, 2, 4))


def _bw_on(fill_hex):
    """Pure black/white text by fill luminance. The 0.179 threshold is the
    equal-contrast point, so the winner always clears 4.5:1 on any fill
    (palette text colors can dip to ~4.3:1 on mid-tone heat cells)."""
    return "000000" if _rel_lum(fill_hex) > 0.179 else "FFFFFF"


def _rag_fill(v, svals, pal):
    """Tercile thresholds over the column's sorted values -> red/amber/green.

    Note: with exactly 2 rows, t1 == svals[0] and t2 == svals[-1], so the
    function yields rag_bad for the lower value and rag_good for the upper
    (rag_mid is never returned in the 2-row case).
    """
    if svals[0] == svals[-1]:
        return pal["rag_mid"]
    n = len(svals)
    t1, t2 = svals[(n - 1) // 3], svals[(2 * (n - 1)) // 3]
    if v <= t1:
        return pal["rag_bad"]
    return pal["rag_mid"] if v <= t2 else pal["rag_good"]


def build_heatmap_slide(prs, p, pal, ctx):
    """Markdown table whose numeric body cells get heat fills, normalized
    per column: bg at the column min -> accent1 at the max. '- Scale: rag'
    switches to terciled red/amber/green chips. Non-numeric body cells get
    a plain surface fill; cell text picks its color by fill luminance."""
    rows = p.get("table_rows", [])
    if len(rows) < 2:
        warn("heatmap-table needs a markdown table; rendering as table")
        return B.build_table_slide(prs, p, pal, ctx)
    slide = B._blank_slide(prs, pal, pal["bg"])
    B._heading(slide, p, pal)
    rag = p.get("scale", "").strip().lower() == "rag"

    n_cols = max(len(r) for r in rows)
    col_vals = []  # per column: sorted numeric body values
    for c in range(n_cols):
        nums = [_cell_value(r[c]) for r in rows[1:] if c < len(r)]
        col_vals.append(sorted(v for v in nums if v is not None))

    n_body_rows = len(rows) - 1  # exclude header
    if n_body_rows > 12:
        warn(f"heatmap-table with {n_body_rows} rows — consider splitting")

    first_col_w = 3.2
    col_w = (11.9 - first_col_w) / max(n_cols - 1, 1)
    # rows + 0.06 gaps must clear the footer zone (1.9 + total <= 6.6)
    row_h = min(0.72, 4.7 / len(rows) - 0.06)
    top = 1.9
    for r, row in enumerate(rows):
        y = top + r * (row_h + 0.06)
        for c in range(n_cols):
            cell = row[c] if c < len(row) else ""
            x = 0.7 + (first_col_w + (c - 1) * col_w if c else 0)
            w = col_w if c else first_col_w
            if r == 0:  # header
                B.add_tb(slide, cell, x + (0.12 if c == 0 else 0),
                         y + row_h / 2 - 0.16, w, 0.35, size=13, bold=True,
                         color=pal["accent1"],
                         align=PP_ALIGN.CENTER if c else PP_ALIGN.LEFT,
                         font=pal["font_body"])
                continue
            v = _cell_value(cell)
            if v is None or not col_vals[c]:
                fill, text_color = pal["surface"], pal["text"]
            elif rag:
                fill = _rag_fill(v, col_vals[c], pal)
                text_color = _bw_on(fill)
            else:
                lo, hi = col_vals[c][0], col_vals[c][-1]
                t = (v - lo) / (hi - lo) if hi > lo else 0.5
                fill = _lerp_hex(pal["bg"], pal["accent1"], t)
                text_color = _bw_on(fill)
            B.add_rect(slide, x + 0.03, y, w - 0.06, row_h, fill)
            # when rows are very dense (row_h < 0.35) shrink the text box
            # height to match and drop font size to 10 so text doesn't overflow
            cell_tb_h = min(0.35, row_h)
            cell_font = 10 if row_h < 0.35 else 12
            B.add_tb(slide, cell, x + (0.12 if c == 0 else 0),
                     y + row_h / 2 - cell_tb_h / 2, w, cell_tb_h,
                     size=cell_font, color=text_color,
                     align=PP_ALIGN.CENTER if c else PP_ALIGN.LEFT,
                     font=pal["font_body"])
    return slide


LAYOUTS["heatmap-table"] = build_heatmap_slide


# ── tornado (sensitivity) ────────────────────────────────────────────────────
def build_tornado_slide(prs, p, pal, ctx):
    """Sensitivity tornado built from shapes. Requires **Series:** Low, High
    (2 names) so '- Driver: -12, +18' rows parse as 2-lists; bars hang off a
    central label gutter — left = values[0] (accent2), right = values[1]
    (accent1), shared symmetric scale, value labels at the outer ends. Rows
    sort by |left|+|right| descending unless '- Sort: off'."""
    rows = [(lbl, v) for lbl, v in p.get("data", [])
            if isinstance(v, list) and len(v) == 2]
    if len(rows) < 2:
        warn("tornado needs **Series:** Low, High and 2+ "
             "'- Driver: low, high' **Data:** rows")
        return B.build_bullet_slide(prs, p, pal, ctx)
    if p.get("sort", "").strip().lower() not in ("off", "no", "false"):
        rows.sort(key=lambda r: abs(r[1][0]) + abs(r[1][1]), reverse=True)
    slide = B._blank_slide(prs, pal, pal["bg"])
    B._heading(slide, p, pal)

    top, bottom = 2.0, 6.3
    cx, gutter = 6.65, 1.25   # spine center + half-width of the label gutter
    val_w = 0.72              # outer value-label zone per side
    half_w = cx - gutter - 0.7 - val_w
    vmax = max(max(abs(v[0]), abs(v[1])) for _, v in rows) or 1
    row_h = min(0.78, (bottom - top) / len(rows))
    bar_h = max(0.04, min(0.42, row_h - 0.16))

    names = [s.strip() for s in p.get("series", "").split(",") if s.strip()]
    if len(names) >= 2:  # side headers above the two bar areas
        B.add_tb(slide, names[0], 0.7, top - 0.42, cx - gutter - 0.75, 0.3,
                 size=11, bold=True, color=pal["text_muted"],
                 align=PP_ALIGN.RIGHT, font=pal["font_label"])
        B.add_tb(slide, names[1], cx + gutter + 0.05, top - 0.42,
                 12.6 - cx - gutter, 0.3, size=11, bold=True,
                 color=pal["text_muted"], font=pal["font_label"])
    chart_h = row_h * len(rows)
    for gx in (cx - gutter, cx + gutter):  # spine edges the bars hang off
        B.add_rect(slide, gx - 0.009, top, 0.018, chart_h, pal["text_muted"])
    for i, (label, (lo, hi)) in enumerate(rows):
        y = top + i * row_h + (row_h - bar_h) / 2
        wl = max(half_w * abs(lo) / vmax, 0.02)
        wr = max(half_w * abs(hi) / vmax, 0.02)
        B.add_rect(slide, cx - gutter - wl, y, wl, bar_h, pal["accent2"])
        B.add_rect(slide, cx + gutter, y, wr, bar_h, pal["accent1"])
        B.add_tb(slide, label, cx - gutter + 0.05, y + bar_h / 2 - 0.16,
                 2 * gutter - 0.1, 0.35, size=12, bold=True, color=pal["text"],
                 align=PP_ALIGN.CENTER, font=pal["font_body"])
        B.add_tb(slide, _fmt_num(lo), cx - gutter - wl - val_w - 0.05,
                 y + bar_h / 2 - 0.16, val_w, 0.32, size=12,
                 color=pal["text_muted"], align=PP_ALIGN.RIGHT,
                 font=pal["font_label"])
        B.add_tb(slide, _fmt_num(hi, signed=hi > 0), cx + gutter + wr + 0.05,
                 y + bar_h / 2 - 0.16, val_w, 0.32, size=12,
                 color=pal["text_muted"], font=pal["font_label"])
    return slide


LAYOUTS["tornado"] = build_tornado_slide


# ── football field (valuation ranges) ────────────────────────────────────────
def _parse_marker(spec):
    """'Current price, 47' -> ('Current price', 47.0) or None."""
    label, sep, val = (spec or "").rpartition(",")
    m = _CELL_NUM_RX.search(val)
    if not sep or not m or not label.strip():
        return None
    return label.strip().strip('"'), float(m.group().replace(",", ""))


def build_football_field_slide(prs, p, pal, ctx):
    """Valuation football field. Requires **Series:** Low, High so
    '- Method: low, high' rows parse as 2-lists; each row is a floating
    rounded bar low->high on a shared value axis with nice-interval
    gridlines. Optional '- Marker: label, value' draws a dashed vertical
    reference line (skipped with a warning outside the data range)."""
    rows = [(lbl, v[0], v[1]) for lbl, v in p.get("data", [])
            if isinstance(v, list) and len(v) == 2 and v[0] < v[1]]
    if len(rows) < 2:
        warn("football-field needs **Series:** Low, High and 2+ "
             "'- Method: low, high' **Data:** rows (low < high)")
        return B.build_bullet_slide(prs, p, pal, ctx)
    slide = B._blank_slide(prs, pal, pal["bg"])
    B._heading(slide, p, pal)

    import math
    from charts import _nice_ceil
    gmin = min(lo for _, lo, _ in rows)
    gmax = max(hi for _, _, hi in rows)
    step = _nice_ceil((gmax - gmin) / 5 or 1)
    a0 = math.floor(gmin / step) * step
    a1 = math.ceil(gmax / step) * step
    if a1 <= a0:
        a1 = a0 + step

    # axis labels must stay above the footer zone (top <= 6.0in) or
    # pptx_lint's page-number heuristics flag the digit-only tick labels
    L, R, top, bottom = 3.5, 12.5, 2.0, 5.85

    def vx(v):
        return L + (v - a0) / (a1 - a0) * (R - L)

    for k in range(int(round((a1 - a0) / step)) + 1):  # gridlines + axis
        v = a0 + k * step
        B.add_rect(slide, vx(v) - 0.0075, top, 0.015, bottom - top,
                   pal["surface"])
        B.add_tb(slide, _fmt_plain(v), vx(v) - 0.6, bottom + 0.08, 1.2, 0.3,
                 size=11, color=pal["text_muted"], align=PP_ALIGN.CENTER,
                 font=pal["font_label"])

    from palettes import hex_rgb
    row_h = min(0.9, (bottom - top) / len(rows))
    bar_h = max(0.04, min(0.46, row_h - 0.22))
    for i, (label, lo, hi) in enumerate(rows):
        y = top + i * row_h + (row_h - bar_h) / 2
        x0, x1 = vx(lo), vx(hi)
        sl, st, sw, sh = B._sc(x0, y, x1 - x0, bar_h)
        bar = slide.shapes.add_shape(5, Inches(sl), Inches(st),
                                     Inches(sw), Inches(sh))  # 5 = ROUNDED_RECT
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_rgb(pal["accent1"])
        bar.line.fill.background()
        B.add_tb(slide, label, 0.7, y + bar_h / 2 - 0.17, 2.3, 0.38, size=13,
                 bold=True, color=pal["text"], align=PP_ALIGN.RIGHT,
                 font=pal["font_body"])
        B.add_tb(slide, _fmt_plain(lo), x0 - 0.78, y + bar_h / 2 - 0.16, 0.7,
                 0.32, size=12, color=pal["text_muted"], align=PP_ALIGN.RIGHT,
                 font=pal["font_label"])
        B.add_tb(slide, _fmt_plain(hi), x1 + 0.05, y + bar_h / 2 - 0.16, 0.7,
                 0.32, size=12, color=pal["text_muted"],
                 font=pal["font_label"])
    _football_marker(slide, p, pal, top, bottom, vx, gmin, gmax)
    return slide


def _football_marker(slide, p, pal, top, bottom, vx, gmin, gmax):
    """Dashed vertical reference line + label for '- Marker: label, value'."""
    spec = p.get("marker")
    if not spec:
        return
    parsed = _parse_marker(spec)
    if not parsed:
        warn(f"Marker needs 'label, value': {spec!r} — skipped")
        return
    label, value = parsed
    if not gmin <= value <= gmax:
        warn(f"Marker {label!r} value {value:g} outside range "
             f"[{gmin:g}, {gmax:g}] — skipped")
        return
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    from palettes import hex_rgb
    x = vx(value)
    l, t, _, h = B._sc(x, top - 0.12, 0, bottom - top + 0.12)
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(l), Inches(t),
                                      Inches(l), Inches(t + h))
    conn.line.color.rgb = hex_rgb(pal["accent3"])
    conn.line.width = Pt(1.75)
    conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    # clamp 2.4in-wide label so it stays within canvas (right edge <= 12.63)
    label_left = max(0.7, min(x - 1.2, 12.63 - 2.4))
    B.add_tb(slide, f"{label} · {_fmt_plain(value)}", label_left, top - 0.46,
             2.4, 0.3, size=11, bold=True, color=pal["text_muted"],
             align=PP_ALIGN.CENTER, font=pal["font_label"])


LAYOUTS["football-field"] = build_football_field_slide


# ── driver tree (value decomposition) ────────────────────────────────────────
def _tree_leaf_order(nid, children, seen):
    """DFS leaf ids under nid in outline order (cycle-safe)."""
    if nid in seen:
        return []
    seen.add(nid)
    kids = children.get(nid, [])
    if not kids:
        return [nid]
    out = []
    for k in kids:
        out.extend(_tree_leaf_order(k, children, seen))
    return out


def build_driver_tree_slide(prs, p, pal, ctx):
    """Left-to-right value decomposition: root box at left, children in
    depth columns to the right, elbow connectors parent -> child. Boxes show
    Label + Value + optional Delta (green-ish '+', red-ish '-')."""
    nodes = [n for n in p.get("nodes", []) if n.get("id", "").strip()][:12]
    if not nodes:
        warn("driver-tree needs '- Node:' rows; rendering bullet-list")
        return B.build_bullet_slide(prs, p, pal, ctx)
    slide = B._blank_slide(prs, pal, pal["bg"])
    B._heading(slide, p, pal)

    by_id = {n["id"]: n for n in nodes}
    children = {}
    for n in nodes:
        par = n.get("parent", "").strip()
        if par and par in by_id:
            children.setdefault(par, []).append(n["id"])
    roots = [n["id"] for n in nodes if not n.get("parent", "").strip()]
    root = roots[0] if roots else nodes[0]["id"]

    depths, queue = {root: 1}, [root]  # BFS from root (validator caps at 3)
    while queue:
        nid = queue.pop(0)
        for ch in children.get(nid, []):
            if ch not in depths:
                depths[ch] = depths[nid] + 1
                queue.append(ch)
    skipped = [i for i in by_id if i not in depths]
    if skipped:
        warn("driver-tree: nodes unreachable from root skipped: "
             + ", ".join(skipped))
    max_d = max(depths.values())

    top, bottom = 1.95, 6.45
    gap = 0.55  # connector corridor between columns
    col_w = min(3.4, (11.93 - gap * (max_d - 1)) / max_d)
    leaves = _tree_leaf_order(root, children, set())
    row_h = (bottom - top) / max(len(leaves), 1)
    box_h = max(0.38, min(0.95, row_h - 0.10))
    ys = {nid: top + r * row_h + row_h / 2 for r, nid in enumerate(leaves)}

    def center_y(nid):  # parents center on their children
        if nid not in ys:
            ys[nid] = sum(center_y(k) for k in children[nid]) / len(children[nid])
        return ys[nid]
    center_y(root)

    for nid, d in depths.items():
        _tree_box(slide, pal, by_id[nid],
                  0.7 + (d - 1) * (col_w + gap), ys[nid] - box_h / 2,
                  col_w, box_h)
    _tree_connectors(slide, pal, children, depths, ys, col_w, gap)
    return slide


def _tree_box(slide, pal, node, x, y, col_w, box_h):
    """Surface card with accent edge stripe: Label bold + Value + Delta."""
    B.add_rect(slide, x, y, col_w, box_h, pal["surface"])
    B.add_rect(slide, x, y, 0.06, box_h, pal["accent1"])
    delta = node.get("delta", "").strip()
    d_color = (pal["rag_good"] if delta.startswith("+")
               else pal["rag_bad"] if delta.startswith("-")
               else pal["text_muted"])
    if box_h >= 0.62:  # two lines: label / value + delta
        B.add_tb(slide, node.get("label", node["id"]), x + 0.18, y + 0.07,
                 col_w - 0.30, 0.32, size=13, bold=True, color=pal["text"],
                 font=pal["font_body"])
        B.add_tb(slide, node.get("value", ""), x + 0.18, y + box_h - 0.36,
                 col_w - 1.2, 0.3, size=12, color=pal["text_muted"],
                 font=pal["font_body"])
        if delta:
            B.add_tb(slide, delta, x + col_w - 1.0, y + box_h - 0.36, 0.88,
                     0.3, size=12, bold=True, color=d_color,
                     align=PP_ALIGN.RIGHT, font=pal["font_body"])
    else:  # compact: single centered row, value + delta right-aligned
        ty = y + box_h / 2 - 0.15
        B.add_tb(slide, node.get("label", node["id"]), x + 0.16, ty,
                 col_w - 1.9, 0.3, size=11, bold=True, color=pal["text"],
                 font=pal["font_body"])
        # label box ends at (x + col_w - 1.74in); value box starts at
        # (x + col_w - 1.78in) — intentional 0.04in overlap so text regions
        # butt against each other without a visible gap at compact density.
        B.add_tb(slide, node.get("value", ""), x + col_w - 1.78, ty, 0.95,
                 0.3, size=11, color=pal["text_muted"], align=PP_ALIGN.RIGHT,
                 font=pal["font_body"])
        if delta:
            B.add_tb(slide, delta, x + col_w - 0.82, ty, 0.7, 0.3, size=11,
                     bold=True, color=d_color, align=PP_ALIGN.RIGHT,
                     font=pal["font_body"])


def _tree_connectors(slide, pal, children, depths, ys, col_w, gap):
    """Elbow connectors parent right edge -> child left edge, text_muted."""
    from palettes import hex_rgb
    for par, kids in children.items():
        if par not in depths:
            continue
        px = 0.7 + (depths[par] - 1) * (col_w + gap) + col_w
        for kid in kids:
            if kid not in depths:
                continue
            kx = 0.7 + (depths[kid] - 1) * (col_w + gap)
            x1, y1, _, _ = B._sc(px, ys[par], 0, 0)
            x2, y2, _, _ = B._sc(kx, ys[kid], 0, 0)
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.ELBOW, Inches(x1), Inches(y1),
                Inches(x2), Inches(y2))
            conn.line.color.rgb = hex_rgb(pal["text_muted"])
            conn.line.width = Pt(1.25)


LAYOUTS["driver-tree"] = build_driver_tree_slide


# ── stakeholder map (support x influence) ────────────────────────────────────
def build_stakeholder_map_slide(prs, p, pal, ctx):
    """2x2 matrix engine with stakeholder axis defaults; items carrying
    TargetX/TargetY get a thin arrow from current position to a hollow
    (outlined) target circle — the engagement move to make."""
    q = dict(p)
    q.setdefault("x_axis", "Support →")
    q.setdefault("y_axis", "Influence →")
    slide = build_matrix_slide(prs, q, pal, ctx)

    L, T, W, H = 2.0, 1.85, 9.6, 4.6  # build_matrix_slide plot area
    from charts import add_arrowhead_to_connector
    from palettes import hex_rgb

    def place(fx, fy):  # fraction -> clamped plot coordinates
        return (max(L + 0.09, min(L + W - 0.09, L + fx * W)),
                max(T + 0.09, min(T + H - 0.09, T + (1 - fy) * H)))

    for item in p.get("matrix_items", [])[:12]:
        if not (item.get("targetx") or item.get("targety")):
            continue
        try:
            fx, fy = float(item.get("x", "")), float(item.get("y", ""))
            tx, ty = float(item["targetx"]), float(item["targety"])
        except (KeyError, ValueError):
            warn(f"stakeholder-map item {item.get('name', '')!r} needs "
                 "numeric X/Y and both TargetX/TargetY — arrow skipped")
            continue
        x0, y0 = place(fx, fy)
        x1, y1 = place(tx, ty)
        # trim both arrow ends so it clears the dot and the hollow ring
        ax0, ay0, ax1, ay1 = x0, y0, x1, y1
        dist = ((x1 - x0) ** 2 + (y1 - y0) ** 2) ** 0.5
        if dist > 0.3:
            ux, uy = (x1 - x0) / dist, (y1 - y0) / dist
            ax0, ay0 = x0 + ux * 0.13, y0 + uy * 0.13
            ax1, ay1 = x1 - ux * 0.12, y1 - uy * 0.12
        bx, by, _, _ = B._sc(ax0, ay0, 0, 0)
        ex, ey, _, _ = B._sc(ax1, ay1, 0, 0)
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                          Inches(bx), Inches(by),
                                          Inches(ex), Inches(ey))
        conn.line.color.rgb = hex_rgb(pal["text_muted"])
        conn.line.width = Pt(1.5)
        add_arrowhead_to_connector(conn)
        ring = B.add_circle(slide, x1 - 0.09, y1 - 0.09, 0.18, pal["bg"],
                            line_hex=pal["accent1"], line_pt=1.5)
        ring.fill.background()  # hollow: future position, not a second dot
    return slide


LAYOUTS["stakeholder-map"] = build_stakeholder_map_slide


# ── raci ─────────────────────────────────────────────────────────────────────
RACI_CHIP_KEYS = {"R": "accent1", "A": "accent2", "C": "accent3",
                  "I": "text_muted"}


def build_raci_slide(prs, p, pal, ctx):
    """RACI chart: activity rows x person columns, colored letter chips
    (R=accent1, A=accent2, C=accent3, I=text_muted; blank = no chip)."""
    rows = p.get("table_rows", [])
    if len(rows) < 2:
        warn("raci needs a markdown table; rendering as table")
        return B.build_table_slide(prs, p, pal, ctx)
    slide = B._blank_slide(prs, pal, pal["bg"])
    B._heading(slide, p, pal)
    if len(rows) - 1 > 12:
        warn(f"raci with {len(rows) - 1} activities — consider splitting")

    n_cols = max(len(r) for r in rows)
    first_col_w = 3.6
    col_w = (11.9 - first_col_w) / max(n_cols - 1, 1)
    row_h = max(0.3, min(0.62, 4.4 / len(rows) - 0.06))
    top = 1.9
    for r, row in enumerate(rows):
        y = top + r * (row_h + 0.06)
        if r > 0:
            B.add_rect(slide, 0.7, y, 11.9, row_h,
                       pal["surface"] if r % 2 else pal["bg_deep"])
        for c in range(n_cols):
            cell = (row[c] if c < len(row) else "").strip()
            x = 0.7 + (first_col_w + (c - 1) * col_w if c else 0)
            w = col_w if c else first_col_w
            if r == 0:  # header: activity column + person names
                B.add_tb(slide, cell, x + (0.15 if c == 0 else 0),
                         y + row_h / 2 - 0.16, w, 0.35, size=14, bold=True,
                         color=pal["accent1"],
                         align=PP_ALIGN.CENTER if c else PP_ALIGN.LEFT,
                         font=pal["font_body"])
            elif c == 0:
                B.add_tb(slide, cell, x + 0.15, y + row_h / 2 - 0.16,
                         w - 0.2, 0.35, size=13, color=pal["text"],
                         font=pal["font_body"])
            elif cell:
                key = RACI_CHIP_KEYS.get(cell.upper())
                if not key:  # validator errors on these; render muted text
                    B.add_tb(slide, cell, x, y + row_h / 2 - 0.16, w, 0.35,
                             size=12, color=pal["text_muted"],
                             align=PP_ALIGN.CENTER, font=pal["font_body"])
                    continue
                d = max(0.2, min(0.42, row_h - 0.14))
                fill = pal[key]
                B.add_circle(slide, x + w / 2 - d / 2, y + row_h / 2 - d / 2,
                             d, fill)
                B.add_tb(slide, cell.upper(), x + w / 2 - d / 2,
                         y + row_h / 2 - 0.15, d, 0.3, size=12, bold=True,
                         color=_label_color_on(fill, pal),
                         align=PP_ALIGN.CENTER, font=pal["font_body"])
    legend_y = min(6.45, top + len(rows) * (row_h + 0.06) + 0.08)
    B.add_tb(slide,
             "R = Responsible · A = Accountable · C = Consulted · I = Informed",
             0.7, legend_y, 11.9, 0.3, size=11, color=pal["text_muted"],
             font=pal["font_label"])
    return slide


LAYOUTS["raci"] = build_raci_slide
