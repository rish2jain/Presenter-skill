#!/usr/bin/env python3
"""
profile_template.py — Analyze a .pptx template and output a JSON summary
of slide layouts, placeholders, and theme colors.

Usage:
    python scripts/profile_template.py template.pptx
    python scripts/profile_template.py template.pptx --generate-config
"""
import argparse
import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from template_helpers import build_template_map

def ph_type_name(ph_type):
    """PP_PLACEHOLDER enum member -> 'TITLE', 'BODY', 'PICTURE', ..."""
    try:
        return ph_type.name
    except AttributeError:
        return str(ph_type)


def emu_to_inches(emu):
    return round(emu / 914400, 3)


def extract_theme_colors(prs):
    """Read the color scheme (dk1/lt1/accent1-6...) from the master's theme."""
    from lxml import etree
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    colors = {}
    try:
        theme_part = prs.slide_masters[0].part.part_related_by(RT.THEME)
        root = etree.fromstring(theme_part.blob)
        scheme = root.find(f"{A}themeElements/{A}clrScheme")
        for child in scheme:
            name = etree.QName(child).localname
            srgb = child.find(f"{A}srgbClr")
            sysc = child.find(f"{A}sysClr")
            if srgb is not None:
                colors[name] = srgb.get("val")
            elif sysc is not None:
                colors[name] = sysc.get("lastClr", sysc.get("val"))
    except Exception as e:
        print(f"[WARN] Could not extract theme colors: {e}", file=sys.stderr)
    return colors


def analyze(pptx_path):
    prs = Presentation(str(pptx_path))
    mapping = build_template_map(prs)
    result = {
        "slide_width_in": emu_to_inches(prs.slide_width),
        "slide_height_in": emu_to_inches(prs.slide_height),
        "slide_count": len(prs.slides),
        "theme_colors": extract_theme_colors(prs),
        "layout_map": mapping["layout_map"],
        "layout_scores": mapping["layout_scores"],
        "layout_index": mapping["layout_index"],
        "slide_layouts": [],
    }

    for i, layout in enumerate(prs.slide_layouts):
        ph_list = []
        for ph in layout.placeholders:
            ph_list.append({
                "idx": ph.placeholder_format.idx,
                "type": ph_type_name(ph.placeholder_format.type),
                "name": ph.name,
                "left_in": emu_to_inches(ph.left),
                "top_in": emu_to_inches(ph.top),
                "width_in": emu_to_inches(ph.width),
                "height_in": emu_to_inches(ph.height),
            })
        result["slide_layouts"].append({
            "index": i,
            "name": layout.name,
            "placeholder_count": len(ph_list),
            "placeholders": ph_list,
        })

    return result


def main():
    parser = argparse.ArgumentParser(description="Profile a .pptx template.")
    parser.add_argument("pptx", help="Path to .pptx template")
    parser.add_argument("--generate-config", action="store_true",
                        help="Write advisory JSON alongside the .pptx (for agent planning; not auto-loaded by build)")
    args = parser.parse_args()

    data = analyze(args.pptx)
    print(json.dumps(data, indent=2))

    if args.generate_config:
        out = Path(args.pptx).with_suffix(".config.json")
        out.write_text(json.dumps(data, indent=2), encoding="utf-8")
        print(f"\nConfig saved: {out}", file=sys.stderr)


if __name__ == "__main__":
    main()
