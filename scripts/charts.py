"""Native python-pptx chart helpers, styled for dark/light palette themes.

The critical detail: default charts render with an opaque white chart/plot
area, which looks broken on dark slides. make_chart_transparent() clears both.
"""
import logging
from importlib.metadata import PackageNotFoundError, version
from lxml import etree
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

from palettes import hex_rgb

logger = logging.getLogger(__name__)

CHART_TYPES = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "hbar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "stacked-100": XL_CHART_TYPE.COLUMN_STACKED_100,
}

LABEL_MODES = ("pct", "abs", "both")


def round_to_sum(values, total=100, decimals=0):
    """Largest-remainder rounding: displayed values sum to exactly `total`.

    Precondition: values should be percentage-normalized shares that already
    sum approximately to `total` (callers pass v/total*100 fractions).
    Ties go to later elements ([33.33]*3 -> [33, 33, 34]). Any negative
    value falls back to plain rounding — largest-remainder doesn't apply.
    """
    import math
    if any(v < 0 for v in values):
        return [round(v, decimals) if decimals else round(v) for v in values]
    factor = 10 ** decimals
    scaled = [v * factor for v in values]
    floors = [math.floor(s) for s in scaled]
    short = int(round(total * factor)) - sum(floors)
    # distribute the shortfall to the largest remainders (later index wins ties)
    order = sorted(range(len(values)),
                   key=lambda i: (scaled[i] - floors[i], i), reverse=True)
    if short >= 0:
        for i in order[:short]:
            floors[i] += 1
    else:
        for i in order[short:]:  # smallest remainders lose first
            floors[i] -= 1
    return [f / factor for f in floors] if decimals else floors

_NO_FILL_SPPR = (
    '<c:spPr xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" '
    'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    "<a:noFill/><a:ln><a:noFill/></a:ln></c:spPr>"
)


def _clear_fill(element, insert_after=None):
    """Ensure element has a <c:spPr> with noFill (replacing any existing fill)."""
    spPr = element.find(qn("c:spPr"))
    if spPr is not None:
        element.remove(spPr)
    spPr = etree.fromstring(_NO_FILL_SPPR)
    if insert_after is not None:
        insert_after.addnext(spPr)
    else:
        element.append(spPr)
    return spPr


def make_chart_transparent(chart):
    """Clear the chart-area and plot-area backgrounds (white-box fix)."""
    chart_space = chart._chartSpace
    chart_el = chart_space.find(qn("c:chart"))
    # chartSpace schema: spPr must directly follow <c:chart>
    _clear_fill(chart_space, insert_after=chart_el)
    plot_area = chart_el.find(qn("c:plotArea"))
    if plot_area is not None:
        _clear_fill(plot_area)


def _style_axes(chart, pal):
    chart.font.size = Pt(12)
    chart.font.color.rgb = hex_rgb(pal["text_muted"])
    chart.font.name = pal["font_body"]
    for axis_name in ("value_axis", "category_axis"):
        try:
            axis = getattr(chart, axis_name)
        except ValueError:
            continue  # pie/doughnut have no axes
        axis.tick_labels.font.size = Pt(12)
        axis.tick_labels.font.color.rgb = hex_rgb(pal["text_muted"])
        axis.format.line.color.rgb = hex_rgb(pal["surface"])
        if axis_name == "value_axis":
            axis.has_major_gridlines = True
            axis.major_gridlines.format.line.color.rgb = hex_rgb(pal["surface"])


def add_native_chart(slide, pal, chart_kind, categories, series,
                     left, top, w, h, series_name=""):
    """Add a palette-styled native chart.

    series: list of (value, ...) floats for a single series, or
            dict {name: values} for multiple series.
    """
    xl_type = CHART_TYPES.get(chart_kind, XL_CHART_TYPE.COLUMN_CLUSTERED)
    chart_data = ChartData()
    chart_data.categories = categories
    if isinstance(series, dict):
        for name, values in series.items():
            chart_data.add_series(name, values)
        n_series = len(series)
    else:
        chart_data.add_series(series_name or "Series 1", series)
        n_series = 1

    frame = slide.shapes.add_chart(
        xl_type, Inches(left), Inches(top), Inches(w), Inches(h), chart_data
    )
    chart = frame.chart
    chart.has_title = False

    make_chart_transparent(chart)
    _style_axes(chart, pal)

    accents = pal.get("chart_series") or [pal["accent1"], pal["accent2"],
                                          pal["accent3"]]
    if chart_kind in ("pie", "doughnut"):
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        for i, point in enumerate(chart.plots[0].series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = hex_rgb(accents[i % len(accents)])
    elif chart_kind == "line":
        chart.has_legend = n_series > 1
        for i, s in enumerate(chart.plots[0].series):
            s.smooth = True
            line = s.format.line
            line.color.rgb = hex_rgb(accents[i % len(accents)])
            line.width = Pt(2.5)
    else:
        chart.has_legend = n_series > 1
        for i, s in enumerate(chart.plots[0].series):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = hex_rgb(accents[i % len(accents)])

    return chart


# ── benchmark / target reference line ────────────────────────────────────────
def _nice_ceil(v):
    """Round up to a 'nice' axis maximum (1/2/2.5/5 x 10^k)."""
    import math
    if v <= 0:
        return 1
    exp = math.floor(math.log10(v))
    for mult in (1, 2, 2.5, 5, 10):
        cand = mult * 10 ** exp
        if cand >= v:
            return cand
    return 10 ** (exp + 1)


def _plot_box(left, top, w, h):
    """Approximate plot-area box within the chart frame: (px, py, pw, ph).
    Insets are approximations — confirm placement in visual QA."""
    return left + 0.09 * w, top + 0.04 * h, w * 0.88, h * 0.80


def _draw_ref_line(slide, pal, value, label, axis_max, left, top, w, h,
                   align=None):
    """Dashed horizontal line + label at `value` on a 0..axis_max value axis.
    Shared geometry for add_benchmark_line and add_value_line."""
    px, py, pw, ph = _plot_box(left, top, w, h)
    y = py + (1 - value / axis_max) * ph

    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(px), Inches(y), Inches(px + pw), Inches(y))
    conn.line.color.rgb = hex_rgb(pal["accent3"])
    conn.line.width = Pt(1.75)
    conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    from pptx.enum.text import PP_ALIGN
    if align is None:
        align = PP_ALIGN.LEFT
    tb_left = px + pw - 2.45 if align == PP_ALIGN.RIGHT else px + 0.05
    tb = slide.shapes.add_textbox(Inches(tb_left), Inches(y - 0.32),
                                  Inches(2.4), Inches(0.28))
    tf = tb.text_frame
    run = tf.paragraphs[0].add_run()
    tf.paragraphs[0].alignment = align
    run.text = label
    run.font.size = Pt(11)
    run.font.bold = True
    # text_muted clears 4.5:1 on every palette (accent3 does not)
    run.font.color.rgb = hex_rgb(pal["text_muted"])
    run.font.name = pal["font_label"]


def add_benchmark_line(slide, chart, pal, spec, left, top, w, h):
    """Dashed reference line across a bar/column/line chart.

    spec: '120 Industry average'. Sets an explicit value-axis maximum so the
    line height is computable; plot-area insets are approximations — confirm
    placement in visual QA.
    """
    import re
    m = re.match(r"\s*\$?([\d,.]+)\s*(.*)", spec)
    if not m:
        return
    value = float(m.group(1).replace(",", ""))
    label = m.group(2).strip().strip('"') or "Benchmark"

    data_max = value
    for s in chart.plots[0].series:
        m = max((v for v in s.values if v is not None), default=None)
        if m is not None:
            data_max = max(data_max, m)
    axis_max = _nice_ceil(data_max * 1.05)
    va = chart.value_axis
    va.maximum_scale = float(axis_max)
    va.minimum_scale = 0.0

    _draw_ref_line(slide, pal, value, label, axis_max, left, top, w, h)


def add_value_line(slide, chart, pal, spec, left, top, w, h):
    """Dashed labelled reference line: '- Value-Line: Target, 40'.

    Unlike add_benchmark_line this never widens the axis: if the value falls
    outside the computed axis range the line is skipped with a warning.
    """
    from helpers import warn
    label, _, value_s = spec.rpartition(",")
    try:
        value = float(value_s.replace("$", "").replace(",", "").strip())
    except ValueError:
        warn(f"Value-Line not parsable: {spec!r} "
             "(expected '- Value-Line: <label>, <value>')")
        return
    label = label.strip().strip('"') or "Value"

    va = chart.value_axis
    axis_max = va.maximum_scale
    if axis_max is None:
        data_max = 0
        for s in chart.plots[0].series:
            m = max((v for v in s.values if v is not None), default=None)
            if m is not None:
                data_max = max(data_max, m)
        axis_max = _nice_ceil(data_max * 1.05)
        va.maximum_scale = float(axis_max)
        va.minimum_scale = 0.0
    axis_min = va.minimum_scale or 0.0
    if not (axis_min <= value <= axis_max):
        warn(f"Value-Line {label!r} at {value:g} is outside the axis range "
             f"{axis_min:g}-{axis_max:g} — line skipped")
        return

    from pptx.enum.text import PP_ALIGN
    _draw_ref_line(slide, pal, value, label, axis_max, left, top, w, h,
                   align=PP_ALIGN.RIGHT)


def _fmt_label_num(v):
    return f"{v:,.0f}" if v == int(v) else f"{v:,.1f}"


def add_stacked_100_labels(chart, mode, pal):
    """Per-point data labels on a stacked-100 chart: mode pct | abs | both.

    Percentages are computed per column via round_to_sum so each column's
    labels sum to exactly 100.
    """
    series_list = list(chart.plots[0].series)
    if not series_list:
        return
    for j, col in enumerate(zip(*[s.values for s in series_list])):
        col = [v or 0 for v in col]
        total = sum(col)
        pcts = (round_to_sum([v / total * 100 for v in col])
                if total else [0] * len(col))
        for i, s in enumerate(series_list):
            if mode == "pct":
                text = f"{pcts[i]:.0f}%"
            elif mode == "abs":
                text = _fmt_label_num(col[i])
            else:
                text = f"{pcts[i]:.0f}% ({_fmt_label_num(col[i])})"
            tf = s.points[j].data_label.text_frame
            tf.text = text
            run = tf.paragraphs[0].runs[0]
            run.font.size = Pt(11)
            run.font.bold = True
            # bg_deep reads on every palette's accent fills (funnel precedent)
            run.font.color.rgb = hex_rgb(pal["bg_deep"])
            run.font.name = pal["font_label"]


def _pptx_version_string():
    try:
        return version("python-pptx")
    except PackageNotFoundError:
        return None


def add_arrowhead_to_connector(conn, *, warn=False):
    """Add a triangle tailEnd to a connector line, when supported.

    python-pptx has no public connector arrowhead API (as of 1.0.x). When
    ``LineFormat._get_or_add_ln`` exists we append ``a:tailEnd`` via that
    private helper; otherwise this is a no-op.
    """
    line = conn.line
    getter = getattr(line, "_get_or_add_ln", None)
    if getter is None:
        if warn:
            logger.warning(
                "Skipping connector arrowhead: LineFormat._get_or_add_ln unavailable "
                "(python-pptx %s). CAGR arrow renders without tail triangle.",
                _pptx_version_string() or "unknown",
            )
        return False

    # Private-API workaround — remove when python-pptx exposes a public tailEnd
    # setter (e.g. line.end_arrow_type). See add_arrowhead_to_connector docstring.
    ln = getter()
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    return True


def _add_arrowhead(conn):
    # Delegates to add_arrowhead_to_connector; see there for private-API notes.
    add_arrowhead_to_connector(conn, warn=True)


def add_cagr_arrow(slide, chart, pal, left, top, w, h, axis_max=None):
    """CAGR arrow from the first to the last column of a single-series chart.

    Uses the same plot-box approximation as add_benchmark_line — confirm
    placement in visual QA.
    """
    values = [v for v in chart.plots[0].series[0].values if v is not None]
    if len(values) < 2 or values[0] <= 0 or values[-1] <= 0:
        return
    periods = len(values) - 1
    cagr = (values[-1] / values[0]) ** (1 / periods) - 1

    va = chart.value_axis
    if axis_max is None:
        existing = va.maximum_scale  # set earlier by add_benchmark_line, or None
        new_max = _nice_ceil(max(values) * 1.05)
        axis_max = max(existing, new_max) if existing else new_max
    va.maximum_scale = float(axis_max)
    va.minimum_scale = 0.0

    px, pw = left + 0.09 * w, w * 0.88
    py, ph = top + 0.04 * h, h * 0.80
    n = len(values)
    x0 = px + pw * (0.5 / n)
    x1 = px + pw * ((n - 0.5) / n)
    y0 = py + (1 - values[0] / axis_max) * ph - 0.22
    y1 = py + (1 - values[-1] / axis_max) * ph - 0.22

    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0), Inches(x1), Inches(y1))
    conn.line.color.rgb = hex_rgb(pal["accent1"])
    conn.line.width = Pt(2.0)
    _add_arrowhead(conn)

    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(
        Inches((x0 + x1) / 2 - 1.1), Inches(min(y0, y1) - 0.34),
        Inches(2.2), Inches(0.3))
    run = tb.text_frame.paragraphs[0].add_run()
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run.text = f"CAGR {cagr:+.1%}"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = hex_rgb(pal["text"])
    run.font.name = pal["font_body"]
