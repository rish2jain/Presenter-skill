"""Layout builders — one function per layout name in the outline vocabulary.

Every documented layout in references/generation-guide.md has a builder here;
LAYOUT_MAP at the bottom is the single source of truth for valid layout names.
Builder signature: build_x(prs, p, pal, ctx) -> slide
  p   = parsed slide dict from build_deck.parse_outline
  ctx = {"outline_dir": Path, "assets_dir": Path}
"""
import re

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from charts import add_native_chart
from helpers import (_add_filled_shape, add_picture_cover, add_soft_shadow,
                     parse_visual, resolve_image_path, set_fill_alpha,
                     set_slide_bg_gradient, warn)
from helpers import add_circle as _add_circle
from helpers import add_overlay as _add_overlay
from helpers import add_picture_contain as _add_picture_contain
from helpers import add_rect as _add_rect
from helpers import add_tb as _add_tb

# Builders design on a 13.33 x 7.5 canvas; set_canvas() rescales everything
# to the actual slide size (e.g. a user template that is 10 x 7.5).
SLIDE_W, SLIDE_H = 13.33, 7.5
_SCALE = {"x": 1.0, "y": 1.0, "font": 1.0}

# Density: compact (default — double-density: bullets flow into two columns,
# exec points into a 2-column card grid) vs comfortable (single column).
# Geometry chosen so the last row clears the footer zone (y >= 7.0).
DENSITY_TABLE = {
    "compact": {
        "bullet_max": 14, "bullet_cols": 2, "bullet_step": 0.74,
        "bullet_y": 1.70, "bullet_size": 14,
        "twocol_max": 8, "twocol_step": 0.585, "twocol_size": 13,
        "exec_max": 8, "exec_cols": 2, "exec_h": 1.02, "exec_gap": 0.12,
        "exec_y": 1.75, "exec_size": 13,
        "cmp_max": 9, "cmp_step": 0.46, "cmp_size": 12,
    },
    "comfortable": {
        "bullet_max": 7, "bullet_cols": 1, "bullet_step": 0.74,
        "bullet_y": 1.70, "bullet_size": 16,
        "twocol_max": 5, "twocol_step": 0.88, "twocol_size": 16,
        "exec_max": 5, "exec_cols": 1, "exec_h": 0.95, "exec_gap": 0.10,
        "exec_y": 1.75, "exec_size": 14,
        "cmp_max": 6, "cmp_step": 0.70, "cmp_size": 14,
    },
}
DEFAULT_DENSITY = "compact"
_DENSITY = {"mode": DEFAULT_DENSITY}


def set_density(mode):
    _DENSITY["mode"] = mode if mode in DENSITY_TABLE else DEFAULT_DENSITY


def density():
    return DENSITY_TABLE[_DENSITY["mode"]]


def set_canvas(prs):
    _SCALE["x"] = prs.slide_width / Inches(1) / SLIDE_W
    _SCALE["y"] = prs.slide_height / Inches(1) / SLIDE_H
    _SCALE["font"] = min(_SCALE["x"], _SCALE["y"], 1.0)


def _sc(left, top, w, h):
    return (left * _SCALE["x"], top * _SCALE["y"],
            w * _SCALE["x"], h * _SCALE["y"])


def _fs(size):
    return max(round(size * _SCALE["font"]), 10)  # never below 10pt


def add_tb(slide, text, left, top, w, h, size=16, **kw):
    return _add_tb(slide, text, *_sc(left, top, w, h), size=_fs(size), **kw)


def add_rect(slide, left, top, w, h, *args, **kw):
    return _add_rect(slide, *_sc(left, top, w, h), *args, **kw)


def add_round_rect(slide, left, top, w, h, *args, **kw):
    return _add_filled_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                             *_sc(left, top, w, h), *args, **kw)


def add_overlay(slide, left, top, w, h, *args, **kw):
    return _add_overlay(slide, *_sc(left, top, w, h), *args, **kw)


def add_picture_contain(slide, img_path, left, top, w, h, **kw):
    return _add_picture_contain(slide, img_path, *_sc(left, top, w, h), **kw)


def add_circle(slide, left, top, d, *args, **kw):
    l, t, w, _ = _sc(left, top, d, d)
    return _add_circle(slide, l, t, w, *args, **kw)


def _blank_slide(prs, pal, bg_hex, accent_bar=True, gradient=False):
    layout = ctx_blank_layout(prs)
    slide = prs.slides.add_slide(layout)
    if gradient and pal["dark"]:
        set_slide_bg_gradient(slide, pal["bg"], pal["bg_deep"])
    else:
        fill = slide.background.fill
        fill.solid()
        from palettes import hex_rgb
        fill.fore_color.rgb = hex_rgb(bg_hex)
    if accent_bar:
        add_rect(slide, 0, 0, SLIDE_W, 0.05, pal["accent1"])
    return slide


# ── Icon bullets ─────────────────────────────────────────────────────────────
ICON_RX = re.compile(r"^icon:([a-z0-9-]+)\s+(.*)$")


def split_icon(text):
    """'icon:rocket Launch fast' -> ('rocket', 'Launch fast'); else (None, text)."""
    m = ICON_RX.match(text.strip())
    return (m.group(1), m.group(2)) if m else (None, text)


def _icon_path(name, pal, ctx):
    try:
        from fetch_icon import fetch_icon
        from pathlib import Path
        cache = Path(ctx.get("assets_dir", "assets")) / "icons"
        return fetch_icon(name, pal["accent1"], cache)
    except Exception as e:
        warn(f"icon '{name}': {e}")
        return None


def _draw_marker(slide, pal, ctx, icon, x, y):
    """Bullet marker: icon-in-circle when given, accent square otherwise.
    Returns the x offset where the bullet text should start."""
    if icon:
        path = _icon_path(icon, pal, ctx)
        if path:
            circle_fill = pal["surface"] if pal["dark"] else pal["bg_deep"]
            add_circle(slide, x, y - 0.06, 0.46, circle_fill)
            add_picture_contain(slide, path, x + 0.10, y + 0.04, 0.26, 0.26)
            return x + 0.66
    add_rect(slide, x, y + 0.05, 0.12, 0.12, pal["accent1"])
    return x + 0.30


def ctx_blank_layout(prs):
    """Find the most blank-like layout (fewest placeholders); index 6 in the
    default template, but arbitrary templates differ."""
    return min(prs.slide_layouts, key=lambda l: len(l.placeholders))


def _heading(slide, p, pal, w=11.9, size=32, align=PP_ALIGN.LEFT):
    add_tb(slide, p.get("heading", ""), 0.7, 0.55, w, 1.1,
           size=size, bold=True, color=pal["accent1"],
           font=pal["font_title"], align=align)


def _place_visual(slide, prs, p, pal, ctx, left, top, w, h):
    """Fill a box with the slide's visual: chart, image, or styled placeholder."""
    kind, value = parse_visual(p.get("visual", ""))
    if kind == "chart" and p.get("data"):
        cats = [d[0] for d in p["data"]]
        series_field = p.get("series", "")
        multi_series = bool(series_field and "," in series_field)
        if multi_series:
            names = [n.strip() for n in series_field.split(",") if n.strip()]
            series_dict = {
                names[i]: [
                    row[1][i] if isinstance(row[1], list) else row[1]
                    for row in p["data"]
                ]
                for i in range(len(names))
            }
            chart = add_native_chart(slide, pal, value, cats, series_dict,
                                     *_sc(left, top, w, h))
        else:
            vals = [
                row[1][0] if isinstance(row[1], list) else row[1]
                for row in p["data"]
            ]
            chart = add_native_chart(slide, pal, value, cats, vals,
                                     *_sc(left, top, w, h),
                                     series_name=series_field or p.get("heading", ""))
            if p.get("benchmark") and value in ("bar", "column", "line", "area"):
                from charts import add_benchmark_line
                add_benchmark_line(slide, chart, pal, p["benchmark"],
                                   *_sc(left, top, w, h))
            if p.get("cagr") and value in ("bar", "column", "line"):
                from charts import add_cagr_arrow
                try:
                    user_max = float(p["axis_max"]) if p.get("axis_max") else None
                except (ValueError, TypeError):
                    user_max = None
                    warn(f"Axis-Max not numeric: {p.get('axis_max')!r} "
                         f"(chart:{value}, CAGR)")
                add_cagr_arrow(slide, chart, pal, *_sc(left, top, w, h),
                               axis_max=user_max)
        # CAGR path passes axis_max into add_cagr_arrow; don't overwrite after.
        cagr_sets_axis = (
            p.get("cagr") and not multi_series
            and value in ("bar", "column", "line"))
        if p.get("axis_max") and not cagr_sets_axis:
            try:
                chart.value_axis.maximum_scale = float(p["axis_max"])
                chart.value_axis.minimum_scale = 0.0
            except (ValueError, TypeError):
                warn(f"Axis-Max not numeric: {p['axis_max']!r}")
        if p.get("labels_mode") and value == "stacked-100":
            from charts import LABEL_MODES, add_stacked_100_labels
            if p["labels_mode"] in LABEL_MODES:
                add_stacked_100_labels(chart, p["labels_mode"], pal)
            else:
                warn(f"Labels mode {p['labels_mode']!r} not one of "
                     f"{'|'.join(LABEL_MODES)} — labels skipped")
        # after Axis-Max so the line is placed against the final axis range
        if p.get("value_line") and value in ("bar", "column", "line"):
            from charts import add_value_line
            add_value_line(slide, chart, pal, p["value_line"],
                           *_sc(left, top, w, h))
        return
    if kind == "image":
        path = resolve_image_path(value, ctx)
        if path:
            add_picture_contain(slide, path, left, top, w, h,
                                alt=p.get("caption") or p.get("heading"))
            return
        warn(f"Image not found: {value}")
    _visual_placeholder(slide, pal, left, top, w, h)


def _visual_placeholder(slide, pal, left, top, w, h):
    add_rect(slide, left, top, w, h, pal["surface"],
             line_hex=pal["accent1"], line_pt=1.0)
    add_tb(slide, "[ Add Visual ]", left, top + h / 2 - 0.3, w, 0.6,
           size=14, color=pal["text_muted"], align=PP_ALIGN.CENTER,
           font=pal["font_label"])


GHOST_KEEP_REAL = {"title", "section-divider", "closing", "agenda"}


def build_ghost_slide(prs, p, pal, ctx):
    """Skeleton slide: real heading + dashed placeholder describing the
    planned exhibit. Used by build_deck --ghost for storyline alignment."""
    slide = _blank_slide(prs, pal, pal["bg"])
    _heading(slide, p, pal)
    kind, value = parse_visual(p.get("visual", ""))
    label = f"[ {p.get('layout', 'bullet-list')} ]"
    if kind:
        label += f"   planned visual — {kind}: {value}"
    elif any(p.get(k) for k in ("data", "bars", "matrix_items", "tiles",
                                "milestones", "table_rows", "steps")):
        label += "   planned exhibit from structured rows"
    box = add_rect(slide, 0.7, 1.9, 11.9, 4.6, pal["surface"],
                   line_hex=pal["text_muted"], line_pt=1.0)
    set_fill_alpha(box, 40)
    add_tb(slide, label, 0.9, 2.05, 11.5, 0.4, size=14, bold=True,
           color=pal["text_muted"], font=pal["font_label"])
    for i, b in enumerate(p.get("bullets", [])[:6]):
        add_tb(slide, f"—  {split_icon(b)[1]}", 1.0, 2.7 + i * 0.55, 11.2,
               0.5, size=12, color=pal["text_muted"], font=pal["font_body"])
    return slide


# ── Layouts ──────────────────────────────────────────────────────────────────
def _slide_title(p, default=""):
    """Title text: explicit - Title: wins, then ## Slide N: heading."""
    return p.get("title") or p.get("heading") or default


def build_title_slide(prs, p, pal, ctx):
    # No accent line under the title — a known hallmark of AI-generated
    # slides; hierarchy comes from size contrast and whitespace instead.
    slide = _blank_slide(prs, pal, pal["bg_deep"], gradient=True)
    add_rect(slide, 0, 0, SLIDE_W, 0.07, pal["accent1"])

    has_hero = False
    kind, value = parse_visual(p.get("visual", ""))
    if kind == "image":
        path = resolve_image_path(value, ctx)
        if path:
            add_picture_contain(slide, path, 9.3, 1.4, 3.4, 4.7,
                                alt=_slide_title(p, "hero image"))
            has_hero = True

    title_w = 8.2 if has_hero else 11.7
    add_tb(slide, _slide_title(p, "Presentation Title"),
           0.8, 2.1, title_w, 1.8, size=46, bold=True,
           color=pal["accent1"], font=pal["font_title"])
    if p.get("subtitle"):
        add_tb(slide, p["subtitle"], 0.8, 4.15, title_w, 0.8, size=20,
               color=pal["text_muted"], font=pal["font_label"])
    return slide


def build_section_divider_slide(prs, p, pal, ctx):
    slide = _blank_slide(prs, pal, pal["bg_deep"], accent_bar=False, gradient=True)
    add_rect(slide, 0, 3.0, 0.18, 1.5, pal["accent1"])
    add_tb(slide, p.get("heading", p.get("title", "")), 0.8, 2.9, 11.7, 1.7,
           size=44, bold=True, color=pal["text"], font=pal["font_title"])
    if p.get("subtitle"):
        add_tb(slide, p["subtitle"], 0.8, 4.6, 11.7, 0.8, size=18,
               color=pal["text_muted"], font=pal["font_label"])
    return slide


def build_exec_summary_slide(prs, p, pal, ctx):
    d = density()
    slide = _blank_slide(prs, pal, pal["bg"])
    _heading(slide, p, pal)
    points = p.get("bullets", [])[:d["exec_max"]]
    cols = d.get("exec_cols", 1)
    per_col = -(-len(points) // cols) if points else 1
    gap_x = 0.3
    card_w = (11.9 - gap_x * (cols - 1)) / cols
    card_h = d["exec_h"]
    for i, point in enumerate(points):
        icon, text = split_icon(point)
        col, row = divmod(i, per_col)  # column-major: read down, then across
        x = 0.7 + col * (card_w + gap_x)
        y = d["exec_y"] + row * (card_h + d["exec_gap"])
        card = add_rect(slide, x, y, card_w, card_h, pal["surface"])
        add_soft_shadow(card)
        add_rect(slide, x, y, 0.08, card_h, pal["accent1"])
        text_x = x + 0.35
        if icon:
            text_x = _draw_marker(slide, pal, ctx, icon, x + 0.28,
                                  y + card_h / 2 - 0.18) + 0.12
        add_tb(slide, text, text_x, y + card_h / 2 - 0.30,
               x + card_w - text_x - 0.2, 0.62,
               size=d["exec_size"], color=pal["text"], font=pal["font_body"],
               accent=pal["accent1"])
    return slide


def build_bullet_slide(prs, p, pal, ctx):
    d = density()
    slide = _blank_slide(prs, pal, pal["bg"])
    _heading(slide, p, pal)
    bullets = p.get("bullets", [])[:d["bullet_max"]]
    cols = d.get("bullet_cols", 1)
    per_col = -(-len(bullets) // cols) if bullets else 1
    col_gap = 0.4
    col_w = (11.93 - col_gap * (cols - 1)) / cols
    for i, bullet in enumerate(bullets):
        icon, text = split_icon(bullet)
        col, row = divmod(i, per_col)  # column-major: read down, then across
        x0 = 0.7 + col * (col_w + col_gap)
        y = d["bullet_y"] + row * d["bullet_step"]
        text_x = _draw_marker(slide, pal, ctx, icon, x0, y)
        add_tb(slide, text, text_x, y - 0.06, x0 + col_w - text_x,
               d["bullet_step"] - 0.06,
               size=d["bullet_size"], color=pal["text"], font=pal["font_body"],
               accent=pal["accent1"])
    return slide


def build_two_column_slide(prs, p, pal, ctx):
    slide = _blank_slide(prs, pal, pal["bg"])
    add_tb(slide, p.get("heading", ""), 0.7, 0.55, 6.0, 1.1,
           size=30, bold=True, color=pal["accent1"], font=pal["font_title"])
    d = density()
    bullets = p.get("bullets", [])[:d["twocol_max"]]
    for i, bullet in enumerate(bullets):
        icon, text = split_icon(bullet)
        y = 1.74 + i * d["twocol_step"]
        text_x = _draw_marker(slide, pal, ctx, icon, 0.7, y)
        add_tb(slide, text, text_x, y - 0.06, 6.8 - text_x,
               d["twocol_step"] - 0.08,
               size=d["twocol_size"], color=pal["text"], font=pal["font_body"],
               accent=pal["accent1"])
    _place_visual(slide, prs, p, pal, ctx, 6.9, 0.9, 6.0, 5.9)
    return slide


def _build_cards(prs, p, pal, ctx, per_row, rows):
    slide = _blank_slide(prs, pal, pal["bg"])
    _heading(slide, p, pal)
    cards = p.get("cards", [])[: per_row * rows]
    gap = 0.4
    card_w = (11.9 - gap * (per_row - 1)) / per_row
    card_h = 4.6 if rows == 1 else 2.55
    for i, card in enumerate(cards):
        row, col = divmod(i, per_row)
        x = 0.7 + col * (card_w + gap)
        y = 1.9 + row * (card_h + 0.35)
        shape = add_rect(slide, x, y, card_w, card_h, pal["surface"])
        add_soft_shadow(shape)
        y_text = y + 0.3
        if card.get("icon"):
            path = _icon_path(card["icon"], pal, ctx)
            if path:
                circle_fill = pal["bg_deep"] if pal["dark"] else pal["bg"]
                add_circle(slide, x + 0.25, y + 0.25, 0.55, circle_fill)
                add_picture_contain(slide, path, x + 0.37, y + 0.37, 0.31, 0.31)
                y_text = y + 0.3 if rows == 1 else y + 0.25
                x_text = x + 0.95
            else:
                x_text = x + 0.25
        else:
            x_text = x + 0.25
        name = card.get("name") or card.get("title", "")
        role = card.get("title", "") if card.get("name") else card.get("subtitle", "")
        body = card.get("bio") or card.get("body") or card.get("desc", "")
        add_tb(slide, name, x_text, y_text, card_w - (x_text - x) - 0.25, 0.6,
               size=18, bold=True, color=pal["text"], font=pal["font_title"])
        if role:
            add_tb(slide, role, x_text, y_text + 0.55, card_w - (x_text - x) - 0.25,
                   0.45, size=13, bold=True, color=pal["accent1"],
                   font=pal["font_body"])
        if body:
            add_tb(slide, body, x + 0.25, y + 1.4, card_w - 0.5, card_h - 1.6,
                   size=13, color=pal["text_muted"], font=pal["font_body"],
                   accent=pal["accent1"])
    return slide


def build_cards3_slide(prs, p, pal, ctx):
    return _build_cards(prs, p, pal, ctx, per_row=3, rows=1)


def build_cards4_slide(prs, p, pal, ctx):
    return _build_cards(prs, p, pal, ctx, per_row=2, rows=2)


def build_stat_callout_slide(prs, p, pal, ctx):
    stats = p.get("stats", [])[:3]
    if not stats:
        warn("stat-callout slide has no stats; rendering as bullet-list")
        return build_bullet_slide(prs, p, pal, ctx)
    slide = _blank_slide(prs, pal, pal["bg_deep"])
    _heading(slide, p, pal, size=30, align=PP_ALIGN.CENTER)
    card_w = 3.8
    total_w = len(stats) * card_w + (len(stats) - 1) * 0.45
    start_x = (SLIDE_W - total_w) / 2
    for i, stat in enumerate(stats):
        x = start_x + i * (card_w + 0.45)
        card = add_rect(slide, x, 1.9, card_w, 4.0, pal["surface"],
                        line_hex=pal["accent1"], line_pt=1.5)
        add_soft_shadow(card)
        add_tb(slide, stat.get("value", ""), x + 0.2, 2.4, card_w - 0.4, 1.5,
               size=52, bold=True, color=pal["accent1"],
               align=PP_ALIGN.CENTER, font=pal["font_title"])
        add_tb(slide, stat.get("label", ""), x + 0.2, 4.1, card_w - 0.4, 0.7,
               size=16, bold=True, color=pal["text"],
               align=PP_ALIGN.CENTER, font=pal["font_body"])
        add_tb(slide, stat.get("sublabel", ""), x + 0.2, 4.85, card_w - 0.4, 0.6,
               size=12, color=pal["text_muted"],
               align=PP_ALIGN.CENTER, font=pal["font_label"])
    return slide


def build_timeline_slide(prs, p, pal, ctx):
    slide = _blank_slide(prs, pal, pal["bg"])
    _heading(slide, p, pal)
    items = p.get("items", [])[:5]
    if not items:  # fall back to "Date: text" bullets
        for b in p.get("bullets", [])[:5]:
            date, _, rest = b.partition(":")
            items.append({"date": date.strip(), "title": rest.strip()})
    n = max(len(items), 1)
    line_y = 3.4
    add_rect(slide, 0.9, line_y, 11.5, 0.04, pal["surface"])
    step = 11.5 / n
    for i, item in enumerate(items):
        cx = 0.9 + step * i + step / 2
        add_rect(slide, cx - 0.11, line_y - 0.09, 0.22, 0.22, pal["accent1"])
        add_tb(slide, item.get("date", ""), cx - step / 2 + 0.1, line_y - 0.85,
               step - 0.2, 0.5, size=15, bold=True, color=pal["accent1"],
               align=PP_ALIGN.CENTER, font=pal["font_body"])
        add_tb(slide, item.get("title", ""), cx - step / 2 + 0.1, line_y + 0.4,
               step - 0.2, 0.8, size=14, bold=True, color=pal["text"],
               align=PP_ALIGN.CENTER, font=pal["font_body"])
        if item.get("desc"):
            add_tb(slide, item["desc"], cx - step / 2 + 0.1, line_y + 1.25,
                   step - 0.2, 1.6, size=12, color=pal["text_muted"],
                   align=PP_ALIGN.CENTER, font=pal["font_label"])
    return slide


def build_comparison_slide(prs, p, pal, ctx):
    slide = _blank_slide(prs, pal, pal["bg"])
    _heading(slide, p, pal)
    halves = [
        (0.7, p.get("left_label", "Before"), p.get("visual_left"), p.get("left_bullets", [])),
        (6.93, p.get("right_label", "After"), p.get("visual_right"), p.get("right_bullets", [])),
    ]
    for x, label, visual, bullets in halves:
        col_w = 5.7
        add_tb(slide, label, x, 1.7, col_w, 0.6, size=20, bold=True,
               color=pal["accent1"], align=PP_ALIGN.CENTER, font=pal["font_title"])
        if visual:
            path = resolve_image_path(visual, ctx)
            if path:
                add_picture_contain(slide, path, x, 2.45, col_w, 4.4, alt=label)
                continue
            warn(f"Comparison image not found: {visual}")
        if bullets:
            d = density()
            panel = add_rect(slide, x, 2.45, col_w, 4.4, pal["surface"])
            add_soft_shadow(panel)
            for i, b in enumerate(bullets[:d["cmp_max"]]):
                add_tb(slide, b, x + 0.3, 2.70 + i * d["cmp_step"],
                       col_w - 0.6, d["cmp_step"] - 0.05,
                       size=d["cmp_size"], color=pal["text"], font=pal["font_body"],
                       accent=pal["accent1"])
        else:
            _visual_placeholder(slide, pal, x, 2.45, col_w, 4.4)
    return slide


def build_table_slide(prs, p, pal, ctx):
    from palettes import hex_rgb
    slide = _blank_slide(prs, pal, pal["bg"])
    _heading(slide, p, pal)
    rows = p.get("table_rows", [])
    if not rows:
        warn("table slide has no rows; rendering as bullet-list")
        return build_bullet_slide(prs, p, pal, ctx)
    n_rows, n_cols = len(rows), max(len(r) for r in rows)
    table_h = min(0.55 * n_rows, 5.2)
    tl, tt, tw, th = _sc(0.7, 1.9, 11.9, table_h)
    frame = slide.shapes.add_table(n_rows, n_cols, Inches(tl), Inches(tt),
                                   Inches(tw), Inches(th))
    table = frame.table
    for r, row in enumerate(rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_rgb(
                pal["accent1"] if r == 0 else
                (pal["surface"] if r % 2 else pal["bg_deep"]))
            tf = cell.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = row[c] if c < len(row) else ""
            run.font.size = Pt(_fs(14))
            run.font.name = pal["font_body"]
            run.font.bold = r == 0
            run.font.color.rgb = hex_rgb(pal["bg_deep"] if r == 0 else pal["text"])
    return slide


def build_full_image_slide(prs, p, pal, ctx):
    slide = _blank_slide(prs, pal, pal["bg_deep"], accent_bar=False)
    kind, value = parse_visual(p.get("visual", ""))
    placed = False
    if kind == "image":
        path = resolve_image_path(value, ctx)
        if path:
            add_picture_cover(slide, prs, path,
                              alt=p.get("caption", "full-bleed image"))
            placed = True
        else:
            warn(f"Full-image visual not found: {value}")
    caption = p.get("caption", "")
    if caption:
        if placed:
            add_overlay(slide, 0, 5.2, SLIDE_W, 2.3, "000000", alpha_pct=50)
        add_tb(slide, caption, 0.7, 5.5, 11.9, 1.6, size=22,
               color="F1F5F9", font=pal["font_body"])
    return slide


def build_closing_slide(prs, p, pal, ctx):
    slide = _blank_slide(prs, pal, pal["bg_deep"], accent_bar=False, gradient=True)
    add_rect(slide, 0, 0, SLIDE_W, 0.07, pal["accent1"])
    add_rect(slide, 0, 7.43, SLIDE_W, 0.07, pal["accent2"])
    add_tb(slide, _slide_title(p, "Thank You"), 1.0, 2.4, 11.33, 1.9,
           size=54, bold=True, color=pal["accent1"],
           align=PP_ALIGN.CENTER, font=pal["font_title"])
    if p.get("subtitle"):
        add_tb(slide, p["subtitle"], 1.0, 4.6, 11.33, 0.9, size=20,
               color=pal["text_muted"], align=PP_ALIGN.CENTER,
               font=pal["font_label"])
    if p.get("contact"):
        # accent1 + bold: WCAG large-text (14pt bold) and readable on light themes,
        # where accent2 support tones lack contrast
        add_tb(slide, p["contact"], 1.0, 5.6, 11.33, 0.6, size=14, bold=True,
               color=pal["accent1"], align=PP_ALIGN.CENTER,
               font=pal["font_body"])
    return slide


LAYOUT_MAP = {
    "title": build_title_slide,
    "section-divider": build_section_divider_slide,
    "exec-summary": build_exec_summary_slide,
    "bullet-list": build_bullet_slide,
    "bullets": build_bullet_slide,
    "two-column-split": build_two_column_slide,
    "cards-3": build_cards3_slide,
    "cards-4": build_cards4_slide,
    "stat-callout": build_stat_callout_slide,
    "timeline": build_timeline_slide,
    "comparison": build_comparison_slide,
    "table": build_table_slide,
    "full-image": build_full_image_slide,
    "closing": build_closing_slide,
}

# Consulting layouts (waterfall, matrix-2x2, harvey-scorecard, ...) register
# themselves here; builders_consulting imports this module lazily, so the
# import must come after LAYOUT_MAP is defined.
import builders_consulting  # noqa: E402
LAYOUT_MAP.update(builders_consulting.LAYOUTS)
