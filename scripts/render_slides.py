#!/usr/bin/env python3
"""
render_slides.py — Convert a .pptx to per-slide PNG thumbnails for QA inspection.

Usage:
    python scripts/render_slides.py deck.pptx --out assets/qa-thumbs/
    python scripts/render_slides.py deck.pptx --slides 1,3,5 --out assets/qa-thumbs/
    python scripts/render_slides.py deck.pptx --grid --out assets/qa-thumbs/

Requires: LibreOffice (soffice) + poppler (pdftoppm)
Fallback:  python-pptx + Pillow (lower quality but zero external deps)
"""
import argparse
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path


def _pptx_to_pdf(pptx_path, tmp):
    """Convert to PDF in tmp. Prefers a running unoserver (unoconvert) —
    no cold LibreOffice launch per call — falling back to soffice."""
    if shutil.which("unoconvert"):
        pdf = Path(tmp) / (Path(pptx_path).stem + ".pdf")
        result = subprocess.run(
            ["unoconvert", "--convert-to", "pdf", str(pptx_path), str(pdf)],
            capture_output=True, text=True, timeout=180)
        if result.returncode == 0 and pdf.is_file():
            return pdf
        print(f"[WARN] unoconvert failed ({result.stderr.strip()[:120]}) — "
              "falling back to soffice", file=sys.stderr)
    result = subprocess.run(
        ["soffice", "--headless", "--convert-to", "pdf",
         "--outdir", tmp, str(pptx_path)],
        capture_output=True, text=True)
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice failed: {result.stderr}")
    pdf_files = list(Path(tmp).glob("*.pdf"))
    if not pdf_files:
        raise RuntimeError("No PDF output from LibreOffice")
    return pdf_files[0]


def render_with_libreoffice(pptx_path, out_dir, slide_filter=None):
    """Primary renderer: LibreOffice → PDF → pdftoppm → PNG."""
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory() as tmp:
        pdf_path = _pptx_to_pdf(pptx_path, tmp)

        # Convert PDF pages to PNG — per page when filtering, so 1,3,5
        # renders exactly those slides (not the 1-5 range)
        page_runs = ([(p, p) for p in sorted(slide_filter)]
                     if slide_filter else [(None, None)])
        for first, last in page_runs:
            cmd = ["pdftoppm", "-png", "-r", "150",
                   str(pdf_path), str(out_dir / "slide")]
            if first is not None:
                cmd[3:3] = ["-f", str(first), "-l", str(last)]
            result = subprocess.run(cmd, capture_output=True, text=True)
            if result.returncode != 0:
                raise RuntimeError(f"pdftoppm failed: {result.stderr}")

    # List generated files
    imgs = sorted(out_dir.glob("slide-*.png"))
    if not imgs:
        imgs = sorted(out_dir.glob("slide*.png"))
    return imgs


def render_fallback(pptx_path, out_dir, slide_filter=None):
    """Fallback renderer using python-pptx + Pillow (renders text content only).

    WARNING: this is a text dump, NOT a layout preview. Do not use it to judge
    positioning, overlap, images, or charts — install LibreOffice for real QA,
    and rely on scripts/qa_check.py for programmatic checks.
    """
    print("[WARN] Pillow fallback renders text only — unusable for layout QA. "
          "Use qa_check.py and/or install LibreOffice.", file=sys.stderr)
    from pptx import Presentation
    from pptx.util import Pt
    from PIL import Image, ImageDraw, ImageFont

    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))
    W, H = 1333, 750  # 13.33 x 7.5 inches at 100 DPI

    imgs = []
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        if slide_filter and slide_num not in slide_filter:
            continue

        img = Image.new("RGB", (W, H), (20, 20, 40))
        draw = ImageDraw.Draw(img)

        # Extract text from slide shapes
        y_cursor = 20
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    try:
                        font = ImageFont.load_default()
                        draw.text((20, y_cursor), text, fill=(240, 245, 250), font=font)
                    except Exception:
                        draw.text((20, y_cursor), text, fill=(240, 245, 250))
                    y_cursor += 20
                    if y_cursor > H - 20:
                        break

        out_path = out_dir / f"slide-{slide_num:02d}.png"
        img.save(out_path, "PNG")
        imgs.append(out_path)

    return imgs


def create_thumbnail_grid(imgs, out_path, cols=4):
    """Stitch all slide thumbnails into a single grid image."""
    from PIL import Image

    if not imgs:
        print("No images to stitch.", file=sys.stderr)
        return

    sample = Image.open(imgs[0])
    tw, th = sample.width // 2, sample.height // 2
    rows = (len(imgs) + cols - 1) // cols
    grid = Image.new("RGB", (cols * tw + (cols - 1) * 4, rows * th + (rows - 1) * 4), (30, 30, 40))

    from PIL import ImageDraw, ImageFont
    try:
        font = ImageFont.load_default(size=max(14, th // 12))
    except TypeError:  # Pillow < 10 has no size kwarg
        font = ImageFont.load_default()

    for idx, img_path in enumerate(imgs):
        img = Image.open(img_path).resize((tw, th), Image.LANCZOS)
        draw = ImageDraw.Draw(img)
        label = str(idx + 1)
        pad = 6
        bbox = draw.textbbox((0, 0), label, font=font)
        bw, bh = bbox[2] - bbox[0], bbox[3] - bbox[1]
        draw.rectangle((0, 0, bw + 2 * pad, bh + 2 * pad), fill=(0, 0, 0))
        draw.text((pad, pad), label, fill=(255, 255, 255), font=font)
        row, col = divmod(idx, cols)
        grid.paste(img, (col * (tw + 4), row * (th + 4)))

    grid.save(out_path, "PNG")
    print(f"Grid saved: {out_path}")
    return out_path


def main():
    parser = argparse.ArgumentParser(description="Render .pptx slides as PNG thumbnails.")
    parser.add_argument("pptx", help="Path to .pptx file")
    parser.add_argument("--out", default="assets/qa-thumbs", help="Output directory")
    parser.add_argument("--slides", default=None, help="Comma-separated slide numbers (e.g. 1,3,5)")
    parser.add_argument("--fallback", action="store_true", help="Force use of Pillow fallback renderer")
    parser.add_argument("--grid", action="store_true", help="Also save a stitched grid image")
    args = parser.parse_args()

    pptx_path = Path(args.pptx)
    slide_filter = None
    if args.slides:
        slide_filter = set(int(s.strip()) for s in args.slides.split(","))

    print(f"Rendering {pptx_path} → {args.out}/")

    if not args.fallback and shutil.which("pdftoppm") and (
            shutil.which("soffice") or shutil.which("unoconvert")):
        try:
            imgs = render_with_libreoffice(pptx_path, args.out, slide_filter)
            print(f"Rendered {len(imgs)} slides (LibreOffice)")
        except Exception as e:
            print(f"LibreOffice render failed: {e} — falling back to Pillow", file=sys.stderr)
            imgs = render_fallback(pptx_path, args.out, slide_filter)
    else:
        imgs = render_fallback(pptx_path, args.out, slide_filter)
        print(f"Rendered {len(imgs)} slides (Pillow fallback)")

    if args.grid and imgs:
        grid_path = Path(args.out) / "grid.png"
        create_thumbnail_grid(imgs, grid_path)

    for img in imgs:
        print(f"  {img}")


if __name__ == "__main__":
    main()
