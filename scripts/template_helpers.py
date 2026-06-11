"""Template mapping, placeholder naming, and brand-kit extraction."""
import json
import re
from pathlib import Path

from pptx.enum.shapes import PP_PLACEHOLDER as PH

# Outline layout -> template layout category (matches build_deck.TEMPLATE_WANT)
LAYOUT_CATEGORIES = {
    "title": "title",
    "closing": "title",
    "section-divider": "section",
    "bullet-list": "content",
    "bullets": "content",
    "exec-summary": "content",
    "two-column-split": "picture",
    "full-image": "picture",
    "comparison": "content",
    "stat-callout": "content",
    "timeline": "content",
    "table": "content",
    "cards-3": "content",
    "cards-4": "content",
}


def _layout_score(layout, want):
    types = set()
    for ph in layout.placeholders:
        try:
            types.add(ph.placeholder_format.type)
        except (ValueError, AttributeError):
            pass
    name = layout.name.lower()
    score = 0
    if want == "title":
        if PH.CENTER_TITLE in types:
            score += 4
        elif PH.TITLE in types and PH.SUBTITLE in types:
            score += 3
        if "title" in name:
            score += 1
    elif want == "content":
        if PH.TITLE in types and (PH.BODY in types or PH.OBJECT in types):
            score += 3
        if PH.CENTER_TITLE in types:
            score -= 2
        if "content" in name:
            score += 1
        if "section" in name:
            score -= 2
    elif want == "section":
        if "section" in name:
            score += 3
        if PH.TITLE in types:
            score += 1
    elif want == "picture":
        if PH.PICTURE in types:
            score += 3
        if "picture" in name or "image" in name:
            score += 1
    return score


def find_best_layout(prs, want, config=None):
    """Return (layout, score). Uses config layout_map when present."""
    if config:
        idx = config.get("layout_map", {}).get(want)
        if idx is not None and 0 <= idx < len(prs.slide_layouts):
            layout = prs.slide_layouts[idx]
            return layout, _layout_score(layout, want) + 10
    best, best_score = None, 0
    for layout in prs.slide_layouts:
        s = _layout_score(layout, want)
        if s > best_score:
            best, best_score = layout, s
    return best, best_score


def label_placeholder_names(slide, layout):
    """Copy layout placeholder names onto slide shapes (community workaround)."""
    for shape in slide.placeholders:
        if not shape.is_placeholder:
            continue
        try:
            idx = shape.placeholder_format.idx
            layout_ph = layout.placeholders.get(idx)
            if layout_ph is not None:
                shape.name = layout_ph.name
        except (ValueError, AttributeError, KeyError):
            pass


def build_template_map(prs):
    """Build advisory layout_map JSON from placeholder scoring."""
    categories = ("title", "content", "section", "picture")
    layout_map = {}
    scores = {}
    for cat in categories:
        layout, score = find_best_layout(prs, cat)
        if layout is not None:
            layout_map[cat] = prs.slide_layouts.index(layout)
            scores[cat] = score
    return {
        "layout_map": layout_map,
        "layout_scores": scores,
        "layout_index": [
            {"index": i, "name": lo.name,
             "placeholder_count": len(lo.placeholders)}
            for i, lo in enumerate(prs.slide_layouts)
        ],
    }


def load_template_config(template_path):
    """Load sidecar .config.json if profile_template wrote one."""
    path = Path(template_path).with_suffix(".config.json")
    if not path.is_file():
        return None
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, OSError):
        return None


def _blend(hex_a, hex_b, t):
    """Blend hex_a toward hex_b by fraction t (0..1)."""
    a = [int(hex_a[i:i + 2], 16) for i in (0, 2, 4)]
    b = [int(hex_b[i:i + 2], 16) for i in (0, 2, 4)]
    return "".join(f"{round(x + (y - x) * t):02X}" for x, y in zip(a, b))


def palette_from_theme(theme_colors):
    """Map PPTX theme clrScheme names to a palette dict (best-effort).

    Neutrals (bg/surface/text_muted) are derived as blends of dk1/lt1 — never
    from saturated theme accents, which produces accent-on-accent contrast
    failures (e.g. blue headers on red cards).
    """
    if not theme_colors:
        return None
    accent = theme_colors.get("accent1", "2563EB")
    dark = theme_colors.get("dk1", "0F172A")
    light = theme_colors.get("lt1", "F8FAFC")
    return {
        "bg": dark,
        "bg_deep": _blend(dark, "000000", 0.35),
        "surface": _blend(dark, light, 0.12),
        "accent1": accent,
        "accent2": theme_colors.get("accent2", accent),
        "accent3": theme_colors.get("accent3", accent),
        "text": light,
        "text_muted": _blend(light, dark, 0.32),
        "dark": True,
        "font_title": "Calibri",
        "font_body": "Calibri",
        "font_label": "Calibri Light",
        "motif": "icon-circle",
        "_from_template": True,
    }


def validate_template_mapping(prs, slides_data, config=None):
    """Warn when template lacks good layouts for outline slide types."""
    warnings = []
    for n, p in enumerate(slides_data, 1):
        layout_name = p.get("layout", "bullet-list")
        want = LAYOUT_CATEGORIES.get(layout_name)
        if not want:
            continue
        _, score = find_best_layout(prs, want, config)
        if score < 2:
            warnings.append(
                f"Slide {n}: template has weak match for '{layout_name}' "
                f"(category '{want}', score {score}) — may fall back to blank layout")
    return warnings
