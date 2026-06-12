#!/usr/bin/env python3
"""
pptx_lint.py — deterministic cross-slide consistency lint for a .pptx.

Complements qa_check.py (per-slide defects) with deck-wide checks in the
Macabacus "Deck Check" / UpSlide "Slide Check" class:

  1. Anti-jiggle: recurring elements (page numbers, footers, kickers) must sit
     at identical coordinates on every slide.
  2. Page-number sequence: consecutive, no gaps or duplicates.
  3. Font inventory: more than MAX_FONTS distinct fonts reads as inconsistent.
  4. Color inventory / palette whitelist: with --palette, any explicit run or
     fill color outside the palette (plus known extras) is an error.

Usage:
    python3 scripts/pptx_lint.py deck.pptx [--palette midnight-executive]
"""
import argparse
import re
import sys
from collections import Counter, defaultdict
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation
from qa_check import iter_shapes, _solid_fill_rgb

EMU_IN = 914400
POS_TOL_IN = 0.03           # jiggle tolerance
FOOTER_ZONE = 0.80          # top > 80% of slide height = footer zone
HEADER_ZONE = 0.12          # top < 12% = header zone (kickers, stamps)
MAX_FONTS = 4
MAX_COLORS = 14
PAGENUM_RX = re.compile(r"^(B·)?\d+$")
# waterfall negative red, pure white/black (cards, shadows) are always allowed
EXTRA_ALLOWED = {"D9655B", "FFFFFF", "000000"}
_HEX6 = re.compile(r"[0-9A-Fa-f]{6}")


def _chart_series_allowed(pal, issues):
    """Return uppercased hex colors from chart_series; warn on bad entries."""
    allowed = set()
    series = pal.get("chart_series")
    if series is None:
        return allowed
    if not isinstance(series, list):
        issues["warn"].append(
            f"palette chart_series is not a list ({type(series).__name__}) — "
            "skipped for whitelist")
        return allowed
    for i, item in enumerate(series):
        if not isinstance(item, str):
            issues["warn"].append(
                f"palette chart_series[{i}] is not a string "
                f"({type(item).__name__}) — skipped")
            continue
        if not _HEX6.fullmatch(item):
            issues["warn"].append(
                f"palette chart_series[{i}] is not a 6-digit hex color — skipped")
            continue
        allowed.add(item.upper())
    return allowed


def _shape_box_in(shape):
    try:
        if None in (shape.left, shape.top):
            return None
        return (shape.left / EMU_IN, shape.top / EMU_IN)
    except (AttributeError, TypeError):
        return None


def _iter_text_shapes(prs):
    for n, slide in enumerate(prs.slides, 1):
        for shape in iter_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False) \
                    and shape.text_frame.text.strip():
                yield n, slide, shape


def check_jiggle(prs, issues):
    """Recurring header/footer-zone elements must not move between slides."""
    sh_in = prs.slide_height / EMU_IN
    groups = defaultdict(list)  # role -> [(slide_no, left, top)]
    for n, _slide, shape in _iter_text_shapes(prs):
        box = _shape_box_in(shape)
        if box is None:
            continue
        left, top = box
        in_footer = top > FOOTER_ZONE * sh_in
        in_header = top < HEADER_ZONE * sh_in
        if not (in_footer or in_header):
            continue
        text = shape.text_frame.text.strip()
        role = "page-number" if PAGENUM_RX.match(text) else text[:60]
        groups[role].append((n, left, top))
    for role, occ in groups.items():
        if len(occ) < 3:
            continue
        ref_l, ref_t = Counter(
            (round(l, 2), round(t, 2)) for _, l, t in occ).most_common(1)[0][0]
        for n, left, top in occ:
            if abs(left - ref_l) > POS_TOL_IN or abs(top - ref_t) > POS_TOL_IN:
                issues["error"].append(
                    f"Slide {n}: '{role}' jiggle — at ({left:.2f},{top:.2f})in "
                    f"but majority position is ({ref_l:.2f},{ref_t:.2f})in")


def check_page_sequence(prs, issues):
    """Plain-digit footer-zone page numbers must track slide positions consistently.

    The invariant is: for any two consecutive numbered slides (n1, n2) with
    page numbers (v1, v2), the number delta must equal the slide-position delta:
        v2 - v1  ==  n2 - n1

    This allows gaps that span unnumbered layouts (title, section-divider, etc.)
    while still catching duplicates (v2 == v1 on different slides) and true
    misnumbering (number delta differs from slide delta).
    """
    sh_in = prs.slide_height / EMU_IN
    pages = []  # (slide_no, value)
    for n, _slide, shape in _iter_text_shapes(prs):
        box = _shape_box_in(shape)
        if box is None or box[1] <= FOOTER_ZONE * sh_in:
            continue
        text = shape.text_frame.text.strip()
        if text.isdigit():
            pages.append((n, int(text)))
    for (n1, v1), (n2, v2) in zip(pages, pages[1:]):
        if v2 - v1 != n2 - n1:
            issues["error"].append(
                f"Slide {n2}: page-number sequence broken "
                f"({v1} on slide {n1}, then {v2} {n2 - n1} slide(s) later)")


def check_axis_scales(prs, issues):
    """Same-type native charts whose explicit value-axis maxima differ by
    more than 10x read as a dishonest scale comparison."""
    groups = defaultdict(list)  # chart_type -> [(slide_no, maximum_scale)]
    for n, slide in enumerate(prs.slides, 1):
        for shape in iter_shapes(slide.shapes):
            if not getattr(shape, "has_chart", False):
                continue
            chart = shape.chart
            try:
                max_scale = chart.value_axis.maximum_scale
            except ValueError:
                continue  # pie/doughnut have no value axis
            if max_scale is not None:
                groups[chart.chart_type].append((n, max_scale))
    for ctype, occ in groups.items():
        if len(occ) < 2:
            continue
        (n_lo, lo) = min(occ, key=lambda o: o[1])
        (n_hi, hi) = max(occ, key=lambda o: o[1])
        if lo > 0 and hi > 10 * lo:
            name = getattr(ctype, "name", str(ctype))
            issues["warn"].append(
                f"{name} charts on slides {n_lo} and {n_hi} have value-axis "
                f"maxima {lo:g} vs {hi:g} (>10x apart) — possible dishonest "
                "scale comparison")


def collect_inventory(prs):
    """(fonts, colors): name/hex -> set of slide numbers using it."""
    fonts, colors = defaultdict(set), defaultdict(set)
    for n, slide, shape in _iter_text_shapes(prs):
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                if run.font.name:
                    fonts[run.font.name].add(n)
                try:
                    rgb = run.font.color.rgb
                except (TypeError, AttributeError):
                    rgb = None
                if rgb is not None:
                    colors[str(rgb)].add(n)
    for n, slide in enumerate(prs.slides, 1):
        for shape in iter_shapes(slide.shapes):
            fill = _solid_fill_rgb(shape)
            if fill is not None:
                colors["%02X%02X%02X" % fill].add(n)
    return fonts, colors


def check_inventory(inv, issues, palette_key=None):
    fonts, colors = inv
    if len(fonts) > MAX_FONTS:
        issues["warn"].append(
            f"deck uses {len(fonts)} fonts ({', '.join(sorted(fonts))}) — "
            f"more than {MAX_FONTS} reads as inconsistent")
    if palette_key:
        from palettes import PALETTES
        pal = PALETTES.get(palette_key)
        if pal is None:
            issues["warn"].append(f"unknown palette '{palette_key}' — "
                                  "skipping whitelist check")
            return
        allowed = {v.upper() for v in pal.values()
                   if isinstance(v, str) and _HEX6.fullmatch(v)}
        allowed |= _chart_series_allowed(pal, issues)
        allowed |= EXTRA_ALLOWED
        for hexv, slides in sorted(colors.items()):
            if hexv.upper() not in allowed:
                where = ", ".join(str(s) for s in sorted(slides)[:5])
                issues["error"].append(
                    f"off-palette color {hexv} on slide(s) {where} "
                    f"(palette: {palette_key})")
    elif len(colors) > MAX_COLORS:
        issues["warn"].append(
            f"deck uses {len(colors)} distinct colors — more than "
            f"{MAX_COLORS}; pass --palette to enforce a whitelist")


def installed_fonts():
    """Lowercased family names visible to the renderer, or None if unknowable.

    fc-list is the most truthful source for what LibreOffice will see;
    matplotlib's font_manager is the cross-platform fallback.
    """
    import shutil
    import subprocess
    if shutil.which("fc-list"):
        try:
            out = subprocess.run(["fc-list", ":", "family"],
                                 capture_output=True, text=True, timeout=20)
            if out.returncode == 0:
                names = set()
                for line in out.stdout.splitlines():
                    for fam in line.split(","):
                        names.add(fam.strip().lower())
                if names:
                    return names
        except (OSError, subprocess.TimeoutExpired):
            pass
    try:
        from matplotlib import font_manager
        names = {n.lower() for n in font_manager.get_font_names()}
        return names or None
    except (ImportError, AttributeError):
        return None


def check_fonts_installed(inv, issues):
    installed = installed_fonts()
    if installed is None:
        issues["warn"].append(
            "cannot enumerate installed fonts (no fc-list or matplotlib) — "
            "visual QA may silently substitute fonts")
        return
    fonts, _ = inv
    for name, slides in sorted(fonts.items()):
        if name.lower() not in installed:
            where = ", ".join(str(s) for s in sorted(slides)[:5])
            issues["warn"].append(
                f"font '{name}' (slide(s) {where}) not installed — renderer "
                "will substitute; QA thumbnails won't match PowerPoint")


def lint_deck(prs, palette_key=None):
    issues = {"error": [], "warn": []}
    check_jiggle(prs, issues)
    check_page_sequence(prs, issues)
    check_axis_scales(prs, issues)
    inv = collect_inventory(prs)
    check_inventory(inv, issues, palette_key)
    check_fonts_installed(inv, issues)
    return issues


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx")
    parser.add_argument("--palette", default=None,
                        help="Enforce this palette's colors as a whitelist")
    parser.add_argument("--assets-dir", default="assets",
                        help="Root assets dir; custom palettes load from "
                             "<assets-dir>/palettes/")
    args = parser.parse_args()
    from palettes import load_custom_palettes
    load_custom_palettes(Path(args.assets_dir) / "palettes")
    prs = Presentation(args.pptx)
    issues = lint_deck(prs, args.palette)
    for e in issues["error"]:
        print(f"  [ERROR] {e}")
    for w in issues["warn"]:
        print(f"  [WARN]  {w}")
    print(f"\n{Path(args.pptx).name}: {len(issues['error'])} error(s), "
          f"{len(issues['warn'])} warning(s)")
    sys.exit(1 if issues["error"] else 0)


if __name__ == "__main__":
    main()
