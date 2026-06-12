#!/usr/bin/env python3
"""
qa_check.py — Programmatic QA for a generated .pptx.

Usage:
    python3 scripts/qa_check.py deck.pptx
    python3 scripts/qa_check.py deck.pptx --text
    python3 scripts/qa_check.py deck.pptx --accessibility   # WCAG AA strict mode
    python3 scripts/qa_check.py deck.pptx --integrity       # OOXML schema audit

Checks include bounds, placeholders, contrast, chart transparency, unique slide
titles, projection font sizes, image alt text, and text density. --integrity
adds OOXML schema validation via the optional openxml-audit package (module or
CLI); when it is not installed the step prints an info line and is skipped.
"""
import re
import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL
from pptx.oxml.ns import qn
from pptx.util import Emu

PLACEHOLDER_RX = re.compile(
    r"lorem|ipsum|xxxx|click to add|\[ add visual \]|placeholder"
    r"|this.*(page|slide).*layout", re.I)
ALT_FILENAME_RX = re.compile(
    r".*\.(png|jpe?g|gif|bmp|svg)$|^image\d*$", re.I)
MIN_PT = 10
PROJECTION_BODY_PT = 18  # federal / Section 508 presentation guidance
EMU_IN = 914400
MAX_WORDS_PER_SLIDE = 220


def _luminance(rgb):
    def chan(c):
        c = c / 255
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    r, g, b = rgb
    return 0.2126 * chan(r) + 0.7152 * chan(g) + 0.0722 * chan(b)


def contrast_ratio(rgb1, rgb2):
    l1, l2 = sorted((_luminance(rgb1), _luminance(rgb2)), reverse=True)
    return (l1 + 0.05) / (l2 + 0.05)


def slide_bg_rgb(slide):
    try:
        fill = slide.background.fill
        if fill.type == MSO_FILL.SOLID:
            c = fill.fore_color.rgb
            return (c[0], c[1], c[2])
        if fill.type == MSO_FILL.GRADIENT:
            stops = [s.color.rgb for s in fill.gradient_stops]
            return tuple(sum(c[i] for c in stops) // len(stops) for i in range(3))
    except (TypeError, ValueError, AttributeError):
        pass
    return None


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def _slide_title(slide):
    """Best-effort (shape, text) title: shapes.title, else largest run."""
    try:
        if slide.shapes.title and slide.shapes.title.text.strip():
            return slide.shapes.title, slide.shapes.title.text.strip()
    except (AttributeError, ValueError):
        pass
    best_shape, best, best_size = None, "", 0
    for shape in iter_shapes(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                size = run.font.size.pt if run.font.size else 0
                if size >= best_size:
                    best_size = size
                    best = run.text.strip()
                    best_shape = shape
    return best_shape, best


def _slide_title_text(slide):
    """Best-effort title: largest bold text or shapes.title."""
    return _slide_title(slide)[1]


def _solid_fill_rgb(shape):
    try:
        if shape.fill.type == MSO_FILL.SOLID:
            c = shape.fill.fore_color.rgb
            return (c[0], c[1], c[2])
    except (TypeError, ValueError, AttributeError, NotImplementedError):
        pass
    return None


def effective_bg(text_shape, filled_below, slide_bg):
    """Background a text box actually sits on: the topmost solid-filled shape
    under its center (cards, chevrons, number circles), else the slide bg."""
    try:
        cx = text_shape.left + text_shape.width // 2
        cy = text_shape.top + text_shape.height // 2
    except (TypeError, AttributeError):
        return slide_bg
    bg = slide_bg
    for l, t, w, h, rgb in filled_below:  # document order == z-order
        if l <= cx <= l + w and t <= cy <= t + h:
            bg = rgb
    return bg


def check_contrast(shape, bg_rgb, issues, where, accessibility=False):
    if bg_rgb is None or not getattr(shape, "has_text_frame", False):
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if not run.text.strip():
                continue
            try:
                c = run.font.color.rgb
            except (TypeError, AttributeError):
                continue
            if c is None:
                continue
            ratio = contrast_ratio((c[0], c[1], c[2]), bg_rgb)
            size = run.font.size.pt if run.font.size else 16
            bold = bool(run.font.bold)
            large = size >= 18 or (size >= 14 and bold)
            min_ratio = 4.5 if accessibility and not large else 3.0
            if ratio < min_ratio:
                issues["error"].append(
                    f"{where}: contrast {ratio:.1f}:1 (<{min_ratio}:1) for "
                    f"{size:.0f}pt text: {run.text[:40]!r}")
            elif ratio < 4.5 and not large:
                issues["warn"].append(
                    f"{where}: contrast {ratio:.1f}:1 (<4.5:1) for small "
                    f"{size:.0f}pt text: {run.text[:40]!r}")


def check_bounds(shape, sw, sh, issues, where):
    try:
        l, t, w, h = shape.left, shape.top, shape.width, shape.height
    except (AttributeError, TypeError):
        return
    if None in (l, t, w, h):
        return
    tol = Emu(int(0.02 * EMU_IN))
    if l < -tol or t < -tol or l + w > sw + tol or t + h > sh + tol:
        issues["error"].append(
            f"{where}: shape outside slide bounds "
            f"({l / EMU_IN:.2f},{t / EMU_IN:.2f} {w / EMU_IN:.2f}x{h / EMU_IN:.2f})")


def check_text(shape, issues, where, accessibility=False):
    if not getattr(shape, "has_text_frame", False):
        return False
    tf = shape.text_frame
    has_text = False
    min_size = None
    for para in tf.paragraphs:
        for run in para.runs:
            if not run.text.strip():
                continue
            has_text = True
            if PLACEHOLDER_RX.search(run.text):
                issues["error"].append(
                    f"{where}: leftover placeholder text: {run.text[:60]!r}")
            if run.font.size is not None:
                pt = run.font.size.pt
                min_size = pt if min_size is None else min(min_size, pt)
                if pt < MIN_PT:
                    issues["error"].append(
                        f"{where}: {pt:.0f}pt text (<{MIN_PT}pt): {run.text[:40]!r}")
                elif accessibility and pt < PROJECTION_BODY_PT and not run.font.bold:
                    # projection-size guidance is opt-in: consulting decks are
                    # read documents, and compact density is intentional
                    issues["error"].append(
                        f"{where}: body text {pt:.0f}pt (<{PROJECTION_BODY_PT}pt "
                        f"for projection): {run.text[:40]!r}")

    if has_text and min_size and shape.width and shape.height \
            and len(tf.text.strip()) > 2:  # skip decorative glyphs (e.g. big quote marks)
        size_in = min_size / 72
        chars_per_line = max(int((shape.width / EMU_IN) / (0.52 * size_in)), 1)
        est_lines = sum(
            max(1, -(-len(p.text) // chars_per_line))
            for p in tf.paragraphs if p.text.strip())
        est_h = est_lines * size_in * 1.45
        if est_h > (shape.height / EMU_IN) * 1.30:
            issues["warn"].append(
                f"{where}: text likely overflows box (~{est_lines} lines): "
                f"{tf.text[:50]!r}")
    return has_text


def check_chart_bg(shape, issues, where):
    if not getattr(shape, "has_chart", False):
        return False
    cs = shape.chart._chartSpace
    spPr = cs.find(qn("c:spPr"))
    transparent = spPr is not None and spPr.find(qn("a:noFill")) is not None
    if not transparent:
        issues["error"].append(
            f"{where}: chart has opaque background (white-box-on-dark defect)")
    return True


def check_image_alt(shape, issues, where, accessibility=False):
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False
    descr = shape._element._nvXxPr.cNvPr.get("descr", "")
    if not descr or not descr.strip():
        msg = f"{where}: image missing alt text (descr attribute)"
        (issues["error"] if accessibility else issues["warn"]).append(msg)
    elif accessibility and ALT_FILENAME_RX.match(descr.strip()):
        issues["error"].append(
            f"{where}: alt text is a filename ({descr.strip()!r}) — "
            "describe the image instead")
    return True


def _tc_merged(tc):
    def _span(name):
        v = tc.get(name)
        try:
            return int(v) if v else 1
        except ValueError:
            return 1
    if _span("gridSpan") > 1 or _span("rowSpan") > 1:
        return True
    return tc.get("hMerge") in ("1", "true") or tc.get("vMerge") in ("1", "true")


def check_table_a11y(shape, issues, where):
    """MS Accessibility Checker table rules: marked header row, no merges."""
    if not getattr(shape, "has_table", False):
        return
    tbl = shape.table._tbl
    tblPr = tbl.tblPr
    if tblPr is None or tblPr.get("firstRow") not in ("1", "true"):
        issues["error"].append(
            f"{where}: table not marked with a header row (firstRow) — "
            "screen readers cannot announce column headers")
    if any(_tc_merged(tc) for tc in tbl.iter(qn("a:tc"))):
        issues["warn"].append(
            f"{where}: table has merged cells — screen readers may misread; "
            "prefer simple structure")


def check_reading_order(slide, issues, where):
    """Title shape should be the first text-bearing shape in spTree order."""
    title_shape, _ = _slide_title(slide)
    if title_shape is None:
        return
    first_text = next(
        (s for s in iter_shapes(slide.shapes)
         if getattr(s, "has_text_frame", False)
         and s.text_frame.text.strip()), None)
    if first_text is not None and first_text.shape_id != title_shape.shape_id:
        issues["warn"].append(
            f"{where}: title is not first in reading order — screen readers "
            "announce it after other content")


def check_deck(pptx_path, accessibility=False):
    prs = Presentation(str(pptx_path))
    sw, sh = prs.slide_width, prs.slide_height
    issues = {"error": [], "warn": []}
    titles_seen = {}

    for n, slide in enumerate(prs.slides, 1):
        where = f"Slide {n}"
        has_text = has_visual = False
        shape_count = word_count = 0
        has_picture = any(s.shape_type == MSO_SHAPE_TYPE.PICTURE
                          for s in iter_shapes(slide.shapes))
        bg_rgb = None if has_picture else slide_bg_rgb(slide)

        title = _slide_title_text(slide)
        if not title:
            msg = f"{where}: no detectable slide title"
            (issues["error"] if accessibility else issues["warn"]).append(msg)
        else:
            titles_seen.setdefault(title.strip().casefold(), []).append(n)
        if accessibility:
            check_reading_order(slide, issues, where)

        filled_below = []  # accumulating: shapes earlier in doc order are below
        for shape in iter_shapes(slide.shapes):
            shape_count += 1
            check_bounds(shape, sw, sh, issues, where)
            if check_text(shape, issues, where, accessibility):
                has_text = True
                word_count += len(shape.text_frame.text.split())
            text_bg = effective_bg(shape, filled_below, bg_rgb) \
                if bg_rgb is not None else None
            check_contrast(shape, text_bg, issues, where, accessibility)
            fill_rgb = _solid_fill_rgb(shape)
            if fill_rgb is not None and None not in (
                    getattr(shape, "left", None), getattr(shape, "top", None),
                    getattr(shape, "width", None), getattr(shape, "height", None)):
                filled_below.append((shape.left, shape.top, shape.width,
                                     shape.height, fill_rgb))
            if check_chart_bg(shape, issues, where):
                has_visual = True
            if check_image_alt(shape, issues, where, accessibility):
                has_visual = True
            if accessibility:
                check_table_a11y(shape, issues, where)
            if shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.TABLE):
                has_visual = True
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and not getattr(
                    shape, "has_text_frame", False):
                has_visual = True

        if shape_count == 0:
            issues["warn"].append(f"{where}: slide is empty")
        elif has_text and not has_visual and shape_count <= 2:
            issues["warn"].append(
                f"{where}: text-only slide — add icon, chart, or image")
        if word_count > MAX_WORDS_PER_SLIDE:
            issues["warn"].append(
                f"{where}: {word_count} words — dense slide, consider splitting")

    for key, ns in titles_seen.items():
        if len(ns) > 1:
            where_list = ", ".join(str(s) for s in ns)
            msg = (f"duplicate slide title {key!r} on slides {where_list} — "
                   "titles must be unique for navigation")
            (issues["error"] if accessibility else issues["warn"]).append(msg)

    return issues


def dump_text(pptx_path):
    prs = Presentation(str(pptx_path))
    for n, slide in enumerate(prs.slides, 1):
        print(f"\n=== Slide {n} ===")
        for shape in iter_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                print(shape.text_frame.text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            print(f"[notes] {slide.notes_slide.notes_text_frame.text}")


NUM_TOKEN_RX = re.compile(
    r"[$€£]?\d[\d,.]*(?:%|bps|bn|B|M|k|x|pp)?|FY\d{2,4}|Q[1-4]\s?\d{2,4}",
    re.I)


def dump_numbers(pptx_path):
    """All numeric/period tokens per slide — input for the LLM
    cross-slide consistency check (see references/qa-guide.md)."""
    prs = Presentation(str(pptx_path))
    for n, slide in enumerate(prs.slides, 1):
        tokens = []
        for shape in iter_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                tokens += [t.strip() for t in
                           NUM_TOKEN_RX.findall(shape.text_frame.text)
                           if t.strip()]
        title = _slide_title_text(slide)
        print(f"Slide {n} [{title[:60]}]: {', '.join(tokens) if tokens else '—'}")


def _integrity_via_module(mod, pptx_path):
    """Validate via an importable openxml_audit module. The package's API is
    probed defensively (validate/audit/check, called with the file path and
    expected to return an iterable of error strings/objects). Returns a list
    of error strings, or None when no recognized entry point exists."""
    fn = next((getattr(mod, name) for name in ("validate", "audit", "check")
               if callable(getattr(mod, name, None))), None)
    if fn is None:
        print("  [WARN] openxml-audit installed but exposes no "
              "validate/audit/check entry point — integrity check skipped",
              file=sys.stderr)
        return None
    result = fn(str(pptx_path))
    return [str(item) for item in result] if result else []


def _integrity_via_cli(cli, pptx_path):
    """Validate via an `openxml-audit <file>` CLI on PATH. Non-zero exit →
    each non-empty output line becomes an integrity error."""
    try:
        r = subprocess.run([cli, str(pptx_path)],
                           capture_output=True, text=True, timeout=120)
    except (subprocess.TimeoutExpired, OSError) as exc:
        return [f"openxml-audit CLI failed ({exc.__class__.__name__})"]
    if r.returncode == 0:
        return []
    lines = [ln.strip() for ln in (r.stdout + "\n" + r.stderr).splitlines()
             if ln.strip()]
    return lines or [f"openxml-audit exited {r.returncode}"]


def run_integrity(pptx_path):
    """OOXML schema validation via the optional openxml-audit package.

    Returns a list of error strings (empty = valid). When the package is not
    installed (neither importable module nor CLI on PATH) prints a single
    info line and returns [] — never a failure."""
    try:
        import openxml_audit  # pip name: openxml-audit
    except ImportError:
        openxml_audit = None
    if openxml_audit is not None:
        errors = _integrity_via_module(openxml_audit, pptx_path)
        if errors is not None:
            return errors
        return []  # module present but unusable — warned above, skip
    cli = shutil.which("openxml-audit")
    if cli:
        return _integrity_via_cli(cli, pptx_path)
    print("  [INFO] pip install openxml-audit for OOXML schema validation "
          "(skipped)")
    return []


def main():
    if len(sys.argv) < 2:
        print("Usage: qa_check.py deck.pptx [--text] [--numbers] "
              "[--accessibility] [--integrity]")
        sys.exit(2)
    path = Path(sys.argv[1])
    if "--numbers" in sys.argv:
        dump_numbers(path)
        sys.exit(0)
    if "--text" in sys.argv:
        dump_text(path)
        sys.exit(0)
    accessibility = "--accessibility" in sys.argv
    issues = check_deck(path, accessibility=accessibility)
    if "--integrity" in sys.argv:
        issues["error"].extend(f"integrity: {e}" for e in run_integrity(path))
    for e in issues["error"]:
        print(f"  [ERROR] {e}")
    for w in issues["warn"]:
        print(f"  [WARN]  {w}")
    print(f"\n{path.name}: {len(issues['error'])} error(s), "
          f"{len(issues['warn'])} warning(s)")
    sys.exit(1 if issues["error"] else 0)


if __name__ == "__main__":
    main()
